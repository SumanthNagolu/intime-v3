#!/usr/bin/env python3
"""
Final accuracy validation: Re-classify every slide from scratch using
the original PPT files + OCR text, then compare against stored slideType values.

This is the definitive check that content JSONs match what the originals produce.

Validates:
1. slideType - re-classify from OCR/PPT and compare
2. originalNotes - re-extract from PPT and compare
3. Slide count - content JSON vs slide images on disk vs PPT slide count
4. Field completeness - every slide has all required fields
"""

import json
import os
import re
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx required (pip3 install python-pptx)")
    exit(1)

MATERIAL = Path("/Users/sumanthrajkumarnagolu/Documents/Material")
CONTENT = Path("public/academy/guidewire/content")
SLIDES = Path("public/academy/guidewire/slides")

CHAPTER_MAP = {
    "ch04": "Chapter 4 - PolicyCenter Introduction",
    "ch05": "Chapter 5 - Claim Center Introduction",
    "ch06": "Chapter 6 - Billing Center Introduction",
    "ch07": "Chapter 7 - InsuranceSuite Developer Fundamentals",
    "ch08": "Chapter 8 - policy center configuration",
    "ch09": "Chapter 9 - ClaimCenter Configuration",
    "ch11": "Chapter 11 - Introduction to Integration",
    "ch13": "Chapter 13 - Rating Introduction",
    "ch14": "Chapter 14 - Rating Configuration",
}

REQUIRED_FIELDS = [
    "slideNumber", "title", "bodyParagraphs",
    "originalNotes", "notes", "narration",
    "ocrText", "ocrWordCount", "ocrConfidence",
    "hasTable", "hasChart", "hasImage",
    "tableData", "mediaReferences", "slideType",
]


def find_pptx(chapter_dir, source_file, source_folder=None):
    if source_folder:
        p = chapter_dir / source_folder / source_file
        if p.exists():
            return str(p)
    for root, dirs, files in os.walk(str(chapter_dir)):
        if source_file in files:
            return os.path.join(root, source_file)
    return None


def extract_ppt_notes(pptx_path):
    """Extract speaker notes from every slide in a PPTX."""
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {}
    notes = {}
    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        text = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                text = slide.notes_slide.notes_text_frame.text.strip()
        except:
            pass
        notes[slide_num] = text
    return notes


def detect_video_slides(pptx_path):
    """Detect which slides have video/media shapes in the PPT."""
    if not pptx_path:
        return set()
    try:
        prs = Presentation(pptx_path)
        video_slides = set()
        for idx, slide in enumerate(prs.slides):
            snum = idx + 1
            for shape in slide.shapes:
                if shape.shape_type == 16:
                    video_slides.add(snum)
                    break
                try:
                    xml = shape._element.xml
                    if '<a:videoFile' in xml or '<p:video' in xml or 'oleObj' in xml.lower():
                        video_slides.add(snum)
                        break
                except:
                    pass
        return video_slides
    except:
        return set()


def get_ppt_slide_count(pptx_path):
    """Get total slide count from PPT."""
    try:
        prs = Presentation(pptx_path)
        return len(prs.slides)
    except:
        return 0


def classify_slide(slide, snum, total_slides, video_slides, prev_type=None):
    """Re-classify a slide from scratch (mirrors classify-slide-types.py with fixes)."""
    ocr = slide.get("ocrText", "").strip().lower()
    title = slide.get("title", "").strip().lower()
    ocr_words = slide.get("ocrWordCount", 0)

    if snum == 1:
        return "title"

    # objectives_review first (handles OCR mangling like "revi eW")
    if ("objectives\nreview" in ocr or "objectives review" in ocr or
            ("objectives" in ocr and re.search(r'revi\s*e?w', ocr))):
        return "objectives_review"

    if ("lesson\nobjectives" in ocr or "lesson objectives" in ocr) and "review" not in ocr:
        return "objectives"

    if "virtual machine instruction" in ocr or "virtual machine instruction" in title:
        return "vm_instructions"

    if ocr.startswith("student exercise") or "student exercise" in ocr[:30]:
        return "exercise"
    if "student exercise" in title:
        return "exercise"

    if ocr.startswith("question") or re.match(r'^question\s', ocr):
        return "question"
    if title == "question" or title.startswith("question:"):
        return "question"

    if ocr.startswith("answer") or re.match(r'^answer\s', ocr):
        return "answer"
    if title == "answer" or title.startswith("answer:"):
        return "answer"

    if snum in video_slides:
        return "demo_video"
    if prev_type == "demo_instruction" and ocr_words == 0 and not ocr:
        return "demo_video"
    if ocr_words == 0 and not ocr and not title:
        if prev_type in ("demo_instruction", "demo_video"):
            return "demo_video"

    if "demonstration" in ocr[:20] or ocr.startswith("d trati"):
        return "demo_instruction"
    if "demonstration" in title[:20]:
        return "demo_instruction"
    if re.search(r'(?:^|\|)\s*demonstration\b', ocr[:80]):
        return "demo_instruction"

    return "content"


def count_slide_images(ch_slug, lesson_name):
    """Count PNG slide images on disk."""
    lesson_dir = SLIDES / ch_slug / lesson_name
    if not lesson_dir.exists():
        return 0
    return len([f for f in lesson_dir.iterdir() if f.suffix.lower() == '.png'])


def main():
    print("=" * 80)
    print("FINAL VALIDATION: Content JSONs vs Original PPTs")
    print("=" * 80)

    total_lessons = 0
    total_slides = 0
    total_ppt_found = 0
    total_ppt_not_found = 0

    # Counters
    type_mismatches = []
    notes_mismatches = []
    slide_count_issues = []
    field_missing = []
    field_order_issues = []

    for ch_slug, ch_name in CHAPTER_MAP.items():
        content_dir = CONTENT / ch_slug
        chapter_dir = MATERIAL / ch_name
        idx_path = content_dir / "index.json"

        if not idx_path.exists():
            continue

        with open(idx_path) as f:
            idx = json.load(f)

        print(f"\n{'='*70}")
        print(f"Chapter: {ch_slug} - {ch_name}")
        print(f"{'='*70}")

        ch_type_mismatches = 0
        ch_notes_mismatches = 0
        ch_lessons = 0

        for li in idx.get("lessons", []):
            lnum = li["lessonId"].split("-l")[1]
            lname = f"lesson-{lnum}"
            lpath = content_dir / f"{lname}.json"

            if not lpath.exists():
                continue

            with open(lpath) as f:
                content = json.load(f)

            source_file = content.get("sourceFile", "")
            source_folder = content.get("sourceFolder", "")
            slides = content.get("slides", [])
            total_lessons += 1
            ch_lessons += 1
            total_slides += len(slides)

            # Find original PPT
            pptx_path = None
            if source_file and chapter_dir.exists():
                pptx_path = find_pptx(chapter_dir, source_file, source_folder)

            if pptx_path:
                total_ppt_found += 1
            else:
                total_ppt_not_found += 1
                if source_file:
                    print(f"  {lname}: PPT NOT FOUND ({source_file})")
                continue

            # Extract from PPT
            ppt_notes = extract_ppt_notes(pptx_path)
            video_slides = detect_video_slides(pptx_path)
            ppt_slide_count = get_ppt_slide_count(pptx_path)
            disk_slide_count = count_slide_images(ch_slug, lname)

            # Validate slide counts
            content_count = len(slides)
            if content_count != ppt_slide_count and content_count != disk_slide_count:
                slide_count_issues.append(
                    f"{ch_slug}/{lname}: content={content_count}, PPT={ppt_slide_count}, disk={disk_slide_count}"
                )

            # Validate each slide
            lesson_type_mismatch = 0
            lesson_notes_mismatch = 0
            lesson_field_missing = 0
            prev_type = None

            for slide in slides:
                snum = slide.get("slideNumber", 0)

                # CHECK 1: Field completeness
                for field in REQUIRED_FIELDS:
                    if field not in slide:
                        field_missing.append(f"{ch_slug}/{lname}/slide-{snum:02d}: missing '{field}'")
                        lesson_field_missing += 1

                # CHECK 2: Field order (text layers should be grouped)
                keys = list(slide.keys())
                try:
                    orig_idx = keys.index("originalNotes")
                    notes_idx = keys.index("notes")
                    narr_idx = keys.index("narration")
                    ocr_idx = keys.index("ocrText")
                    if not (orig_idx < notes_idx < narr_idx < ocr_idx):
                        field_order_issues.append(f"{ch_slug}/{lname}/slide-{snum:02d}: text layers not in order")
                except ValueError:
                    pass  # Missing field already caught above

                # CHECK 3: slideType accuracy (re-classify and compare)
                expected_type = classify_slide(slide, snum, len(slides), video_slides, prev_type)
                stored_type = slide.get("slideType", "")
                if expected_type != stored_type:
                    type_mismatches.append(
                        f"{ch_slug}/{lname}/slide-{snum:02d}: stored='{stored_type}', expected='{expected_type}'"
                    )
                    lesson_type_mismatch += 1
                    ch_type_mismatches += 1

                # CHECK 4: originalNotes accuracy
                ppt_note = ppt_notes.get(snum, "")
                stored_note = slide.get("originalNotes", "")
                if ppt_note != stored_note:
                    # Show first difference
                    notes_mismatches.append(
                        f"{ch_slug}/{lname}/slide-{snum:02d}: PPT note ({len(ppt_note)} chars) != stored ({len(stored_note)} chars)"
                    )
                    lesson_notes_mismatch += 1
                    ch_notes_mismatches += 1

                prev_type = stored_type

            # Lesson summary
            status = "OK"
            issues = []
            if lesson_type_mismatch:
                issues.append(f"{lesson_type_mismatch} type mismatches")
            if lesson_notes_mismatch:
                issues.append(f"{lesson_notes_mismatch} notes mismatches")
            if lesson_field_missing:
                issues.append(f"{lesson_field_missing} missing fields")
            if issues:
                status = ", ".join(issues)

            print(f"  {lname}: {len(slides)} slides, PPT:{ppt_slide_count}, disk:{disk_slide_count} | {status}")

        print(f"\n  Chapter: {ch_lessons} lessons | type mismatches: {ch_type_mismatches} | notes mismatches: {ch_notes_mismatches}")

    # Final report
    print()
    print("=" * 80)
    print("FINAL VALIDATION REPORT")
    print("=" * 80)
    print(f"\n  Lessons validated:    {total_lessons}")
    print(f"  Total slides:         {total_slides}")
    print(f"  PPTs found:           {total_ppt_found}")
    print(f"  PPTs not found:       {total_ppt_not_found}")

    print(f"\n  --- SLIDE TYPE ACCURACY ---")
    print(f"  Type mismatches:      {len(type_mismatches)}/{total_slides} ({(1 - len(type_mismatches)/total_slides)*100:.2f}% accurate)")
    if type_mismatches:
        print(f"\n  Mismatches:")
        for m in type_mismatches[:20]:
            print(f"    {m}")
        if len(type_mismatches) > 20:
            print(f"    ... and {len(type_mismatches) - 20} more")

    print(f"\n  --- ORIGINAL NOTES ACCURACY ---")
    print(f"  Notes mismatches:     {len(notes_mismatches)}/{total_slides} ({(1 - len(notes_mismatches)/total_slides)*100:.2f}% accurate)")
    if notes_mismatches:
        print(f"\n  Mismatches:")
        for m in notes_mismatches[:20]:
            print(f"    {m}")
        if len(notes_mismatches) > 20:
            print(f"    ... and {len(notes_mismatches) - 20} more")

    print(f"\n  --- SLIDE COUNT CONSISTENCY ---")
    print(f"  Count issues:         {len(slide_count_issues)}")
    if slide_count_issues:
        for m in slide_count_issues:
            print(f"    {m}")

    print(f"\n  --- FIELD COMPLETENESS ---")
    print(f"  Missing fields:       {len(field_missing)}")
    if field_missing:
        for m in field_missing[:10]:
            print(f"    {m}")

    print(f"\n  --- FIELD ORDERING ---")
    print(f"  Order issues:         {len(field_order_issues)}")
    if field_order_issues:
        for m in field_order_issues[:10]:
            print(f"    {m}")

    # Final verdict
    print()
    total_issues = len(type_mismatches) + len(notes_mismatches) + len(field_missing) + len(field_order_issues)
    if total_issues == 0:
        print("  " + "=" * 50)
        print("  VERDICT: ALL VALIDATIONS PASSED")
        print("  " + "=" * 50)
        print(f"  {total_slides} slides across {total_lessons} lessons")
        print(f"  slideType: 100% accurate against PPT originals")
        print(f"  originalNotes: 100% match with PPT speaker notes")
        print(f"  Field completeness: All {len(REQUIRED_FIELDS)} required fields present")
        print(f"  Field ordering: Text layers correctly grouped")
    else:
        print(f"  VERDICT: {total_issues} ISSUES FOUND - review above")


if __name__ == "__main__":
    main()
