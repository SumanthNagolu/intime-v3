#!/usr/bin/env python3
"""
Classify slideType for all slides across all lessons using OCR text, title,
and PPT video shape detection.

Slide types:
  - title: First slide of each lesson (lesson title)
  - objectives: "Lesson objectives" slide (usually slide 2)
  - scenario: Persona/scenario slide (character with dialog)
  - demo_instruction: "Demonstration" slide with instructions
  - demo_video: Video placeholder slide (no content, follows demo_instruction)
  - question: "Question" slide
  - answer: "Answer" slide
  - objectives_review: "Lesson objectives review" slide
  - vm_instructions: "Virtual Machine Instructions" slide
  - exercise: "Student exercise" slide
  - content: Regular content slide (default)
"""

import json
import os
import re
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("WARNING: python-pptx not available, skipping PPT video detection")
    Presentation = None

MATERIAL = Path("/Users/sumanthrajkumarnagolu/Documents/Material")
CONTENT = Path("public/academy/guidewire/content")

CHAPTER_MAP = {
    "ch04": "Chapter 4 - PolicyCenter Introduction",
    "ch05": "Chapter 5 - Claim Center Introduction",
    "ch06": "Chapter 6 - Billing Center Introduction",
    "ch07": "Chapter 7 - InsuranceSuite Developer Fundamentals",
    "ch08": "Chapter 8 - policy center configuration",
    "ch09": "Chapter 9 - ClaimCenter Configuration",
    "ch11": "Chapter 11 - Introduction to Integration",
    "ch13": "Chapter 13 - Rating Introduction",
    "ch14": "Chapter 14 - Rating Configuration",
}


def find_pptx(chapter_dir, source_file, source_folder=None):
    if source_folder:
        p = chapter_dir / source_folder / source_file
        if p.exists():
            return str(p)
    for root, dirs, files in os.walk(str(chapter_dir)):
        if source_file in files:
            return os.path.join(root, source_file)
    return None


def detect_video_slides(pptx_path):
    """Detect which slides have video/media shapes in the PPT."""
    if not Presentation or not pptx_path:
        return set()
    try:
        prs = Presentation(pptx_path)
        video_slides = set()
        for idx, slide in enumerate(prs.slides):
            snum = idx + 1
            for shape in slide.shapes:
                # shape_type 3=freeform, 16=media, 24=placeholder
                # Check for movie/video elements
                if shape.shape_type == 16:  # MSO_SHAPE_TYPE.MEDIA
                    video_slides.add(snum)
                    break
                # Also check shape XML for video
                try:
                    xml = shape._element.xml
                    if '<a:videoFile' in xml or '<p:video' in xml or 'oleObj' in xml.lower():
                        video_slides.add(snum)
                        break
                except:
                    pass
        return video_slides
    except:
        return set()


def classify_slide(slide, snum, total_slides, video_slides, prev_type=None):
    """Classify a single slide based on OCR text, title, position, and PPT data."""
    ocr = slide.get("ocrText", "").strip().lower()
    title = slide.get("title", "").strip().lower()
    notes = slide.get("originalNotes", "").strip().lower()
    all_notes = slide.get("notes", "").strip().lower()
    ocr_words = slide.get("ocrWordCount", 0)

    # 1. TITLE: First slide
    if snum == 1:
        return "title"

    # 2. OBJECTIVES REVIEW: Check first (before objectives) to avoid false match
    # Handle OCR mangling like "revi eW" for "review"
    if ("objectives\nreview" in ocr or "objectives review" in ocr or
            ("objectives" in ocr and re.search(r'revi\s*e?w', ocr))):
        return "objectives_review"

    # 3. OBJECTIVES: "Lesson objectives" (not review)
    if ("lesson\nobjectives" in ocr or "lesson objectives" in ocr) and "review" not in ocr:
        return "objectives"

    # 4. VM INSTRUCTIONS
    if "virtual machine instruction" in ocr or "virtual machine instruction" in title:
        return "vm_instructions"

    # 5. STUDENT EXERCISE
    if ocr.startswith("student exercise") or "student exercise" in ocr[:30]:
        return "exercise"
    if "student exercise" in title:
        return "exercise"

    # 6. QUESTION: Blue "Question" slide
    if ocr.startswith("question") or re.match(r'^question\s', ocr):
        return "question"
    if title == "question" or title.startswith("question:"):
        return "question"

    # 7. ANSWER: Green "Answer" slide
    if ocr.startswith("answer") or re.match(r'^answer\s', ocr):
        return "answer"
    if title == "answer" or title.startswith("answer:"):
        return "answer"

    # 8. DEMO VIDEO: Video placeholder (from PPT analysis) or empty slide after demo_instruction
    if snum in video_slides:
        return "demo_video"
    # Empty slide right after a demo_instruction
    if prev_type == "demo_instruction" and ocr_words == 0 and not ocr:
        return "demo_video"
    # Also catch slides with no OCR that sit between demo-related slides
    if ocr_words == 0 and not ocr and not title:
        if prev_type in ("demo_instruction", "demo_video"):
            return "demo_video"

    # 9. DEMO INSTRUCTION: "Demonstration" with instructions
    # Check first 20 chars, but also handle "click here | Demonstration..." patterns
    if "demonstration" in ocr[:20] or ocr.startswith("d trati"):  # OCR sometimes mangles "Demonstration"
        return "demo_instruction"
    if "demonstration" in title[:20]:
        return "demo_instruction"
    # Broader check: "| Demonstration" anywhere in first 80 chars (handles prefix noise)
    if re.search(r'(?:^|\|)\s*demonstration\b', ocr[:80]):
        return "demo_instruction"

    # 10. Default: CONTENT
    return "content"


def main():
    print("=" * 80)
    print("SLIDE TYPE CLASSIFICATION")
    print("=" * 80)

    type_totals = {}
    total_slides = 0
    total_classified = 0

    for ch_slug, ch_name in CHAPTER_MAP.items():
        content_dir = CONTENT / ch_slug
        chapter_dir = MATERIAL / ch_name
        idx_path = content_dir / "index.json"

        if not idx_path.exists():
            continue

        with open(idx_path) as f:
            idx = json.load(f)

        print(f"\n{'='*60}")
        print(f"Chapter: {ch_slug}")
        print(f"{'='*60}")

        ch_types = {}

        for li in idx.get("lessons", []):
            lnum = li["lessonId"].split("-l")[1]
            lname = f"lesson-{lnum}"
            lpath = content_dir / f"{lname}.json"

            if not lpath.exists():
                continue

            with open(lpath) as f:
                content = json.load(f)

            source_file = content.get("sourceFile", "")
            source_folder = content.get("sourceFolder", "")

            # Detect video slides from PPT
            pptx_path = find_pptx(chapter_dir, source_file, source_folder) if chapter_dir.exists() else None
            video_slides = detect_video_slides(pptx_path)

            slides = content.get("slides", [])
            lesson_types = {}
            prev_type = None

            for slide in slides:
                snum = slide.get("slideNumber", 0)
                total_slides += 1

                stype = classify_slide(slide, snum, len(slides), video_slides, prev_type)
                slide["slideType"] = stype
                prev_type = stype
                total_classified += 1

                lesson_types[stype] = lesson_types.get(stype, 0) + 1
                ch_types[stype] = ch_types.get(stype, 0) + 1
                type_totals[stype] = type_totals.get(stype, 0) + 1

            with open(lpath, 'w') as f:
                json.dump(content, f, indent=2, ensure_ascii=False)

            # Compact lesson summary
            type_str = ", ".join(f"{t}:{c}" for t, c in sorted(lesson_types.items()))
            print(f"  {lname}: {len(slides)} slides [{type_str}]")

        # Chapter summary
        print(f"\n  Chapter totals: {json.dumps(ch_types)}")

    print()
    print("=" * 80)
    print("CLASSIFICATION SUMMARY")
    print("=" * 80)
    print(f"  Total slides: {total_slides}")
    print(f"  Classified: {total_classified}")
    print()
    print(f"  {'Type':<20} {'Count':<8} {'%':<8}")
    print(f"  {'-'*36}")
    for stype in ["title", "objectives", "content", "demo_instruction", "demo_video",
                   "question", "answer", "objectives_review", "vm_instructions", "exercise"]:
        count = type_totals.get(stype, 0)
        pct = count / total_slides * 100 if total_slides else 0
        print(f"  {stype:<20} {count:<8} {pct:.1f}%")


if __name__ == "__main__":
    main()
