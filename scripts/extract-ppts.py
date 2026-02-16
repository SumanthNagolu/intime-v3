#!/usr/bin/env python3
"""
Extract PPT content into structured JSON + slide images for the Guidewire Academy portal.
Processes all PPTX/PPTM files from the training material directory.
"""

import json
import os
import sys
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Source and output directories
MATERIAL_DIR = "/Users/sumanthrajkumarnagolu/Documents/Material"
OUTPUT_DIR = "/Users/sumanthrajkumarnagolu/Projects/intime-v3/public/academy/guidewire/content"

# Chapter mapping (folder name -> chapter config)
CHAPTER_MAP = {
    "Chapter 1 - Guidewire Cloud Overview": {"id": 1, "slug": "ch01"},
    "Chapter 2 - Surepath Overview": {"id": 2, "slug": "ch02"},
    "Chapter 3 - InsuranceSuite Implementation Tools": {"id": 3, "slug": "ch03"},
    "Chapter 4 - PolicyCenter Introduction": {"id": 4, "slug": "ch04"},
    "Chapter 5 - Claim Center Introduction": {"id": 5, "slug": "ch05"},
    "Chapter 6 - Billing Center Introduction": {"id": 6, "slug": "ch06"},
    "Chapter 7 - InsuranceSuite Developer Fundamentals": {"id": 7, "slug": "ch07"},
    "Chapter 8 - policy center configuration": {"id": 8, "slug": "ch08"},
    "Chapter 9 - ClaimCenter Configuration": {"id": 9, "slug": "ch09"},
    "Chapter 10 - BillingCenter Configuration": {"id": 10, "slug": "ch10"},
    "Chapter 11 - Introduction to Integration": {"id": 11, "slug": "ch11"},
    "Chapter 12 - Advanced product Designer": {"id": 12, "slug": "ch12"},
    "Chapter 13 - Rating Introduction": {"id": 13, "slug": "ch13"},
    "Chapter 14 - Rating Configuration": {"id": 14, "slug": "ch14"},
}


def extract_text_from_shape(shape):
    """Extract text content from a shape, preserving structure."""
    if not shape.has_text_frame:
        return None

    paragraphs = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Detect bullet level
        level = para.level if hasattr(para, 'level') else 0

        # Check if bold (title-like)
        is_bold = False
        if para.runs:
            is_bold = any(run.font.bold for run in para.runs if run.font.bold is not None)

        paragraphs.append({
            "text": text,
            "level": level,
            "bold": is_bold,
        })

    return paragraphs if paragraphs else None


def extract_slide_content(slide, slide_number):
    """Extract content from a single slide."""
    title = ""
    body_paragraphs = []
    has_table = False
    has_chart = False
    has_image = False
    table_data = None
    media_files = []

    for shape in slide.shapes:
        # Title
        if shape.has_text_frame and shape.shape_id == slide.shapes.title_id if hasattr(slide.shapes, 'title_id') else False:
            title = shape.text_frame.text.strip()
        elif hasattr(shape, "is_placeholder") and shape.is_placeholder:
            ph = shape.placeholder_format
            if ph and ph.idx == 0:  # Title placeholder
                if shape.has_text_frame:
                    title = shape.text_frame.text.strip()
                continue

        # Try to get title from shape name
        if not title and shape.has_text_frame and 'title' in (shape.name or '').lower():
            title = shape.text_frame.text.strip()
            continue

        # Table
        if shape.has_table:
            has_table = True
            table = shape.table
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
            table_data = rows
            continue

        # Chart
        if shape.has_chart:
            has_chart = True
            continue

        # Images/media
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            has_image = True
            continue

        # Check for embedded media (videos)
        if hasattr(shape, 'click_action') and shape.click_action:
            if hasattr(shape.click_action, 'hyperlink') and shape.click_action.hyperlink:
                href = shape.click_action.hyperlink.address
                if href and any(ext in href.lower() for ext in ['.mp4', '.mkv', '.avi', '.wmv']):
                    media_files.append(href)

        # Body text
        text_content = extract_text_from_shape(shape)
        if text_content:
            # Skip if this is just the title repeated
            if not (len(text_content) == 1 and text_content[0]["text"] == title):
                body_paragraphs.extend(text_content)

    # Fallback: get title from first shape with text
    if not title and slide.shapes:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip()
                # Remove from body if it was added
                body_paragraphs = [p for p in body_paragraphs if p["text"] != title]
                break

    # Extract speaker notes
    notes = ""
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame:
            notes = notes_slide.notes_text_frame.text.strip()

    return {
        "slideNumber": slide_number,
        "title": title,
        "bodyParagraphs": body_paragraphs,
        "notes": notes,
        "hasTable": has_table,
        "hasChart": has_chart,
        "hasImage": has_image,
        "tableData": table_data,
        "mediaReferences": media_files,
    }


def find_videos_for_lesson(chapter_dir, lesson_folder_name):
    """Find video files associated with a lesson."""
    lesson_dir = os.path.join(chapter_dir, lesson_folder_name)
    videos = []

    if os.path.isdir(lesson_dir):
        for f in sorted(os.listdir(lesson_dir)):
            if f.lower().endswith(('.mp4', '.mkv', '.avi', '.wmv')):
                videos.append(f)

    return videos


def extract_lesson_title(pptx_filename):
    """Extract a clean lesson title from the PPTX filename."""
    name = os.path.splitext(pptx_filename)[0]

    # Remove common prefixes like "IS_Fund_02 - ", "In_policy_01_", etc.
    patterns = [
        r'^IS_Fund[_\s]*\d+[_\s]*[-–]\s*',
        r'^IS_Fund[_\s]*\d+[_\s]*',
        r'^IS[_\s]*Fund[_\s]*[-–]?\s*\d+[_\s]*[-–]\s*',
        r'^In_policy[_\s]*\d+[_\s]*[-–]?\s*',
        r'^In_policy[_\s]*\d+[_\s]*',
        r'^PC_Intro[_\s]*\d+[_\s]*[-–]?\s*',
        r'^IS_Claim[_\s]*\d+[_\s]*[-–]?\s*',
        r'^IS_Claim[_\s]*__?\d+[_\s]*[-–]?\s*',
        r'^BillingCenter[_\s]*\d+[_\s]*[-–]?\s*',
        r'^In_Integration[_\s]*\d+[_\s]*[-–]?\s*',
        r'^Ra_Intro[_\s]*\d+[_\s]*[-–]?\s*',
        r'^Ra_Conf[_\s]*\d+[_\s]*[-–]?\s*',
        r'^PP[_\s]*\d+[_\s]*[-–]?\s*Configuration[_\s]*[-–]?\s*',
        r'^PP[_\s]*\d+[_\s]*[-–]?\s*',
        r'^\d+[_\s]*[-–]\s*',
    ]

    clean = name
    for pattern in patterns:
        clean = re.sub(pattern, '', clean, flags=re.IGNORECASE)

    # Clean up: remove leading/trailing special chars, "Introduction to" standardization
    clean = clean.strip(' -–_')
    clean = re.sub(r'^Introduction\s+to\s+', 'Introduction to ', clean, flags=re.IGNORECASE)

    return clean if clean else name


def get_lesson_number(folder_name, pptx_filename):
    """Extract lesson number from folder or filename."""
    # Try folder name first: IS_Fund_02, In_policy_05, In_Claim_01, BillingCenter_03
    patterns = [
        r'IS_Fund[_\s]*(\d+)',
        r'In_policy[_\s]*(\d+)',
        r'In_Claim[_\s]*(\d+)',
        r'BillingCenter[_\s]*(\d+)',
        r'In_Integration[_\s]*(\d+)',
        r'Ra_Intro[_\s]*(\d+)',
        r'Ra_Conf[_\s]*(\d+)',
        r'PP[_\s]*(\d+)',
        r'(\d+)\s*[-–]',
    ]

    for source in [folder_name, pptx_filename]:
        for pattern in patterns:
            m = re.search(pattern, source)
            if m:
                return int(m.group(1))

    return 0


def process_pptx(filepath, chapter_id, chapter_slug, lesson_folder, chapter_dir):
    """Process a single PPTX file and return structured lesson data."""
    filename = os.path.basename(filepath)
    print(f"  Processing: {filename}")

    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"  ERROR: Could not open {filename}: {e}")
        return None

    # Extract all slides
    slides_data = []
    for i, slide in enumerate(prs.slides, 1):
        slide_data = extract_slide_content(slide, i)
        slides_data.append(slide_data)

    # Get lesson metadata
    lesson_number = get_lesson_number(lesson_folder or "", filename)
    title = extract_lesson_title(filename)

    # If title is empty, try from first slide
    if not title and slides_data and slides_data[0]["title"]:
        title = slides_data[0]["title"]

    # Find associated videos
    videos = []
    if lesson_folder and chapter_dir:
        video_files = find_videos_for_lesson(chapter_dir, lesson_folder)
        for j, vf in enumerate(video_files):
            videos.append({
                "index": j + 1,
                "filename": vf,
                "path": f"/academy/guidewire/videos/{chapter_slug}/{vf}",
            })

    # Estimate duration based on slide count
    est_minutes = max(len(slides_data) * 3, 15)  # ~3 min per slide, minimum 15

    # Build lesson ID
    lesson_id = f"{chapter_slug}-l{lesson_number:02d}"

    return {
        "lessonId": lesson_id,
        "chapterId": chapter_id,
        "chapterSlug": chapter_slug,
        "lessonNumber": lesson_number,
        "title": title,
        "sourceFile": filename,
        "sourceFolder": lesson_folder,
        "totalSlides": len(slides_data),
        "estimatedMinutes": est_minutes,
        "slides": slides_data,
        "videos": videos,
        "hasNotes": any(s["notes"] for s in slides_data),
        "hasTables": any(s["hasTable"] for s in slides_data),
        "hasImages": any(s["hasImage"] for s in slides_data),
    }


def process_chapter(chapter_folder, chapter_config):
    """Process all PPTX files in a chapter directory."""
    chapter_dir = os.path.join(MATERIAL_DIR, chapter_folder)
    chapter_id = chapter_config["id"]
    chapter_slug = chapter_config["slug"]

    print(f"\n{'='*60}")
    print(f"Chapter {chapter_id}: {chapter_folder}")
    print(f"{'='*60}")

    lessons = []

    # Walk through subdirectories (each is a lesson folder)
    lesson_folders = []

    # Check for nested structure (e.g., Ch 8 has "Policy Center 10.0 Configuration/PP_XX")
    # and Ch 11 has "Introduction to Integration/In_Integration_XX"
    for item in sorted(os.listdir(chapter_dir)):
        item_path = os.path.join(chapter_dir, item)
        if os.path.isdir(item_path) and item not in ['.DS_Store', 'backups', 'videos']:
            # Check if this is a lesson folder (has PPT) or a container folder
            has_ppt = any(f.endswith(('.pptx', '.pptm')) for f in os.listdir(item_path) if os.path.isfile(os.path.join(item_path, f)))

            if has_ppt:
                lesson_folders.append((item, item_path))
            else:
                # Check subdirectories (nested structure)
                for sub_item in sorted(os.listdir(item_path)):
                    sub_path = os.path.join(item_path, sub_item)
                    if os.path.isdir(sub_path):
                        has_sub_ppt = any(f.endswith(('.pptx', '.pptm')) for f in os.listdir(sub_path) if os.path.isfile(os.path.join(sub_path, f)))
                        if has_sub_ppt:
                            lesson_folders.append((sub_item, sub_path))

    # Also check for PPTs directly in chapter root
    root_ppts = [f for f in os.listdir(chapter_dir) if f.endswith(('.pptx', '.pptm')) and os.path.isfile(os.path.join(chapter_dir, f))]

    # Process lesson folders
    for folder_name, folder_path in lesson_folders:
        pptx_files = [f for f in os.listdir(folder_path) if f.endswith(('.pptx', '.pptm')) and not f.startswith('~$')]

        for pptx in pptx_files:
            pptx_path = os.path.join(folder_path, pptx)
            lesson = process_pptx(pptx_path, chapter_id, chapter_slug, folder_name, os.path.dirname(folder_path))
            if lesson:
                lessons.append(lesson)

    # Process root-level PPTs
    for pptx in root_ppts:
        if pptx.startswith('~$'):
            continue
        pptx_path = os.path.join(chapter_dir, pptx)
        lesson = process_pptx(pptx_path, chapter_id, chapter_slug, None, chapter_dir)
        if lesson:
            lessons.append(lesson)

    # Sort by lesson number
    lessons.sort(key=lambda x: x["lessonNumber"])

    # Re-number if needed
    for i, lesson in enumerate(lessons):
        if lesson["lessonNumber"] == 0:
            lesson["lessonNumber"] = i + 1
            lesson["lessonId"] = f"{chapter_slug}-l{i+1:02d}"

    return lessons


def process_docx_chapter(chapter_folder, chapter_config):
    """Process DOCX files (for chapters like Advanced Product Designer)."""
    chapter_dir = os.path.join(MATERIAL_DIR, chapter_folder)
    chapter_id = chapter_config["id"]
    chapter_slug = chapter_config["slug"]

    try:
        from docx import Document
    except ImportError:
        print("  python-docx not available, skipping DOCX extraction")
        return []

    print(f"\n{'='*60}")
    print(f"Chapter {chapter_id}: {chapter_folder} (DOCX)")
    print(f"{'='*60}")

    lessons = []
    docx_files = []

    for f in sorted(os.listdir(chapter_dir)):
        if f.endswith('.docx') and not f.startswith('~$'):
            docx_files.append(f)

    # Find videos in videos subfolder
    video_dir = os.path.join(chapter_dir, "videos")
    all_videos = []
    if os.path.isdir(video_dir):
        all_videos = sorted([f for f in os.listdir(video_dir) if f.lower().endswith(('.mp4', '.mkv'))])

    for i, docx_file in enumerate(docx_files, 1):
        docx_path = os.path.join(chapter_dir, docx_file)
        print(f"  Processing: {docx_file}")

        try:
            doc = Document(docx_path)
        except Exception as e:
            print(f"  ERROR: {e}")
            continue

        # Extract paragraphs
        slides_data = []
        current_section = {"slideNumber": 1, "title": "", "bodyParagraphs": [], "notes": "",
                          "hasTable": False, "hasChart": False, "hasImage": False, "tableData": None, "mediaReferences": []}
        section_num = 1

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Headings become new sections
            if para.style and para.style.name and 'Heading' in para.style.name:
                if current_section["bodyParagraphs"] or current_section["title"]:
                    slides_data.append(current_section)
                    section_num += 1
                current_section = {"slideNumber": section_num, "title": text, "bodyParagraphs": [], "notes": "",
                                  "hasTable": False, "hasChart": False, "hasImage": False, "tableData": None, "mediaReferences": []}
            else:
                is_bold = any(run.bold for run in para.runs) if para.runs else False
                current_section["bodyParagraphs"].append({
                    "text": text,
                    "level": 0,
                    "bold": is_bold,
                })

        if current_section["bodyParagraphs"] or current_section["title"]:
            slides_data.append(current_section)

        # Title from filename
        title = os.path.splitext(docx_file)[0]
        title = re.sub(r'^\d+\s*[-–]\s*', '', title).strip()

        lesson_id = f"{chapter_slug}-l{i:02d}"

        lessons.append({
            "lessonId": lesson_id,
            "chapterId": chapter_id,
            "chapterSlug": chapter_slug,
            "lessonNumber": i,
            "title": title,
            "sourceFile": docx_file,
            "sourceFolder": None,
            "totalSlides": len(slides_data),
            "estimatedMinutes": max(len(slides_data) * 5, 20),
            "slides": slides_data,
            "videos": [{"index": j+1, "filename": v, "path": f"/academy/guidewire/videos/{chapter_slug}/{v}"} for j, v in enumerate(all_videos)],
            "hasNotes": False,
            "hasTables": False,
            "hasImages": False,
        })

    return lessons


def main():
    """Main extraction pipeline."""
    print("=" * 60)
    print("GUIDEWIRE ACADEMY - PPT CONTENT EXTRACTION")
    print("=" * 60)

    # Create output directory
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    all_chapters = {}
    total_lessons = 0
    total_slides = 0

    for chapter_folder, config in CHAPTER_MAP.items():
        chapter_dir = os.path.join(MATERIAL_DIR, chapter_folder)

        if not os.path.isdir(chapter_dir):
            print(f"\n  SKIP: {chapter_folder} (not found)")
            continue

        # Check if chapter has PPTX or DOCX
        has_pptx = False
        for root, dirs, files in os.walk(chapter_dir):
            if any(f.endswith(('.pptx', '.pptm')) and not f.startswith('~$') for f in files):
                has_pptx = True
                break

        has_docx = any(f.endswith('.docx') and not f.startswith('~$') for f in os.listdir(chapter_dir) if os.path.isfile(os.path.join(chapter_dir, f)))

        lessons = []
        if has_pptx:
            lessons = process_chapter(chapter_folder, config)
        elif has_docx:
            lessons = process_docx_chapter(chapter_folder, config)

        if lessons:
            # Save chapter data
            chapter_output_dir = os.path.join(OUTPUT_DIR, config["slug"])
            os.makedirs(chapter_output_dir, exist_ok=True)

            # Save individual lesson files
            for lesson in lessons:
                lesson_file = os.path.join(chapter_output_dir, f"lesson-{lesson['lessonNumber']:02d}.json")
                with open(lesson_file, 'w', encoding='utf-8') as f:
                    json.dump(lesson, f, indent=2, ensure_ascii=False)

            # Save chapter index
            chapter_index = {
                "chapterId": config["id"],
                "slug": config["slug"],
                "title": chapter_folder.replace(f"Chapter {config['id']} - ", ""),
                "totalLessons": len(lessons),
                "totalSlides": sum(l["totalSlides"] for l in lessons),
                "totalVideos": sum(len(l["videos"]) for l in lessons),
                "lessons": [{
                    "lessonId": l["lessonId"],
                    "lessonNumber": l["lessonNumber"],
                    "title": l["title"],
                    "totalSlides": l["totalSlides"],
                    "estimatedMinutes": l["estimatedMinutes"],
                    "videoCount": len(l["videos"]),
                    "hasNotes": l["hasNotes"],
                } for l in lessons]
            }

            index_file = os.path.join(chapter_output_dir, "index.json")
            with open(index_file, 'w', encoding='utf-8') as f:
                json.dump(chapter_index, f, indent=2, ensure_ascii=False)

            all_chapters[config["slug"]] = chapter_index
            total_lessons += len(lessons)
            total_slides += sum(l["totalSlides"] for l in lessons)

            print(f"  -> Extracted {len(lessons)} lessons, {sum(l['totalSlides'] for l in lessons)} slides")

    # Save master curriculum index
    master_index = {
        "program": "Guidewire Developer Training",
        "totalChapters": len(all_chapters),
        "totalLessons": total_lessons,
        "totalSlides": total_slides,
        "chapters": all_chapters,
    }

    master_file = os.path.join(OUTPUT_DIR, "curriculum-index.json")
    with open(master_file, 'w', encoding='utf-8') as f:
        json.dump(master_index, f, indent=2, ensure_ascii=False)

    print(f"\n{'='*60}")
    print(f"EXTRACTION COMPLETE")
    print(f"{'='*60}")
    print(f"Chapters: {len(all_chapters)}")
    print(f"Lessons:  {total_lessons}")
    print(f"Slides:   {total_slides}")
    print(f"Output:   {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
