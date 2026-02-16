#!/usr/bin/env python3
"""
Extract actual embedded images from PPTX slides for the academy lesson viewer.
Each slide's primary image is saved as a PNG in:
  public/academy/guidewire/slides/{chapterSlug}/lesson-{nn}/slide-{nn}.png

Falls back to LibreOffice PDF export if available, or composites embedded images.
"""

import os
import json
import sys
from io import BytesIO

try:
    from pptx import Presentation
    from pptx.util import Emu
    from PIL import Image
except ImportError:
    print("ERROR: python-pptx and Pillow required")
    print("  pip3 install python-pptx Pillow")
    sys.exit(1)

MATERIAL_DIR = "/Users/sumanthrajkumarnagolu/Library/CloudStorage/OneDrive-IntimeeSolutions/Intime E-Solutions/Trainings/Guidewire/Material"
CONTENT_DIR = "/Users/sumanthrajkumarnagolu/Projects/intime-v3/public/academy/guidewire/content"
SLIDES_DIR = "/Users/sumanthrajkumarnagolu/Projects/intime-v3/public/academy/guidewire/slides"

CHAPTER_MAP = {
    "Chapter 4 - PolicyCenter Introduction": {"id": 4, "slug": "ch04"},
    "Chapter 5 - Claim Center Introduction": {"id": 5, "slug": "ch05"},
    "Chapter 6 - Billing Center Introduction": {"id": 6, "slug": "ch06"},
    "Chapter 7 - InsuranceSuite Developer Fundamentals": {"id": 7, "slug": "ch07"},
    "Chapter 8 - policy center configuration": {"id": 8, "slug": "ch08"},
    "Chapter 9 - ClaimCenter Configuration": {"id": 9, "slug": "ch09"},
    "Chapter 11 - Introduction to Integration": {"id": 11, "slug": "ch11"},
    "Chapter 12 - Advanced product Designer": {"id": 12, "slug": "ch12"},
    "Chapter 13 - Rating Introduction": {"id": 13, "slug": "ch13"},
    "Chapter 14 - Rating Configuration": {"id": 14, "slug": "ch14"},
}

# Target image width
TARGET_WIDTH = 1280


def extract_slide_images(pptx_path, output_dir, force=False):
    """Extract embedded images from PPTX slides."""
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"    ERROR opening PPTX: {e}")
        return 0

    slide_width = prs.slide_width or Emu(12192000)
    slide_height = prs.slide_height or Emu(6858000)
    aspect = slide_height / slide_width
    img_width = TARGET_WIDTH
    img_height = int(img_width * aspect)

    os.makedirs(output_dir, exist_ok=True)
    count = 0

    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        img_path = os.path.join(output_dir, f"slide-{str(slide_num).zfill(2)}.png")

        if os.path.exists(img_path) and not force:
            count += 1
            continue

        # Collect all picture shapes sorted by area (largest first)
        pictures = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    img_blob = shape.image.blob
                    content_type = shape.image.content_type
                    area = (shape.width or 0) * (shape.height or 0)
                    coverage = area / (slide_width * slide_height) if slide_width and slide_height else 0
                    pictures.append({
                        'blob': img_blob,
                        'content_type': content_type,
                        'left': shape.left or 0,
                        'top': shape.top or 0,
                        'width': shape.width or 0,
                        'height': shape.height or 0,
                        'area': area,
                        'coverage': coverage,
                    })
                except Exception:
                    continue

        pictures.sort(key=lambda p: p['area'], reverse=True)

        if pictures:
            # Strategy 1: If one image covers >50% of the slide, use it directly
            main_pic = pictures[0]
            if main_pic['coverage'] > 0.5:
                try:
                    img = Image.open(BytesIO(main_pic['blob']))
                    img = img.convert('RGB')
                    # Resize to target width maintaining aspect
                    w, h = img.size
                    if w > 0:
                        new_h = int(TARGET_WIDTH * h / w)
                        img = img.resize((TARGET_WIDTH, new_h), Image.LANCZOS)
                    img.save(img_path, 'PNG', optimize=True)
                    count += 1
                    continue
                except Exception:
                    pass

            # Strategy 2: Composite multiple images onto a slide-sized canvas
            try:
                canvas = Image.new('RGB', (img_width, img_height), (255, 255, 255))

                for pic in pictures:
                    try:
                        pic_img = Image.open(BytesIO(pic['blob'])).convert('RGBA')

                        # Scale position and size to canvas coordinates
                        scale_x = img_width / slide_width
                        scale_y = img_height / slide_height
                        x = int(pic['left'] * scale_x)
                        y = int(pic['top'] * scale_y)
                        w = int(pic['width'] * scale_x)
                        h = int(pic['height'] * scale_y)

                        if w > 0 and h > 0:
                            pic_img = pic_img.resize((w, h), Image.LANCZOS)
                            canvas.paste(pic_img, (x, y), pic_img if pic_img.mode == 'RGBA' else None)
                    except Exception:
                        continue

                canvas.save(img_path, 'PNG', optimize=True)
                count += 1
                continue
            except Exception:
                pass

        # Strategy 3: No images - render text content
        try:
            from PIL import ImageDraw, ImageFont
            canvas = Image.new('RGB', (img_width, img_height), (250, 248, 244))  # cream bg
            draw = ImageDraw.Draw(canvas)

            try:
                title_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 28)
                body_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 18)
                small_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 13)
            except Exception:
                title_font = ImageFont.load_default()
                body_font = ImageFont.load_default()
                small_font = ImageFont.load_default()

            # Header bar
            draw.rectangle([0, 0, img_width, 50], fill=(23, 23, 23))
            draw.text((16, 14), f"Slide {slide_num}", fill=(201, 169, 97), font=small_font)

            y = 70
            has_title = False

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue

                    if not has_title:
                        has_title = True
                        draw.text((30, y), text[:90], fill=(23, 23, 23), font=title_font)
                        y += 45
                        draw.line([(30, y), (img_width - 30, y)], fill=(201, 169, 97), width=2)
                        y += 15
                    else:
                        level = para.level if para.level else 0
                        indent = 30 + level * 25
                        if level > 0:
                            draw.ellipse([indent - 10, y + 6, indent - 5, y + 11], fill=(120, 120, 120))

                        # Word wrap
                        words = text.split()
                        line = ""
                        max_w = img_width - indent - 30
                        for word in words:
                            test = f"{line} {word}".strip()
                            bbox = draw.textbbox((0, 0), test, font=body_font)
                            if bbox[2] - bbox[0] > max_w and line:
                                draw.text((indent, y), line, fill=(60, 60, 60), font=body_font)
                                y += 24
                                line = word
                            else:
                                line = test
                        if line:
                            draw.text((indent, y), line, fill=(60, 60, 60), font=body_font)
                            y += 24
                        y += 4

                    if y > img_height - 40:
                        break
                if y > img_height - 40:
                    break

            if y <= 80:
                draw.text((img_width // 2 - 80, img_height // 2 - 15),
                          f"Slide {slide_num}", fill=(170, 170, 170), font=title_font)

            canvas.save(img_path, 'PNG', optimize=True)
            count += 1
        except Exception as e:
            print(f"    ERROR rendering slide {slide_num}: {e}")

    return count


def main():
    print("=" * 60)
    print("SLIDE IMAGE EXTRACTION (Embedded Images)")
    print("=" * 60)

    total = 0

    for chapter_name, info in CHAPTER_MAP.items():
        slug = info["slug"]
        chapter_dir = os.path.join(MATERIAL_DIR, chapter_name)
        if not os.path.isdir(chapter_dir):
            continue

        index_path = os.path.join(CONTENT_DIR, slug, "index.json")
        if not os.path.exists(index_path):
            continue

        with open(index_path) as f:
            index_data = json.load(f)

        print(f"\n=== {chapter_name} ({slug}) ===")

        for lesson in index_data.get("lessons", []):
            lesson_num = lesson["lessonId"].split("-l")[1]

            lesson_json = os.path.join(CONTENT_DIR, slug, f"lesson-{lesson_num}.json")
            if not os.path.exists(lesson_json):
                continue

            with open(lesson_json) as f:
                lesson_data = json.load(f)

            source_file = lesson_data.get("sourceFile", "")
            source_folder = lesson_data.get("sourceFolder")

            # Find the PPTX file
            pptx_path = None
            if source_folder:
                pptx_path = os.path.join(chapter_dir, source_folder, source_file)
            if not pptx_path or not os.path.exists(pptx_path):
                for root, dirs, files in os.walk(chapter_dir):
                    if source_file in files:
                        pptx_path = os.path.join(root, source_file)
                        break

            if not pptx_path or not os.path.exists(pptx_path):
                continue

            out_dir = os.path.join(SLIDES_DIR, slug, f"lesson-{lesson_num}")
            print(f"  Lesson {lesson_num}: {source_file}")

            count = extract_slide_images(pptx_path, out_dir, force=True)
            total += count
            print(f"    -> {count} slides extracted")

    print(f"\n{'=' * 60}")
    print(f"Total: {total} slide images extracted")
    print(f"Output: {SLIDES_DIR}")


if __name__ == "__main__":
    main()
