#!/usr/bin/env python3
"""
Generate InTime Training Program Demo PPT
Explains how our training works - general methodology + Guidewire specific.
For internal team to understand how training is delivered.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ──────────────────────────────────────────────
# COLORS
# ──────────────────────────────────────────────
BLACK = RGBColor(0x17, 0x17, 0x17)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xC9, 0xA9, 0x61)
GOLD_DARK = RGBColor(0xA8, 0x8B, 0x4A)
CREAM = RGBColor(0xFD, 0xFB, 0xF7)
CHARCOAL_100 = RGBColor(0xF5, 0xF5, 0xF5)
CHARCOAL_200 = RGBColor(0xE5, 0xE5, 0xE5)
CHARCOAL_400 = RGBColor(0xA3, 0xA3, 0xA3)
CHARCOAL_500 = RGBColor(0x73, 0x73, 0x73)
CHARCOAL_600 = RGBColor(0x52, 0x52, 0x52)
CHARCOAL_700 = RGBColor(0x40, 0x40, 0x40)
GREEN = RGBColor(0x0A, 0x87, 0x54)
RED = RGBColor(0xDC, 0x26, 0x26)
BLUE = RGBColor(0x03, 0x69, 0xA1)
DARK_BG = RGBColor(0x0F, 0x0F, 0x0F)
AMBER = RGBColor(0xD9, 0x77, 0x06)
DARK_CARD = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_GREEN = RGBColor(0xF0, 0xF9, 0xF4)
LIGHT_BLUE = RGBColor(0xEF, 0xF6, 0xFF)
LIGHT_GOLD = RGBColor(0xFD, 0xF8, 0xED)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ──────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────

def add_gold_accent_line(slide, top=Inches(0)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left=0, top=top,
        width=prs.slide_width, height=Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=BLACK, spacing=Pt(6), font_name='Calibri', bullet_char=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet_char} {item}" if bullet_char else item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
    return txBox

def add_card(slide, left, top, width, height, bg_color=WHITE, border_color=CHARCOAL_200):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left=left, top=top, width=width, height=height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.05
    return shape

def add_divider(slide, left, top, width):
    d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left=left, top=top,
                                width=width, height=Inches(0.02))
    d.fill.solid()
    d.fill.fore_color.rgb = CHARCOAL_200
    d.line.fill.background()

def add_section_label(slide, left, top, text, color=CHARCOAL_400):
    add_text_box(slide, left, top, Inches(8), Inches(0.3),
                 text.upper(), font_size=10, color=color, bold=True)

def add_page_number(slide, num, total):
    add_text_box(slide, Inches(12.3), Inches(7.1), Inches(0.9), Inches(0.3),
                 f"{num}/{total}", font_size=9, color=CHARCOAL_400,
                 alignment=PP_ALIGN.RIGHT)

def add_numbered_step(slide, left, top, num, title, desc, width=Inches(2.6),
                      height=Inches(2.0), num_color=GOLD, bg=WHITE):
    add_card(slide, left, top, width, height, bg_color=bg)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1),
                 Inches(0.4), Inches(0.4), str(num), font_size=24, color=num_color, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.5),
                 width - Inches(0.3), Inches(0.3), title, font_size=12, color=BLACK, bold=True)
    add_divider(slide, left + Inches(0.15), top + Inches(0.82), width - Inches(0.3))
    add_text_box(slide, left + Inches(0.15), top + Inches(0.9),
                 width - Inches(0.3), height - Inches(1.0), desc, font_size=9, color=CHARCOAL_600)

TOTAL_SLIDES = 15

# ══════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gold_accent_line(slide)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
             "I N T I M E   A C A D E M Y", font_size=16, color=GOLD, bold=True)
add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
             "How Our Training\nProgram Works", font_size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.8),
             "The methodology, the structure, and the Guidewire program",
             font_size=22, color=CHARCOAL_400)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    left=Inches(1), top=Inches(5.2), width=Inches(2), height=Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = GOLD
shape.line.fill.background()

add_text_box(slide, Inches(1), Inches(5.5), Inches(7), Inches(0.4),
             "Internal Team Briefing  |  February 2026",
             font_size=12, color=CHARCOAL_500)
add_page_number(slide, 1, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 2: THE PROBLEM WITH TRADITIONAL TRAINING
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "WHY WE BUILT THIS")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "Traditional IT Training is Broken", font_size=36, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.8), Inches(9), Inches(0.3),
             "Every other training program has the same problems. We designed ours to fix all of them.",
             font_size=14, color=CHARCOAL_600)

problems = [
    ("Group Classes\nSkip Topics", "Trainer covers 60-70% of syllabus.\nRuns out of time, skips 'less important'\ntopics. Students have gaps.",
     "We cover 100%.\nEvery lesson. Every topic.\nNo exceptions."),
    ("Theory Heavy,\nNo Real Project", "Students learn concepts but never\nbuild anything real. Can't answer\n'tell me about your project' in interviews.",
     "Our labs ARE the project.\nEach assignment builds on previous.\nStudents build a full implementation."),
    ("One Pace\nFits None", "Fast learners bored. Slow learners lost.\nBusy US/Canada students miss sessions\nand fall behind permanently.",
     "Self-paced option for busy people.\nLive group option for structure.\nBoth complete 100% of material."),
    ("No Interview\nPreparation", "Students know Guidewire but can't\narticulate it in interviews. No resume\nbuilding, no mock interviews.",
     "Profile created. Custom project\nassigned. Interview questions after\nevery lesson. Mock interview prep."),
]

pw, ph, ptop = Inches(2.95), Inches(3.5), Inches(2.3)
for i, (problem, desc, solution) in enumerate(problems):
    left = Inches(0.6) + i * (pw + Inches(0.15))
    add_card(slide, left, ptop, pw, ph)
    # Problem
    add_text_box(slide, left + Inches(0.15), ptop + Inches(0.1),
                 pw - Inches(0.3), Inches(0.5), problem, font_size=12, color=RED, bold=True)
    add_text_box(slide, left + Inches(0.15), ptop + Inches(0.6),
                 pw - Inches(0.3), Inches(1.0), desc, font_size=9, color=CHARCOAL_600)
    # Solution
    add_divider(slide, left + Inches(0.15), ptop + Inches(1.7), pw - Inches(0.3))
    add_text_box(slide, left + Inches(0.15), ptop + Inches(1.75),
                 pw - Inches(0.3), Inches(0.2), "OUR FIX", font_size=8, color=GREEN, bold=True)
    add_text_box(slide, left + Inches(0.15), ptop + Inches(1.95),
                 pw - Inches(0.3), Inches(1.3), solution, font_size=9, color=GREEN)

add_text_box(slide, Inches(1), Inches(6.15), Inches(11), Inches(0.4),
             "Bottom line: Our students finish knowing MORE than people with 1-2 years of real experience,\n"
             "because they've touched every part of the system systematically.",
             font_size=12, color=CHARCOAL_500, bold=True)
add_page_number(slide, 2, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 3: OUR APPROACH - THE INTIME METHOD
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "THE INTIME METHOD", color=CHARCOAL_500)
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "100+ Sequential Lessons.\nOne Full Project.", font_size=40, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(2.2), Inches(9), Inches(0.4),
             "We broke down the entire technology into sequential lessons. Each lesson has 4 components. "
             "All assignments connect into one full project.",
             font_size=15, color=CHARCOAL_400)

# 4 components of each lesson
components = [
    ("1", "EXPLANATION\nPPT", "Detailed slide deck explaining\nthe concept with diagrams,\narchitecture, and examples.\n\nStudents read and understand\nthe theory first.", GOLD),
    ("2", "DEMO\nVIDEO", "Embedded video showing\nexactly how it works.\nStep-by-step walkthrough.\n\nStudents see it done before\nthey try it themselves.", BLUE),
    ("3", "INTERVIEW\nQUESTIONS", "Real interview questions for\nthis specific topic.\nStudents practice answering.\n\nBy the end, they've covered\n500+ interview questions.", GREEN),
    ("4", "HANDS-ON\nASSIGNMENT", "Practical lab exercise.\nBuild that feature in\nthe project environment.\n\nEach assignment is a user story\nin the full project.", AMBER),
]

cw, ch, ctop = Inches(2.7), Inches(3.3), Inches(3.0)
for i, (num, title, desc, accent) in enumerate(components):
    left = Inches(0.7) + i * (cw + Inches(0.25))
    add_card(slide, left, ctop, cw, ch, bg_color=DARK_CARD, border_color=CHARCOAL_700)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        left=left + Inches(0.15), top=ctop + Inches(0.15),
        width=Inches(0.45), height=Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    add_text_box(slide, left + Inches(0.15), ctop + Inches(0.15),
                 Inches(0.45), Inches(0.45), num, font_size=18, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.7), ctop + Inches(0.15),
                 cw - Inches(0.85), Inches(0.5), title, font_size=12, color=accent, bold=True)
    add_divider(slide, left + Inches(0.15), ctop + Inches(0.72), cw - Inches(0.3))
    add_text_box(slide, left + Inches(0.15), ctop + Inches(0.8),
                 cw - Inches(0.3), Inches(2.3), desc, font_size=10, color=CHARCOAL_400)

add_text_box(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.4),
             "Every single lesson follows this pattern. No lesson is theory-only. No lesson skips interview prep.\n"
             "By the time a student finishes, they've built a complete project AND prepared for every possible interview question.",
             font_size=12, color=CHARCOAL_500, bold=True)
add_page_number(slide, 3, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 4: THE FULL PROJECT APPROACH
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "PROJECT-BASED LEARNING")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "Labs = Full Project. Modules = Epics. Lessons = User Stories.",
             font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(10), Inches(0.4),
             "Think of it like building software with Agile. The entire training IS one implementation project. "
             "Each module is an epic. Each lesson's assignment is a user story that adds to the project.",
             font_size=13, color=CHARCOAL_600)

# Visual: Project breakdown
add_card(slide, Inches(0.6), Inches(2.3), Inches(11.8), Inches(1.2), bg_color=LIGHT_GOLD)
add_text_box(slide, Inches(0.8), Inches(2.35), Inches(11), Inches(0.25),
             "THE FULL PROJECT (Example: Insurance Company Implementation)", font_size=11, color=GOLD_DARK, bold=True)
add_divider(slide, Inches(0.8), Inches(2.6), Inches(11.4))
add_text_box(slide, Inches(0.8), Inches(2.65), Inches(11.4), Inches(0.8),
             "Student builds a complete insurance application from scratch. Every assignment adds a real feature. "
             "By Day 40, they have a working system with policies, claims, users, rules, integrations - everything.\n"
             "When an interviewer asks 'Tell me about your project,' they have a REAL answer with REAL implementation details.",
             font_size=11, color=CHARCOAL_600)

# Module = Epic breakdown
epics = [
    ("EPIC 1", "PolicyCenter\nBusiness", "Days 1-9\n~18 lessons", "Set up accounts, create policies,\nsubmissions, renewals, cancellations,\nusers, roles, product model, rating", GREEN),
    ("EPIC 2", "ClaimCenter\nBusiness", "Days 10-13\n~16 lessons", "Claims process, intake, setup,\nadjudication, financials, payments,\ncontacts, vendors, service requests", BLUE),
    ("EPIC 3", "Developer\nFundamentals", "Days 14-20\n~25 lessons", "Configuration, data model, Gosu,\nUI architecture, widgets, views,\nentities, typelists, popups", AMBER),
    ("EPIC 4", "ClaimCenter\nConfiguration", "Days 21-26\n~17 lessons", "UI config, Gosu rules, assignment\nrules, financials config, search,\npermissions, activities", GREEN),
    ("EPIC 5", "PolicyCenter\nConfiguration", "Days 27-30\n~14 lessons", "Data model, wizards, revisioning,\nunderwriting issues, permissions,\nactivities, job lifecycle", BLUE),
    ("EPIC 6", "Integration", "Days 31-40\n~21 lessons", "REST/SOAP APIs, plugins, messaging,\nbundles, queries, templates,\nbatch processes, auth", GOLD),
]

# 3 columns x 2 rows
ew, eh = Inches(3.8), Inches(1.6)
for i, (epic, name, days, desc, accent) in enumerate(epics):
    col = i % 3
    row = i // 3
    left = Inches(0.6) + col * (ew + Inches(0.15))
    top = Inches(3.7) + row * (eh + Inches(0.15))
    add_card(slide, left, top, ew, eh)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.08),
                 Inches(0.8), Inches(0.2), epic, font_size=9, color=accent, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.3),
                 Inches(1.3), Inches(0.5), name, font_size=12, color=BLACK, bold=True)
    add_text_box(slide, left + Inches(1.5), top + Inches(0.08),
                 Inches(1), Inches(0.35), days, font_size=8, color=CHARCOAL_400)
    add_divider(slide, left + Inches(0.15), top + Inches(0.82), ew - Inches(0.3))
    add_text_box(slide, left + Inches(0.15), top + Inches(0.88),
                 ew - Inches(0.3), Inches(0.65), desc, font_size=8, color=CHARCOAL_600)

add_text_box(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.3),
             "Doing assignments itself IS the huge project. Students don't realize they've built a full system until it's done.",
             font_size=12, color=CHARCOAL_500, bold=True)
add_page_number(slide, 4, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 5: TWO DELIVERY MODELS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "DELIVERY MODELS")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "Two Ways to Deliver the Same Program", font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(9), Inches(0.3),
             "Same 100+ lessons. Same assignments. Same project. Different pacing to suit different students.",
             font_size=14, color=CHARCOAL_600)

# Both models - dedicated coordinator
add_card(slide, Inches(0.6), Inches(2.2), Inches(11.8), Inches(0.8), bg_color=LIGHT_GOLD)
add_text_box(slide, Inches(0.8), Inches(2.25), Inches(11.4), Inches(0.2),
             "BOTH MODELS: DEDICATED COORDINATOR", font_size=10, color=GOLD_DARK, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.5), Inches(11.4), Inches(0.4),
             "Every student is assigned a dedicated coordinator on Day 1. They handle scheduling, progress tracking, "
             "trainer assignments, and are the student's single point of contact throughout the program.",
             font_size=11, color=CHARCOAL_600)

# Model A
add_card(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.8), bg_color=LIGHT_BLUE)
add_text_box(slide, Inches(0.8), Inches(3.3), Inches(1), Inches(0.35),
             "MODEL A", font_size=14, color=BLUE, bold=True)
add_text_box(slide, Inches(2.0), Inches(3.3), Inches(4), Inches(0.35),
             "Self-Paced + 20 Hrs Live Trainer", font_size=18, color=BLACK, bold=True)
add_divider(slide, Inches(0.8), Inches(3.7), Inches(5.4))

add_text_box(slide, Inches(0.8), Inches(3.8), Inches(5.4), Inches(0.2),
             "DESIGNED FOR", font_size=9, color=BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(4.0), Inches(5.4), Inches(0.3),
             "Busy US/Canada resources who can't attend fixed-time classes. Working professionals.",
             font_size=11, color=CHARCOAL_600)

add_text_box(slide, Inches(0.8), Inches(4.35), Inches(5.4), Inches(0.2),
             "HOW IT WORKS", font_size=9, color=BLUE, bold=True)
model_a_steps = [
    "Student gets access to all 100+ lessons on our platform",
    "Goes through lessons at their own pace (PPT → Video → Q&A → Assignment)",
    "Assignments build on each other = full project grows",
    "Coordinator tracks progress and answers logistics questions",
    "After completing self-paced material:",
    "   → 20 hours of live 1-on-1 trainer sessions (2 hrs/day, ~2 weeks)",
    "   → Blocker resolution, project guidance, interview prep",
    "   → Mock interviews and profile finalization",
]
add_bullet_list(slide, Inches(0.8), Inches(4.55), Inches(5.4), Inches(2.0),
                model_a_steps, font_size=10, color=CHARCOAL_600, bullet_char="\u2022", spacing=Pt(2))

add_text_box(slide, Inches(0.8), Inches(6.35), Inches(5.4), Inches(0.5),
             "KEY ADVANTAGE: 100% material covered at own pace.\n20 hrs live trainer ensures mastery + interview readiness.",
             font_size=10, color=BLUE, bold=True)

# Model B
add_card(slide, Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.8), bg_color=LIGHT_GREEN)
add_text_box(slide, Inches(7.0), Inches(3.3), Inches(1), Inches(0.35),
             "MODEL B", font_size=14, color=GREEN, bold=True)
add_text_box(slide, Inches(8.2), Inches(3.3), Inches(4), Inches(0.35),
             "80 Hrs Live Group (8 Weeks)", font_size=18, color=BLACK, bold=True)
add_divider(slide, Inches(7.0), Inches(3.7), Inches(5.4))

add_text_box(slide, Inches(7.0), Inches(3.8), Inches(5.4), Inches(0.2),
             "DESIGNED FOR", font_size=9, color=GREEN, bold=True)
add_text_box(slide, Inches(7.0), Inches(4.0), Inches(5.4), Inches(0.3),
             "Students who need structure, accountability, and daily live interaction with a trainer.",
             font_size=11, color=CHARCOAL_600)

add_text_box(slide, Inches(7.0), Inches(4.35), Inches(5.4), Inches(0.2),
             "HOW IT WORKS", font_size=9, color=GREEN, bold=True)
model_b_steps = [
    "80 hours total: 2 hours/day live sessions over 8 weeks",
    "Profile created Day 1, custom project assigned Day 1",
    "Trainer assigns lessons as homework after each session",
    "Students study material + complete assignments before next session",
    "Come to live session with questions and blockers",
    "Trainer discusses, unblocks, guides project work",
    "Must stay in sync — completing homework keeps pace with group",
    "Interview prep integrated throughout the 8 weeks",
]
add_bullet_list(slide, Inches(7.0), Inches(4.55), Inches(5.4), Inches(2.0),
                model_b_steps, font_size=10, color=CHARCOAL_600, bullet_char="\u2022", spacing=Pt(2))

add_text_box(slide, Inches(7.0), Inches(6.35), Inches(5.4), Inches(0.5),
             "KEY ADVANTAGE: Daily accountability + live interaction.\n8 weeks of structured trainer-led guidance.",
             font_size=10, color=GREEN, bold=True)

add_page_number(slide, 5, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 6: MODEL A DEEP DIVE - SELF-PACED
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "MODEL A: SELF-PACED + 20 HRS LIVE", color=CHARCOAL_500)
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "Self-Paced Learning → 20 Hours Live Trainer", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.7), Inches(9), Inches(0.3),
             "Student completes all material at their own pace, then gets 2 weeks of dedicated 1-on-1 trainer time (2 hrs/day).",
             font_size=14, color=CHARCOAL_400)

steps_a = [
    ("1", "ENROLL &\nCOORDINATOR", "Student enrolls.\nDedicated coordinator\nassigned. Platform\naccess to all 100+\nlessons.",
     "Day 1"),
    ("2", "SELF-PACED\nLEARNING", "Goes through lessons\nsequentially at own\npace. PPT → Video →\nInterview Q's →\nAssignment.",
     "Own Pace"),
    ("3", "BUILD\nPROJECT", "Completes assignment\nfor each lesson.\nAssignments connect\ninto one full project\nthat grows daily.",
     "Own Pace"),
    ("4", "COMPLETE\nMATERIAL", "Finishes all 100+\nlessons and all\nassignments.\nCoordinator tracks\nprogress throughout.",
     "Own Pace"),
    ("5", "20 HRS LIVE\nTRAINER", "2 hrs/day for ~2 weeks.\n1-on-1 with trainer.\nBlocker resolution.\nProject deep-dive.\nQ&A on any topic.",
     "~2 Weeks"),
    ("6", "INTERVIEW\nPREP", "Mock interviews with\ntrainer during live\nsessions. Resume review.\nProject walkthrough\npractice.",
     "During Live"),
    ("7", "PROFILE &\nPLACEMENT", "Complete profile ready.\nProject portfolio done.\nResume finalized.\nHandoff to placement\nteam.",
     "Completion"),
]

sw, sh, stop = Inches(1.6), Inches(3.0), Inches(2.3)
for i, (num, title, desc, timing) in enumerate(steps_a):
    left = Inches(0.4) + i * (sw + Inches(0.12))
    add_card(slide, left, stop, sw, sh, bg_color=DARK_CARD, border_color=CHARCOAL_700)
    # Number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        left=left + Inches(0.1), top=stop + Inches(0.1),
        width=Inches(0.35), height=Inches(0.35))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    add_text_box(slide, left + Inches(0.1), stop + Inches(0.1),
                 Inches(0.35), Inches(0.35), num, font_size=14, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.5), stop + Inches(0.12),
                 sw - Inches(0.6), Inches(0.45), title, font_size=9, color=GOLD, bold=True)
    add_divider(slide, left + Inches(0.1), stop + Inches(0.6), sw - Inches(0.2))
    add_text_box(slide, left + Inches(0.1), stop + Inches(0.68),
                 sw - Inches(0.2), Inches(1.5), desc, font_size=8, color=CHARCOAL_400)
    add_text_box(slide, left + Inches(0.1), stop + sh - Inches(0.35),
                 sw - Inches(0.2), Inches(0.25), timing, font_size=8, color=CHARCOAL_500, bold=True)

    # Arrow between steps
    if i < len(steps_a) - 1:
        arrow_left = left + sw + Inches(0.02)
        add_text_box(slide, arrow_left, stop + Inches(1.3),
                     Inches(0.1), Inches(0.2), "\u2192", font_size=10, color=GOLD)

add_text_box(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.8),
             "Why this works for US/Canada students:\n"
             "\u2022 Work during the day, study evenings/weekends  |  "
             "\u2022 No missed sessions — they set the schedule  |  "
             "\u2022 100% material covered, nothing skipped  |  "
             "\u2022 20 hrs live trainer at the end for mastery + interview prep",
             font_size=11, color=CHARCOAL_400)

add_text_box(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.4),
             "Self-paced ensures 100% coverage. Live trainer at the end ensures mastery, project quality, and interview readiness.",
             font_size=12, color=GOLD, bold=True)
add_page_number(slide, 6, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 7: MODEL B DEEP DIVE - LIVE GROUP
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "MODEL B: LIVE GROUP (80 HRS / 8 WEEKS)")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "80 Hours Live: 2 Hours/Day Over 8 Weeks", font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(10), Inches(0.3),
             "Profile Day 1. Project Day 1. 2-hr live session daily. Complete homework assignments before next session to stay in sync.",
             font_size=14, color=CHARCOAL_600)

# Daily rhythm
add_card(slide, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.5), bg_color=LIGHT_GREEN)
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(5), Inches(0.3),
             "DAILY RHYTHM", font_size=12, color=GREEN, bold=True)
add_divider(slide, Inches(0.8), Inches(2.6), Inches(5.4))

rhythm = [
    ("LIVE SESSION", "2-hour live session with trainer (daily)"),
    ("DISCUSSION", "Trainer covers topics, answers questions, unblocks"),
    ("PROJECT WORK", "Trainer guides project work at appropriate moments"),
    ("HOMEWORK", "Trainer assigns lessons to study + assignments to complete"),
    ("STUDENT TIME", "Study PPTs, watch videos, complete assignments"),
    ("NEXT SESSION", "Come prepared with completed homework, ready for next topics"),
]
for j, (phase, desc) in enumerate(rhythm):
    y = Inches(2.7) + j * Inches(0.3)
    add_text_box(slide, Inches(0.8), y, Inches(1.6), Inches(0.25),
                 phase, font_size=9, color=GREEN, bold=True)
    add_text_box(slide, Inches(2.5), y, Inches(3.5), Inches(0.25),
                 desc, font_size=9, color=CHARCOAL_600)

# 8-week overview
add_card(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.5))
add_text_box(slide, Inches(7.0), Inches(2.3), Inches(5), Inches(0.3),
             "8-WEEK OVERVIEW", font_size=12, color=BLACK, bold=True)
add_divider(slide, Inches(7.0), Inches(2.6), Inches(5.4))

week_plan = [
    ("Week 1-2", "Profile + project setup. Insurance intro. PolicyCenter business."),
    ("Week 3", "ClaimCenter business. Claims, financials, vendors."),
    ("Week 4-5", "Developer fundamentals. Config, data model, Gosu, UI."),
    ("Week 6", "ClaimCenter configuration. Rules, permissions, activities."),
    ("Week 7", "PolicyCenter configuration. Wizards, revisioning, lifecycle."),
    ("Week 8", "Integration + final project review + interview prep."),
]
for j, (week, desc) in enumerate(week_plan):
    y = Inches(2.7) + j * Inches(0.33)
    add_text_box(slide, Inches(7.0), y, Inches(1.2), Inches(0.25),
                 week, font_size=9, color=BLACK, bold=True)
    add_text_box(slide, Inches(8.3), y, Inches(4.1), Inches(0.25),
                 desc, font_size=9, color=CHARCOAL_600)

# What makes it different
add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.3),
             "WHAT MAKES OUR LIVE GROUP DIFFERENT FROM OTHERS", font_size=12, color=BLACK, bold=True)

diffs = [
    ("Profile Day 1", "We don't wait until the end.\nProfile created Day 1.\nProject assigned Day 1.\nStudents know what\nthey're building for.", GREEN),
    ("Homework Model", "Lessons are homework.\nStudents must complete\nbefore next session.\nLive time is for Q&A,\nnot passive lectures.", BLUE),
    ("8 Weeks = Depth", "Not a rushed 2-week bootcamp.\n80 hours of live time.\nStudents absorb, practice,\nand build properly.\nNo shortcuts.", AMBER),
    ("Stay In Sync", "Everyone completes the\nsame homework daily.\nGroup moves together.\nTrainer knows exactly\nwhere everyone is.", GOLD),
]

dw, dh, dtop = Inches(2.85), Inches(1.55), Inches(5.4)
for i, (title, desc, accent) in enumerate(diffs):
    left = Inches(0.6) + i * (dw + Inches(0.15))
    add_card(slide, left, dtop, dw, dh)
    add_text_box(slide, left + Inches(0.15), dtop + Inches(0.08),
                 dw - Inches(0.3), Inches(0.2), title, font_size=10, color=accent, bold=True)
    add_divider(slide, left + Inches(0.15), dtop + Inches(0.3), dw - Inches(0.3))
    add_text_box(slide, left + Inches(0.15), dtop + Inches(0.38),
                 dw - Inches(0.3), Inches(1.1), desc, font_size=8, color=CHARCOAL_600)

add_page_number(slide, 7, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 8: THE STUDENT JOURNEY
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "THE STUDENT JOURNEY")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "From Enrollment to Placement", font_size=32, color=BLACK, bold=True)

journey = [
    ("ENROLL", "Student signs up for\nthe program.\nAssessed for level.\nPlatform access given.",
     "OPT Recruiter / Training Sales\nidentifies the lead, qualifies,\nand enrolls them.", BLUE),
    ("ONBOARD", "Dedicated coordinator\nassigned. Profile created\nDay 1. Project assigned\nbased on resume.",
     "Coordinator handles logistics,\nschedules, progress tracking.\nTrainer sets up project env.", GREEN),
    ("LEARN &\nBUILD", "Self-paced: own schedule.\nLive group: 2 hrs/day.\nCompletes all assignments.\n100% coverage guaranteed.",
     "Coordinator tracks progress.\nTrainer (1-on-1 or group)\nguides and reviews.", AMBER),
    ("INTERVIEW\nPREP", "500+ interview questions.\nMock interviews with trainer.\nProject walkthrough practice.\nResume finalized.",
     "Trainer conducts mock\ninterviews. Profile is\npolished and ready.", GOLD),
    ("PLACEMENT", "Added to our talent pool.\nTop performers: W2 bench.\nOthers: matched to jobs.\nInterview-ready.",
     "Account Manager / Bench Sales\ntake over. Submit to\nclient jobs immediately.", RED),
]

jw, jh = Inches(2.3), Inches(4.0)
for i, (title, student_view, team_view, accent) in enumerate(journey):
    left = Inches(0.5) + i * (jw + Inches(0.15))
    add_card(slide, left, Inches(1.8), jw, jh)
    # Phase title
    add_text_box(slide, left + Inches(0.15), Inches(1.85),
                 jw - Inches(0.3), Inches(0.4), title, font_size=13, color=accent, bold=True)
    add_divider(slide, left + Inches(0.15), Inches(2.3), jw - Inches(0.3))
    # Student view
    add_text_box(slide, left + Inches(0.15), Inches(2.35),
                 jw - Inches(0.3), Inches(0.15), "STUDENT", font_size=7, color=CHARCOAL_400, bold=True)
    add_text_box(slide, left + Inches(0.15), Inches(2.5),
                 jw - Inches(0.3), Inches(1.1), student_view, font_size=9, color=CHARCOAL_600)
    # Team view
    add_divider(slide, left + Inches(0.15), Inches(3.7), jw - Inches(0.3))
    add_text_box(slide, left + Inches(0.15), Inches(3.75),
                 jw - Inches(0.3), Inches(0.15), "OUR TEAM", font_size=7, color=CHARCOAL_400, bold=True)
    add_text_box(slide, left + Inches(0.15), Inches(3.9),
                 jw - Inches(0.3), Inches(0.8), team_view, font_size=8, color=CHARCOAL_500)

    # Arrow
    if i < len(journey) - 1:
        add_text_box(slide, left + jw + Inches(0.02), Inches(2.8),
                     Inches(0.12), Inches(0.2), "\u2192", font_size=12, color=GOLD, bold=True)

add_text_box(slide, Inches(1), Inches(6.1), Inches(11), Inches(0.8),
             "The handoff is seamless. Training Sales → Dedicated Coordinator → Trainer → Account Manager/Bench Sales.\n"
             "Coordinator stays with the student throughout. No one falls through the cracks.",
             font_size=12, color=CHARCOAL_500, bold=True)
add_page_number(slide, 8, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 9: NOW - GUIDEWIRE SPECIFIC
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gold_accent_line(slide)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(0.5),
             "I N T I M E   A C A D E M Y", font_size=14, color=GOLD, bold=True)
add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(1.5),
             "Guidewire InsuranceSuite\nTraining Program", font_size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(4.6), Inches(8), Inches(0.8),
             "40 Days  |  121 Lessons  |  119 Assignments  |  6 Modules\nPolicyCenter + ClaimCenter + Integration",
             font_size=20, color=CHARCOAL_400)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    left=Inches(1), top=Inches(5.8), width=Inches(2), height=Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = GOLD
shape.line.fill.background()

add_text_box(slide, Inches(1), Inches(6.1), Inches(11), Inches(0.4),
             "Now let's look at the specific Guidewire curriculum...",
             font_size=14, color=CHARCOAL_500)
add_page_number(slide, 9, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 10: GW CURRICULUM - BUSINESS MODULES (Days 1-13)
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "GUIDEWIRE CURRICULUM")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "Days 1-13: Business Fundamentals", font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(10), Inches(0.3),
             "Insurance basics, PolicyCenter business operations, and ClaimCenter business operations. "
             "Students learn what the system DOES before learning how to configure it.",
             font_size=13, color=CHARCOAL_600)

# Days 1-9: PolicyCenter
add_card(slide, Inches(0.6), Inches(2.2), Inches(7.5), Inches(4.6))
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(5), Inches(0.3),
             "DAYS 1-9: POLICYCENTER BUSINESS (~34 lessons)", font_size=12, color=GREEN, bold=True)
add_divider(slide, Inches(0.8), Inches(2.6), Inches(7.1))

pc_days = [
    ("Day 1", "Dev Fundamentals, Cloud Overview, Project Resources, User Story Cards"),
    ("Day 2", "PolicyCenter Intro, Accounts, Policy Transactions, Policy File, Product Model"),
    ("Day 3", "Full Application Submission, Policy Tools, Changes & Preemptions, Renewals"),
    ("Day 4", "Cancellation/Reinstatement/Rewrite, Out of Sequence, Users & Groups, Orgs"),
    ("Day 5", "Forms, Underwriting Authority (Intro + Managing), Products & Policy Lines"),
    ("Day 6", "Coverages, Coverage Terms, Product Model Availability, Question Sets"),
    ("Day 7", "Modifiers, Contingencies, Policy Holds & UW Referrals, Documents"),
    ("Day 8", "Notes, Activities, Roles & Permissions, Validation"),
    ("Day 9", "Premium Audits, Rating Basics"),
]

for j, (day, topics) in enumerate(pc_days):
    y = Inches(2.7) + j * Inches(0.42)
    add_text_box(slide, Inches(0.8), y, Inches(0.8), Inches(0.3),
                 day, font_size=9, color=GREEN, bold=True)
    add_text_box(slide, Inches(1.6), y, Inches(6.3), Inches(0.35),
                 topics, font_size=8, color=CHARCOAL_600)

# Days 10-13: ClaimCenter
add_card(slide, Inches(8.4), Inches(2.2), Inches(4.4), Inches(4.6))
add_text_box(slide, Inches(8.6), Inches(2.3), Inches(4), Inches(0.3),
             "DAYS 10-13: CLAIMCENTER\nBUSINESS (~16 lessons)", font_size=11, color=BLUE, bold=True)
add_divider(slide, Inches(8.6), Inches(2.7), Inches(4.0))

cc_days = [
    ("Day 10", "ClaimCenter Intro,\nClaims Process,\nClaim Maintenance,\nOrg Structure,\nLoB Coverage"),
    ("Day 11", "Claim Intake,\nClaim Setup,\nData Model,\nAdjudication"),
    ("Day 12", "Financial Terms,\nCreating Payments,\nPayment Approvals,\nFinancial Holds"),
    ("Day 13", "Managing Contacts,\nManaging Vendors,\nService Requests,\nPermissions & ACLs"),
]

for j, (day, topics) in enumerate(cc_days):
    y = Inches(2.8) + j * Inches(0.95)
    add_text_box(slide, Inches(8.6), y, Inches(0.8), Inches(0.2),
                 day, font_size=9, color=BLUE, bold=True)
    add_text_box(slide, Inches(8.6), y + Inches(0.2), Inches(4), Inches(0.7),
                 topics, font_size=8, color=CHARCOAL_600)

add_text_box(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.3),
             "By Day 13, students understand the full insurance lifecycle: policies, claims, payments, underwriting, and compliance.",
             font_size=12, color=CHARCOAL_500, bold=True)
add_page_number(slide, 10, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 11: GW CURRICULUM - DEVELOPER (Days 14-20)
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "GUIDEWIRE CURRICULUM")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "Days 14-20: Developer Fundamentals", font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(10), Inches(0.3),
             "From understanding the business to configuring the system. Students learn the development toolkit.",
             font_size=14, color=CHARCOAL_600)

dev_days = [
    ("Day 14", "InsuranceSuite Developer Fundamentals",
     ["Introduction to Guidewire Configuration", "Introduction to the Data Model", "Extending the Data Model"],
     "Students go from user to developer. Learn how config works under the hood.", AMBER),
    ("Day 15", "UI Architecture",
     ["The User Interface Architecture", "Atomic Widgets", "Detail Views"],
     "Learn the PCF framework. How screens are built. Widget types and layouts.", AMBER),
    ("Day 16", "Gosu & Rules",
     ["Introduction to Locations", "Introduction to Gosu", "Gosu Rules"],
     "The Guidewire programming language. Locations, rules, and scripting.", AMBER),
    ("Day 17", "Advanced Development",
     ["Enhancements", "Code Generation & Debugging", "Creating New Entities"],
     "Extend the platform. Debug issues. Create custom data structures.", AMBER),
    ("Day 18", "List Views & Types",
     ["List Views", "Editable List Views", "Typelists"],
     "Build data grids. Editable tables. Custom type definitions.", AMBER),
    ("Day 19", "UI Components",
     ["Popups", "Validation", "Input Sets"],
     "Modal dialogs. Form validation. Reusable input components.", AMBER),
    ("Day 20", "Advanced UI",
     ["Partial Page Update", "Subtypes", "Modes", "Entity Names"],
     "AJAX updates. Polymorphic entities. Display modes. Naming conventions.", AMBER),
]

for j, (day, title, lessons, note, accent) in enumerate(dev_days):
    y = Inches(2.1) + j * Inches(0.72)
    add_card(slide, Inches(0.6), y, Inches(11.8), Inches(0.65))
    add_text_box(slide, Inches(0.8), y + Inches(0.05),
                 Inches(0.8), Inches(0.25), day, font_size=11, color=accent, bold=True)
    add_text_box(slide, Inches(1.6), y + Inches(0.05),
                 Inches(2.5), Inches(0.25), title, font_size=11, color=BLACK, bold=True)
    lessons_text = "  |  ".join(lessons)
    add_text_box(slide, Inches(4.2), y + Inches(0.05),
                 Inches(4.5), Inches(0.25), lessons_text, font_size=8, color=CHARCOAL_600)
    add_text_box(slide, Inches(8.8), y + Inches(0.05),
                 Inches(3.4), Inches(0.5), note, font_size=8, color=CHARCOAL_500)

add_text_box(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.3),
             "After Day 20, students can build custom UI, write Gosu code, extend data models, and debug Guidewire applications.",
             font_size=12, color=CHARCOAL_500, bold=True)
add_page_number(slide, 11, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 12: GW CURRICULUM - CONFIG + INTEGRATION (Days 21-40)
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "GUIDEWIRE CURRICULUM")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "Days 21-40: Configuration & Integration", font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(10), Inches(0.3),
             "Apply developer skills to configure both ClaimCenter and PolicyCenter, then build integrations.",
             font_size=14, color=CHARCOAL_600)

# CC Config (Days 21-26)
add_card(slide, Inches(0.6), Inches(2.2), Inches(3.8), Inches(3.0), bg_color=LIGHT_BLUE)
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(3.4), Inches(0.3),
             "DAYS 21-26\nClaimCenter Config", font_size=12, color=BLUE, bold=True)
add_divider(slide, Inches(0.8), Inches(2.75), Inches(3.4))
cc_config = [
    "Day 21: UI Config, LoB, Claim Intake Config",
    "Day 22: Gosu Rules, Assignment Rules, Validation",
    "Day 23: Claim Setup Rules, Contacts, Financials",
    "Day 24: Financial Holds, Approvals, Vendor Services",
    "Day 25: Search, Claim History, Permissions",
    "Day 26: Activity Rules, Exposure Rules, Archiving",
]
add_bullet_list(slide, Inches(0.8), Inches(2.85), Inches(3.4), Inches(2.2),
                cc_config, font_size=9, color=CHARCOAL_600, bullet_char="\u2022", spacing=Pt(3))

# PC Config (Days 27-30)
add_card(slide, Inches(4.6), Inches(2.2), Inches(3.8), Inches(3.0), bg_color=LIGHT_GREEN)
add_text_box(slide, Inches(4.8), Inches(2.3), Inches(3.4), Inches(0.3),
             "DAYS 27-30\nPolicyCenter Config", font_size=12, color=GREEN, bold=True)
add_divider(slide, Inches(4.8), Inches(2.75), Inches(3.4))
pc_config = [
    "Day 27: Data Model, Location Groups, Wizards",
    "Day 28: Contacts, Revisioning, UW Issues",
    "Day 29: Approving UW Issues, Validation Classes,",
    "            Revisioning Contacts & Locations",
    "Day 30: Permissions, Creating Activities,",
    "            Assigning Activities, Job Lifecycle",
]
add_bullet_list(slide, Inches(4.8), Inches(2.85), Inches(3.4), Inches(2.2),
                pc_config, font_size=9, color=CHARCOAL_600, bullet_char="\u2022", spacing=Pt(3))

# Integration (Days 31-40)
add_card(slide, Inches(8.6), Inches(2.2), Inches(4.2), Inches(3.0), bg_color=LIGHT_GOLD)
add_text_box(slide, Inches(8.8), Inches(2.3), Inches(3.8), Inches(0.3),
             "DAYS 31-40\nIntegration", font_size=12, color=GOLD_DARK, bold=True)
add_divider(slide, Inches(8.8), Inches(2.75), Inches(3.8))
integration = [
    "Day 31: Integration Intro, Gosu for Integration",
    "Day 32: Gosu Queries, Bundles & Transactions",
    "Day 33: Templates, XML Modeler",
    "Day 34: Integration Views",
    "Day 35: RESTful & SOAP Web Services",
    "Day 36: Plugins, Messaging Architecture",
    "Day 37: Triggering & Creating Messages",
    "Day 38: Payload Transformation, Sending",
    "Day 39: Acknowledging Messages, Authentication",
    "Day 40: Batch Processes, Document Management",
]
add_bullet_list(slide, Inches(8.8), Inches(2.85), Inches(3.8), Inches(2.2),
                integration, font_size=8, color=CHARCOAL_600, bullet_char="\u2022", spacing=Pt(1))

# Summary stats
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.3),
             "PROGRAM TOTALS", font_size=12, color=BLACK, bold=True)

stats = [
    ("40", "DAYS", "Complete curriculum\ncovering everything"),
    ("121", "LESSONS", "Each with PPT, video,\ninterview questions"),
    ("119", "ASSIGNMENTS", "Hands-on labs building\nthe full project"),
    ("6", "MODULES", "Business → Dev →\nConfig → Integration"),
    ("500+", "INTERVIEW Q's", "Prepared for every\npossible question"),
    ("1", "FULL PROJECT", "Complete implementation\nready for portfolio"),
]

sw, sh = Inches(1.85), Inches(1.2)
for i, (num, label, desc) in enumerate(stats):
    left = Inches(0.5) + i * (sw + Inches(0.15))
    add_card(slide, left, Inches(5.85), sw, sh)
    add_text_box(slide, left + Inches(0.1), Inches(5.9),
                 sw - Inches(0.2), Inches(0.35), num, font_size=24, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(6.25),
                 sw - Inches(0.2), Inches(0.2), label, font_size=9, color=BLACK, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(6.45),
                 sw - Inches(0.2), Inches(0.5), desc, font_size=7, color=CHARCOAL_500,
                 alignment=PP_ALIGN.CENTER)

add_page_number(slide, 12, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 13: PROFILE & INTERVIEW PREP
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "THE FINAL PHASE", color=CHARCOAL_500)
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "Profile, Project & Interview Prep", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.7), Inches(9), Inches(0.3),
             "This is where we turn a student into a candidate. Most training programs stop at teaching. We go further.",
             font_size=14, color=CHARCOAL_400)

final_cards = [
    ("PROFILE\nCREATION", [
        "Resume reviewed and enhanced",
        "Guidewire-specific keywords added",
        "Project experience written up",
        "Skills matrix documented",
        "LinkedIn profile optimized",
        "Professional headshot & branding",
    ], "Profile is ready to submit\nto clients on Day 1 after\ngraduation.", BLUE),
    ("CUSTOM\nPROJECT", [
        "Project assigned based on resume",
        "Tailored to student's background",
        "Full implementation showcase",
        "Real business scenarios used",
        "Code review by trainer",
        "Ready for 'tell me about your project'",
    ], "When interviewer asks about\nproject, student has a REAL\nanswer with REAL details.", GREEN),
    ("INTERVIEW\nPREPARATION", [
        "500+ topic-specific questions covered",
        "Mock interviews with trainer",
        "Behavioral question practice",
        "Technical scenario walkthroughs",
        "Common gotchas reviewed",
        "Confidence building exercises",
    ], "Student has already answered\nevery possible question\nbefore the real interview.", GOLD),
    ("PLACEMENT\nREADINESS", [
        "Added to InTime talent pool",
        "Top performers: W2 bench offer",
        "Matched to active job openings",
        "Interview coaching continues",
        "Support through placement process",
        "Ongoing career guidance",
    ], "Seamless handoff to our\nAccount Manager and\nBench Sales teams.", AMBER),
]

fw, fh, ftop = Inches(2.85), Inches(4.1), Inches(2.3)
for i, (title, items, note, accent) in enumerate(final_cards):
    left = Inches(0.6) + i * (fw + Inches(0.15))
    add_card(slide, left, ftop, fw, fh, bg_color=DARK_CARD, border_color=CHARCOAL_700)
    add_text_box(slide, left + Inches(0.15), ftop + Inches(0.1),
                 fw - Inches(0.3), Inches(0.4), title, font_size=12, color=accent, bold=True)
    add_divider(slide, left + Inches(0.15), ftop + Inches(0.55), fw - Inches(0.3))
    add_bullet_list(slide, left + Inches(0.15), ftop + Inches(0.65),
                    fw - Inches(0.3), Inches(2.0),
                    items, font_size=9, color=CHARCOAL_400, bullet_char="\u2022", spacing=Pt(2))
    add_text_box(slide, left + Inches(0.15), ftop + fh - Inches(0.75),
                 fw - Inches(0.3), Inches(0.6), note, font_size=8, color=accent)

add_text_box(slide, Inches(1), Inches(6.7), Inches(11), Inches(0.4),
             "The difference between our graduates and others: They can talk about their project, answer any question, "
             "and have a profile ready to submit. That's what placement teams need.",
             font_size=12, color=CHARCOAL_500, bold=True)
add_page_number(slide, 13, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 14: WHAT OUR TEAM NEEDS TO KNOW
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "FOR OUR TEAM")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "How Each Team Interacts With Training", font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(9), Inches(0.3),
             "Every team in InTime touches the training program. Here's how you connect.",
             font_size=14, color=CHARCOAL_600)

team_roles = [
    ("Training Sales", "YOUR JOB",
     "Sell the program to leads.\nExplain both models.\nHandle enrollment.\nFollow up with prospects.",
     "KEY TALKING POINTS",
     "100+ lessons, nothing skipped\nFull project, not just theory\nSelf-paced or live group\n1-on-1 trainer time\nInterview prep included\nPlacement support after",
     GREEN),
    ("OPT Recruiter", "YOUR JOB",
     "Source OPT/CPT students.\nQualify for training.\nHandoff to Training Sales.\nFeed experienced to AM team.",
     "HOW TO QUALIFY",
     "IT background or interest?\nWilling to invest in training?\nAvailable for program duration?\nUS work authorization (OPT/CPT)?\nCareer changers welcome\nFresh grads are ideal",
     BLUE),
    ("Training Team\n(Lead + Trainer +\nCoordinator)", "YOUR JOB",
     "Deliver the program.\nTrack student progress.\nEnsure 100% completion.\nPrepare for placement.",
     "SUCCESS METRICS",
     "100% lesson completion\n100% assignment completion\nProject quality score 4+/5\nMock interview readiness\n90%+ student satisfaction\nPlacement within 30 days",
     AMBER),
    ("Account Manager\n& Bench Sales", "YOUR JOB",
     "Receive graduates.\nMatch to open jobs.\nSubmit to clients.\nSupport through interviews.",
     "WHAT YOU GET",
     "Profile-ready candidates\nFull project experience\n500+ interview Q's practiced\nGuidewire implementation skills\nResume updated & polished\nConfident, prepared candidates",
     GOLD),
]

tw, th = Inches(2.85), Inches(4.6)
for i, (team, label1, desc1, label2, desc2, accent) in enumerate(team_roles):
    left = Inches(0.6) + i * (tw + Inches(0.15))
    add_card(slide, left, Inches(2.1), tw, th)
    add_text_box(slide, left + Inches(0.15), Inches(2.15),
                 tw - Inches(0.3), Inches(0.5), team, font_size=11, color=accent, bold=True)
    add_divider(slide, left + Inches(0.15), Inches(2.65), tw - Inches(0.3))

    add_text_box(slide, left + Inches(0.15), Inches(2.7),
                 tw - Inches(0.3), Inches(0.15), label1, font_size=7, color=CHARCOAL_400, bold=True)
    add_text_box(slide, left + Inches(0.15), Inches(2.85),
                 tw - Inches(0.3), Inches(1.0), desc1, font_size=9, color=CHARCOAL_600)

    add_divider(slide, left + Inches(0.15), Inches(4.0), tw - Inches(0.3))
    add_text_box(slide, left + Inches(0.15), Inches(4.05),
                 tw - Inches(0.3), Inches(0.15), label2, font_size=7, color=CHARCOAL_400, bold=True)
    add_text_box(slide, left + Inches(0.15), Inches(4.2),
                 tw - Inches(0.3), Inches(2.0), desc2, font_size=9, color=CHARCOAL_600)

add_page_number(slide, 14, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 15: CLOSING
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gold_accent_line(slide)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
             "\"We don't just teach Guidewire.\nWe build Guidewire professionals.\n\n"
             "Every lesson. Every assignment.\nEvery interview question.\nOne full project.\"",
             font_size=26, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.4),
             "- The InTime Academy", font_size=14, color=GOLD, alignment=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    left=Inches(5.5), top=Inches(4.1), width=Inches(2), height=Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = GOLD
shape.line.fill.background()

# Key numbers
nums = [
    ("40", "Days"),
    ("121", "Lessons"),
    ("119", "Assignments"),
    ("6", "Modules"),
    ("500+", "Interview Q's"),
    ("1", "Full Project"),
]

nw = Inches(1.7)
for i, (num, label) in enumerate(nums):
    left = Inches(0.8) + i * (nw + Inches(0.15))
    add_text_box(slide, left, Inches(4.5), nw, Inches(0.5),
                 num, font_size=28, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(5.0), nw, Inches(0.3),
                 label, font_size=10, color=CHARCOAL_400, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.8), Inches(9), Inches(0.5),
             "Model A: Self-paced + 20 hrs live trainer.  Model B: 80 hrs live over 8 weeks.\n"
             "Both produce placement-ready Guidewire professionals. Dedicated coordinator throughout.",
             font_size=14, color=CHARCOAL_400, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(6.6), Inches(9), Inches(0.3),
             "I N T I M E   A C A D E M Y  |  Guidewire InsuranceSuite  |  February 2026",
             font_size=11, color=CHARCOAL_500, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 15, TOTAL_SLIDES)

# ──────────────────────────────────────────────
# SAVE
# ──────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)),
                           "InTime_Training_Program.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {TOTAL_SLIDES}")
