#!/usr/bin/env python3
"""
Export PPT slides as PNG images for the academy lesson viewer.
Each slide becomes a PNG in public/academy/guidewire/slides/{chapterSlug}/lesson-{nn}/slide-{nn}.png
"""

import os
import json
import subprocess
import sys

MATERIAL_DIR = "/Users/sumanthrajkumarnagolu/Library/CloudStorage/OneDrive-IntimeeSolutions/Intime E-Solutions/Trainings/Guidewire/Material"
CONTENT_DIR = "/Users/sumanthrajkumarnagolu/Projects/intime-v3/public/academy/guidewire/content"
SLIDES_DIR = "/Users/sumanthrajkumarnagolu/Projects/intime-v3/public/academy/guidewire/slides"

CHAPTER_MAP = {
    "Chapter 4 - PolicyCenter Introduction": {"id": 4, "slug": "ch04"},
    "Chapter 5 - Claim Center Introduction": {"id": 5, "slug": "ch05"},
    "Chapter 6 - Billing Center Introduction": {"id": 6, "slug": "ch06"},
    "Chapter 7 - InsuranceSuite Developer Fundamentals": {"id": 7, "slug": "ch07"},
    "Chapter 8 - policy center configuration": {"id": 8, "slug": "ch08"},
    "Chapter 9 - ClaimCenter Configuration": {"id": 9, "slug": "ch09"},
    "Chapter 11 - Introduction to Integration": {"id": 11, "slug": "ch11"},
    "Chapter 13 - Rating Introduction": {"id": 13, "slug": "ch13"},
    "Chapter 14 - Rating Configuration": {"id": 14, "slug": "ch14"},
}


def export_pptx_to_images(pptx_path, output_dir):
    """Export PPTX slides to PNG images using LibreOffice or python-pptx + Pillow."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("  ERROR: python-pptx or Pillow not available")
        return 0

    prs = Presentation(pptx_path)
    slide_width = prs.slide_width or Emu(9144000)  # default 10 inches
    slide_height = prs.slide_height or Emu(6858000)  # default 7.5 inches

    # Calculate image dimensions (1280px wide)
    img_width = 1280
    img_height = int(img_width * (slide_height / slide_width))

    os.makedirs(output_dir, exist_ok=True)
    count = 0

    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        img_path = os.path.join(output_dir, f"slide-{str(slide_num).zfill(2)}.png")

        if os.path.exists(img_path):
            count += 1
            continue

        # Create a simplified slide rendering
        img = Image.new('RGB', (img_width, img_height), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)

        # Try to get a font
        try:
            title_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 32)
            body_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 20)
            small_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 14)
        except:
            title_font = ImageFont.load_default()
            body_font = ImageFont.load_default()
            small_font = ImageFont.load_default()

        # Draw a header bar
        draw.rectangle([0, 0, img_width, 60], fill=(23, 23, 23))
        draw.text((20, 16), f"Slide {slide_num}", fill=(201, 169, 97), font=small_font)

        # Extract text from shapes
        y_pos = 80
        has_title = False

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # Check if it's a title-like shape
                is_title = (not has_title and
                            (shape.shape_id == slide.shapes[0].shape_id or
                             any(run.font.bold for run in para.runs if run.font.bold is not None) or
                             any(run.font.size and run.font.size >= Pt(20) for run in para.runs if run.font.size)))

                if is_title and not has_title:
                    has_title = True
                    # Draw title
                    draw.text((40, y_pos), text[:80], fill=(23, 23, 23), font=title_font)
                    y_pos += 50
                    # Draw separator line
                    draw.line([(40, y_pos), (img_width - 40, y_pos)], fill=(201, 169, 97), width=2)
                    y_pos += 20
                else:
                    # Draw body text with indent based on level
                    level = para.level if para.level else 0
                    indent = 40 + level * 30

                    if level > 0:
                        # Draw bullet
                        draw.ellipse([indent - 12, y_pos + 7, indent - 6, y_pos + 13], fill=(100, 100, 100))

                    # Wrap text
                    words = text.split()
                    line = ""
                    max_width = img_width - indent - 40
                    for word in words:
                        test_line = f"{line} {word}".strip()
                        bbox = draw.textbbox((0, 0), test_line, font=body_font)
                        if bbox[2] - bbox[0] > max_width and line:
                            draw.text((indent, y_pos), line, fill=(50, 50, 50), font=body_font)
                            y_pos += 28
                            line = word
                        else:
                            line = test_line
                    if line:
                        draw.text((indent, y_pos), line, fill=(50, 50, 50), font=body_font)
                        y_pos += 28

                    y_pos += 4

                if y_pos > img_height - 40:
                    break
            if y_pos > img_height - 40:
                break

        # If no text was found, add a placeholder
        if y_pos <= 100:
            draw.text((img_width // 2 - 100, img_height // 2 - 20),
                      f"Slide {slide_num} - Visual Content",
                      fill=(150, 150, 150), font=body_font)

        # Footer
        draw.rectangle([0, img_height - 30, img_width, img_height], fill=(245, 245, 245))
        draw.text((20, img_height - 24), f"Guidewire Academy", fill=(150, 150, 150), font=small_font)

        img.save(img_path, 'PNG', optimize=True)
        count += 1

    return count


def main():
    print("=" * 60)
    print("SLIDE IMAGE EXPORT")
    print("=" * 60)

    total_slides = 0

    for chapter_name, info in CHAPTER_MAP.items():
        slug = info["slug"]
        chapter_dir = os.path.join(MATERIAL_DIR, chapter_name)
        if not os.path.isdir(chapter_dir):
            continue

        # Read the index.json to know which lessons exist
        index_path = os.path.join(CONTENT_DIR, slug, "index.json")
        if not os.path.exists(index_path):
            continue

        with open(index_path) as f:
            index = json.load(f)

        print(f"\n=== {chapter_name} ({slug}) ===")

        for lesson in index.get("lessons", []):
            lesson_num = lesson["lessonId"].split("-l")[1]  # e.g., "01" from "ch04-l01"

            # Find the corresponding lesson JSON to get sourceFile
            lesson_json = os.path.join(CONTENT_DIR, slug, f"lesson-{lesson_num}.json")
            if not os.path.exists(lesson_json):
                continue

            with open(lesson_json) as f:
                lesson_data = json.load(f)

            source_file = lesson_data.get("sourceFile", "")
            source_folder = lesson_data.get("sourceFolder")

            # Find the PPTX file
            pptx_path = None
            if source_folder:
                pptx_path = os.path.join(chapter_dir, source_folder, source_file)
            if not pptx_path or not os.path.exists(pptx_path):
                # Try searching recursively
                for root, dirs, files in os.walk(chapter_dir):
                    if source_file in files:
                        pptx_path = os.path.join(root, source_file)
                        break

            if not pptx_path or not os.path.exists(pptx_path):
                continue

            out_dir = os.path.join(SLIDES_DIR, slug, f"lesson-{lesson_num}")
            print(f"  Lesson {lesson_num}: {source_file}")

            count = export_pptx_to_images(pptx_path, out_dir)
            total_slides += count
            print(f"    -> {count} slides exported")

    print(f"\n{'=' * 60}")
    print(f"Total: {total_slides} slide images exported")
    print(f"Output: {SLIDES_DIR}")


if __name__ == "__main__":
    main()
