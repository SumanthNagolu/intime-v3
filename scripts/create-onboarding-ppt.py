#!/usr/bin/env python3
"""
Generate InTime Functional Team Onboarding PPT
Complete operations-focused presentation with detailed role breakdowns,
training program, intern scaling model, and expectations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ──────────────────────────────────────────────
# COLORS
# ──────────────────────────────────────────────
BLACK = RGBColor(0x17, 0x17, 0x17)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xC9, 0xA9, 0x61)
GOLD_DARK = RGBColor(0xA8, 0x8B, 0x4A)
CREAM = RGBColor(0xFD, 0xFB, 0xF7)
CHARCOAL_100 = RGBColor(0xF5, 0xF5, 0xF5)
CHARCOAL_200 = RGBColor(0xE5, 0xE5, 0xE5)
CHARCOAL_400 = RGBColor(0xA3, 0xA3, 0xA3)
CHARCOAL_500 = RGBColor(0x73, 0x73, 0x73)
CHARCOAL_600 = RGBColor(0x52, 0x52, 0x52)
CHARCOAL_700 = RGBColor(0x40, 0x40, 0x40)
GREEN = RGBColor(0x0A, 0x87, 0x54)
RED = RGBColor(0xDC, 0x26, 0x26)
BLUE = RGBColor(0x03, 0x69, 0xA1)
DARK_BG = RGBColor(0x0F, 0x0F, 0x0F)
AMBER = RGBColor(0xD9, 0x77, 0x06)
DARK_CARD = RGBColor(0x1A, 0x1A, 0x1A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ──────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────

def add_gold_accent_line(slide, top=Inches(0)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left=0, top=top,
        width=prs.slide_width, height=Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=BLACK, spacing=Pt(6), font_name='Calibri', bullet_char=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet_char} {item}" if bullet_char else item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
    return txBox

def add_card(slide, left, top, width, height, bg_color=WHITE, border_color=CHARCOAL_200):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left=left, top=top, width=width, height=height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.05
    return shape

def add_divider(slide, left, top, width):
    d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left=left, top=top,
                                width=width, height=Inches(0.02))
    d.fill.solid()
    d.fill.fore_color.rgb = CHARCOAL_200
    d.line.fill.background()

def add_section_label(slide, left, top, text, color=CHARCOAL_400):
    add_text_box(slide, left, top, Inches(6), Inches(0.3),
                 text.upper(), font_size=10, color=color, bold=True)

def add_page_number(slide, num, total):
    add_text_box(slide, Inches(12.3), Inches(7.1), Inches(0.9), Inches(0.3),
                 f"{num}/{total}", font_size=9, color=CHARCOAL_400,
                 alignment=PP_ALIGN.RIGHT)

def add_role_card(slide, left, top, width, height, role_name, role_type, cost,
                  daily_activities, hiring_profile, key_metrics,
                  bg_color=WHITE, accent_color=BLACK):
    """Compact role card with all key info"""
    card = add_card(slide, left, top, width, height, bg_color=bg_color)

    # Role name
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1),
                 width - Inches(0.4), Inches(0.3),
                 role_name, font_size=13, color=accent_color, bold=True)

    # Type badge (no cost)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.35),
                 width - Inches(0.4), Inches(0.2),
                 role_type, font_size=9, color=CHARCOAL_500)

    add_divider(slide, left + Inches(0.2), top + Inches(0.58), width - Inches(0.4))

    # Daily activities
    add_text_box(slide, left + Inches(0.2), top + Inches(0.63),
                 width - Inches(0.4), Inches(0.15),
                 "DAILY", font_size=7, color=CHARCOAL_400, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.78),
                 width - Inches(0.4), Inches(1.0),
                 daily_activities, font_size=8, color=CHARCOAL_600)

    # What we look for
    add_text_box(slide, left + Inches(0.2), top + height - Inches(1.35),
                 width - Inches(0.4), Inches(0.15),
                 "WE LOOK FOR", font_size=7, color=CHARCOAL_400, bold=True)
    add_text_box(slide, left + Inches(0.2), top + height - Inches(1.18),
                 width - Inches(0.4), Inches(0.6),
                 hiring_profile, font_size=8, color=CHARCOAL_600)

    # Key metrics
    add_text_box(slide, left + Inches(0.2), top + height - Inches(0.55),
                 width - Inches(0.4), Inches(0.15),
                 "KEY METRIC", font_size=7, color=CHARCOAL_400, bold=True)
    add_text_box(slide, left + Inches(0.2), top + height - Inches(0.38),
                 width - Inches(0.4), Inches(0.3),
                 key_metrics, font_size=9, color=accent_color, bold=True)

    return card

TOTAL_SLIDES = 20

# ══════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gold_accent_line(slide)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
             "I N T I M E", font_size=16, color=GOLD, bold=True)
add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
             "Welcome to the Team", font_size=52, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(3.8), Inches(8), Inches(0.8),
             "Our Business, Every Role, and What We Expect From You",
             font_size=22, color=CHARCOAL_400)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    left=Inches(1), top=Inches(5), width=Inches(2), height=Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = GOLD
shape.line.fill.background()

add_text_box(slide, Inches(1), Inches(5.3), Inches(7), Inches(0.4),
             "February 2026  |  US IT Staffing  |  Hyderabad Operations  |  Phase 1 Operations",
             font_size=12, color=CHARCOAL_500)
add_page_number(slide, 1, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 2: WHAT IS INTIME
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "WHO WE ARE")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
             "InTime is a US IT Staffing Company.", font_size=36, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.8), Inches(9), Inches(0.5),
             "We connect talented IT professionals with companies that need them in the United States.",
             font_size=16, color=CHARCOAL_600)

what_cards = [
    ("TRAIN", "Academy Program",
     "We train people in IT technologies\n(Guidewire, Java, etc.) from scratch.\n"
     "2-month intensive programs.\nStudents pay for training."),
    ("RECRUIT", "Find IT Talent",
     "We source experienced IT professionals\nfrom India, OPT students from US/Canada.\n"
     "Build a qualified talent pool ready\nto deploy on client projects."),
    ("PLACE", "Deploy on Projects",
     "We place consultants at US companies\non contract (hourly bill rate).\n"
     "We also do direct hire placements\n(percentage of first-year salary)."),
    ("DELIVER", "Ensure Success",
     "We manage ongoing client relationships,\nconsultant performance, and renewals.\n"
     "Happy clients = repeat business.\nThis is where revenue compounds."),
]

cw, ch, ctop, cstart = Inches(2.7), Inches(2.8), Inches(2.7), Inches(0.7)
for i, (tag, title, desc) in enumerate(what_cards):
    left = cstart + i * (cw + Inches(0.25))
    add_card(slide, left, ctop, cw, ch)
    add_text_box(slide, left + Inches(0.2), ctop + Inches(0.2),
                 Inches(1), Inches(0.25), tag, font_size=10, color=GOLD, bold=True)
    add_text_box(slide, left + Inches(0.2), ctop + Inches(0.5),
                 cw - Inches(0.4), Inches(0.35), title, font_size=15, color=BLACK, bold=True)
    add_text_box(slide, left + Inches(0.2), ctop + Inches(1.0),
                 cw - Inches(0.4), Inches(1.7), desc, font_size=11, color=CHARCOAL_600)

add_text_box(slide, Inches(1), Inches(5.9), Inches(11), Inches(0.8),
             "Hyderabad office runs all operations. US clients pay us. "
             "We earn the spread between what clients pay and what consultants cost.\n"
             "Training = immediate revenue. Placements = recurring revenue.",
             font_size=13, color=CHARCOAL_500)
add_page_number(slide, 2, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 3: THE PROBLEM
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "THE PROBLEM WE SOLVE", color=CHARCOAL_500)
add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
             "Staffing agencies run on chaos.\nWe eliminate that chaos.",
             font_size=40, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(2.8), Inches(8), Inches(0.5),
             "Every staffing agency has the same problems:", font_size=16, color=CHARCOAL_400)

problems = [
    ("12+ Tools", "ATS, CRM, spreadsheets, email,\ntimesheets, payroll - all disconnected.\nNothing talks to each other."),
    ("No Visibility", "Managers can't see what's\nhappening in real time.\nReactive firefighting daily."),
    ("40% Time Wasted", "Recruiters spend 40% of their\ntime on admin work instead\nof actually recruiting."),
    ("Data Silos", "Critical information is stuck\nin different tools and people's\nheads. No single source of truth."),
]

pw, ph, ptop, pstart = Inches(2.6), Inches(2.0), Inches(3.6), Inches(0.7)
for i, (title, desc) in enumerate(problems):
    left = pstart + i * (pw + Inches(0.25))
    add_card(slide, left, ptop, pw, ph, bg_color=DARK_CARD, border_color=CHARCOAL_700)
    add_text_box(slide, left + Inches(0.2), ptop + Inches(0.2),
                 pw - Inches(0.4), Inches(0.35), title, font_size=14, color=GOLD, bold=True)
    add_text_box(slide, left + Inches(0.2), ptop + Inches(0.7),
                 pw - Inches(0.4), Inches(1.2), desc, font_size=11, color=CHARCOAL_400)

add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
             "InTime Solution: One unified platform + one well-structured team = Zero chaos.",
             font_size=16, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 3, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 4: SEVEN FUNNELS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "OUR BUSINESS MODEL")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "Seven Interconnected Funnels", font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(9), Inches(0.3),
             "Three supply sources feed one talent pool, which deploys through four demand channels.",
             font_size=14, color=CHARCOAL_600)

# Supply
add_text_box(slide, Inches(1), Inches(2.2), Inches(5), Inches(0.3),
             "SUPPLY SIDE (Where Talent Comes From)", font_size=11, color=GREEN, bold=True)

supply = [
    ("1", "Academy", "Train new IT talent from scratch.\n2-month certification programs.\nStudents pay for training."),
    ("2", "Recruiting", "Source experienced professionals.\nBuild qualified talent pool.\nOPT, LinkedIn, referrals."),
    ("3", "Campaigns", "Attract inbound candidates\nthrough marketing & outreach.\nJob ads, social, events."),
]
fw, fh, ftop = Inches(2.4), Inches(1.45), Inches(2.6)
for i, (num, title, desc) in enumerate(supply):
    left = Inches(1) + i * (fw + Inches(0.2))
    add_card(slide, left, ftop, fw, fh, bg_color=RGBColor(0xF0, 0xF9, 0xF4))
    add_text_box(slide, left + Inches(0.15), ftop + Inches(0.1),
                 Inches(0.3), Inches(0.3), num, font_size=18, color=GREEN, bold=True)
    add_text_box(slide, left + Inches(0.5), ftop + Inches(0.1),
                 fw - Inches(0.6), Inches(0.3), title, font_size=14, color=BLACK, bold=True)
    add_text_box(slide, left + Inches(0.15), ftop + Inches(0.5),
                 fw - Inches(0.3), Inches(0.9), desc, font_size=10, color=CHARCOAL_600)

# Talent Pool
add_text_box(slide, Inches(3.5), Inches(4.2), Inches(6), Inches(0.4),
             "TALENT POOL  (Qualified, profiled, ready to deploy)",
             font_size=14, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)

# Demand
add_text_box(slide, Inches(1), Inches(4.7), Inches(5), Inches(0.3),
             "DEMAND SIDE (Where Revenue Comes From)", font_size=11, color=BLUE, bold=True)

demand = [
    ("4", "Account\nAcquisition", "Win new clients.\nVendor partnerships.\nMSA negotiations."),
    ("5", "Job\nFulfillment", "Fill client positions.\nSubmit candidates.\nClose placements."),
    ("6", "Bench\nSales", "Sell available W2\nconsultants to\nclient projects."),
    ("7", "Delivery\nMgmt", "Ensure success.\nClient satisfaction.\nRenewals & growth."),
]
fw2, fh2, ftop2 = Inches(2.4), Inches(1.4), Inches(5.1)
for i, (num, title, desc) in enumerate(demand):
    left = Inches(0.7) + i * (fw2 + Inches(0.2))
    add_card(slide, left, ftop2, fw2, fh2, bg_color=RGBColor(0xEF, 0xF6, 0xFF))
    add_text_box(slide, left + Inches(0.15), ftop2 + Inches(0.08),
                 Inches(0.3), Inches(0.3), num, font_size=18, color=BLUE, bold=True)
    add_text_box(slide, left + Inches(0.5), ftop2 + Inches(0.05),
                 fw2 - Inches(0.6), Inches(0.5), title, font_size=12, color=BLACK, bold=True)
    add_text_box(slide, left + Inches(0.15), ftop2 + Inches(0.6),
                 fw2 - Inches(0.3), Inches(0.7), desc, font_size=9, color=CHARCOAL_600)

# Flywheel
add_text_box(slide, Inches(10.5), Inches(2.2), Inches(2.5), Inches(4.3),
             "THE FLYWHEEL\n\n"
             "Graduates\n     |\nTalent Pool\n     |\nPlacements\n     |\n"
             "Client Success\n     |\nRenewals\n     |\nReputation\n     |\nMore Growth",
             font_size=10, color=CHARCOAL_500, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 4, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 5: TRAINING PROGRAM DEEP DIVE
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "TRAINING PROGRAM - OUR REVENUE ENGINE")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "How the Academy Works", font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(10), Inches(0.4),
             "40 days. 121 lessons. 119 assignments. Every lesson has explanation PPTs with embedded demo videos, "
             "interview questions, and hands-on assignments. Labs build a FULL Guidewire project end-to-end.",
             font_size=13, color=CHARCOAL_600)

# Curriculum overview card
add_card(slide, Inches(0.7), Inches(2.2), Inches(7.5), Inches(2.6))
add_text_box(slide, Inches(0.9), Inches(2.3), Inches(5), Inches(0.3),
             "40-DAY CURRICULUM (6 MODULES)", font_size=12, color=BLACK, bold=True)
add_divider(slide, Inches(0.9), Inches(2.6), Inches(7.1))

modules = [
    ("Days 1-9", "PolicyCenter Business", "Insurance intro, accounts, policies, transactions,\nproduct model, underwriting, rating basics"),
    ("Days 10-13", "ClaimCenter Business", "Claims process, intake, setup, adjudication,\nfinancials, payments, contacts, vendors"),
    ("Days 14-20", "Developer Fundamentals", "Configuration, data model, UI architecture,\nGosu, widgets, views, entities, typelists"),
    ("Days 21-26", "ClaimCenter Config", "UI config, rules, financials config,\nsearch, permissions, activities, archiving"),
    ("Days 27-30", "PolicyCenter Config", "Data model, wizards, revisioning,\nunderwriting issues, permissions, activities"),
    ("Days 31-40", "Integration", "REST/SOAP APIs, plugins, messaging,\nbundles, queries, batch processes, auth"),
]

for j, (days, module, desc) in enumerate(modules):
    y = Inches(2.7) + j * Inches(0.35)
    add_text_box(slide, Inches(0.9), y, Inches(1.2), Inches(0.3),
                 days, font_size=9, color=GOLD, bold=True)
    add_text_box(slide, Inches(2.1), y, Inches(2.0), Inches(0.3),
                 module, font_size=9, color=BLACK, bold=True)
    add_text_box(slide, Inches(4.2), y, Inches(3.8), Inches(0.3),
                 desc.split('\n')[0], font_size=8, color=CHARCOAL_600)

# How it's structured
add_card(slide, Inches(8.5), Inches(2.2), Inches(4.2), Inches(2.6))
add_text_box(slide, Inches(8.7), Inches(2.3), Inches(3.8), Inches(0.3),
             "EACH LESSON INCLUDES", font_size=12, color=BLACK, bold=True)
add_divider(slide, Inches(8.7), Inches(2.6), Inches(3.8))
lesson_items = [
    "Detailed explanation PPT slides",
    "Embedded demo videos",
    "Interview prep questions",
    "Hands-on assignment (lab)",
    "Assignments = full project build",
    "Each module = epic, each lesson = user story",
    "All labs connect into one complete project",
]
add_bullet_list(slide, Inches(8.7), Inches(2.7), Inches(3.8), Inches(1.8),
                lesson_items, font_size=9, color=CHARCOAL_600, bullet_char="\u2022", spacing=Pt(2))

# Two delivery models
add_text_box(slide, Inches(1), Inches(5.0), Inches(5), Inches(0.3),
             "TWO DELIVERY MODELS", font_size=12, color=BLACK, bold=True)

# Model A - Self-paced
add_card(slide, Inches(0.7), Inches(5.35), Inches(5.8), Inches(1.6))
add_text_box(slide, Inches(0.9), Inches(5.4), Inches(3), Inches(0.25),
             "MODEL A: SELF-PACED", font_size=11, color=BLUE, bold=True)
add_text_box(slide, Inches(0.9), Inches(5.65), Inches(5.4), Inches(1.2),
             "Designed for busy US/Canada resources.\n"
             "Students complete 121 lessons + assignments at own pace.\n"
             "Profile created, custom project assigned based on resume.\n"
             "1-on-1 trainer sessions at scheduled times for project + interview prep.\n"
             "Ensures EVERYTHING is covered (unlike group programs that skip topics).",
             font_size=9, color=CHARCOAL_600)

# Model B - Live group
add_card(slide, Inches(6.8), Inches(5.35), Inches(5.8), Inches(1.6))
add_text_box(slide, Inches(7.0), Inches(5.4), Inches(3), Inches(0.25),
             "MODEL B: LIVE GROUP (2 WEEKS)", font_size=11, color=GREEN, bold=True)
add_text_box(slide, Inches(7.0), Inches(5.65), Inches(5.4), Inches(1.2),
             "Small group with live trainer for 2 weeks intensive.\n"
             "Lessons assigned as homework - students study & do assignments.\n"
             "Come to sessions with questions/blockers for discussion.\n"
             "Profile created Day 1, project assigned Day 1.\n"
             "Trainer guides students through project at appropriate times.",
             font_size=9, color=CHARCOAL_600)

add_page_number(slide, 5, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 6: SUPPLY SIDE - TRAINING POD ROLES
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "TEAM STRUCTURE")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "Training Team: 4 People", font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(9), Inches(0.3),
             "Mission: Enroll students, deliver 40-day Guidewire program, and produce placement-ready professionals.",
             font_size=14, color=CHARCOAL_600)

rw, rh = Inches(2.8), Inches(4.1)
rtop = Inches(2.1)

# Training Lead
add_role_card(slide, Inches(0.6), rtop, rw, rh,
    role_name="Training Lead",
    role_type="Senior Hire",
    cost="",
    daily_activities="Morning standup with trainer (15 min)\nReview student progress dashboards\nWeekly 1:1 with each student\nCoordinate with Screening on pipeline\nTraining quality audit (sit in sessions)",
    hiring_profile="3-5 yrs training ops or tech recruiting\nManaged 5+ direct reports\nUnderstands IT staffing industry\nStrong on accountability & tracking",
    key_metrics="Batch fill: 100% | Retention: 90%+",
    accent_color=GREEN)

# Trainer
add_role_card(slide, Inches(3.55), rtop, rw, rh,
    role_name="Trainer (Guidewire)",
    role_type="Senior Hire - CRITICAL",
    cost="",
    daily_activities="Deliver 40-day curriculum (121 lessons)\nLive sessions or self-paced guidance\nReview assignments & project work\nDocument student progress daily\n1-on-1 sessions for interview prep",
    hiring_profile="5+ yrs Guidewire implementation\n3+ full lifecycle projects\nTeaching/mentoring experience\nPatient, structured communicator\nTHIS HIRE MAKES OR BREAKS US",
    key_metrics="100% sessions delivered | 4.5+/5 rating",
    accent_color=GREEN)

# Training Sales
add_role_card(slide, Inches(6.5), rtop, rw, rh,
    role_name="Training Sales",
    role_type="Intern",
    cost="",
    daily_activities="Respond to inquiries within 2 hours\n20+ outbound calls/messages daily\nConduct discovery calls with leads\nUpdate CRM with all interactions\n3-5 enrollment closes per week",
    hiring_profile="2-3 yrs education or tech sales\nConsultative selling experience\nComfortable with high-value conversations\nPersistent but not pushy",
    key_metrics="Response < 2hrs | 4-6 enrollments/month",
    accent_color=GREEN)

# Coordinator
add_role_card(slide, Inches(9.45), rtop, rw, rh,
    role_name="Training Coordinator",
    role_type="Intern",
    cost="",
    daily_activities="Confirm session attendance (AM check)\nPrepare session materials/links\nTrack attendance, flag no-shows\nCollect daily feedback from students\nOnboard new students (1-2/week)",
    hiring_profile="1-2 yrs operations/coordination\nExtremely organized, detail-oriented\nGood with students/customer-facing\nComfortable with CRM & spreadsheets",
    key_metrics="100% ready 30 min before | 0 gaps",
    accent_color=GREEN)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.6),
             "The Trainer is our most critical hire. Without a great trainer, the entire model fails.\n"
             "Training Lead manages the team and holds everyone accountable for student outcomes.",
             font_size=12, color=CHARCOAL_500, bold=True)
add_page_number(slide, 6, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 7: SUPPLY SIDE - SCREENING POD ROLES
# ══════════════════════════════════════════════
# SLIDE 7: ACCOUNT MANAGEMENT & RECRUITING TEAM
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "TEAM STRUCTURE")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "Account Management & Recruiting Team", font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(10), Inches(0.4),
             "Full-cycle team. Same team sources candidates (backend) AND places them on client jobs (frontend). No handoffs.",
             font_size=14, color=CHARCOAL_600)

# How this team works
add_card(slide, Inches(0.6), Inches(2.2), Inches(11.8), Inches(1.0), bg_color=RGBColor(0xEF, 0xF6, 0xFF))
add_text_box(slide, Inches(0.8), Inches(2.25), Inches(5.5), Inches(0.25),
             "BACKEND (Sourcing & Screening)", font_size=11, color=GREEN, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(0.6),
             "Recruiters source candidates from job boards, LinkedIn,\nreferrals. Screen, qualify, and build the talent pool.",
             font_size=10, color=CHARCOAL_600)
add_text_box(slide, Inches(6.6), Inches(2.25), Inches(5.5), Inches(0.25),
             "FRONTEND (Placing & Delivery)", font_size=11, color=BLUE, bold=True)
add_text_box(slide, Inches(6.6), Inches(2.5), Inches(5.5), Inches(0.6),
             "Same team submits candidates to client jobs, manages\nclient relationships, and ensures placement success.",
             font_size=10, color=CHARCOAL_600)

rtop_am = Inches(3.45)
rh_am = Inches(3.4)

# Account Manager
add_role_card(slide, Inches(0.6), rtop_am, Inches(3.8), rh_am,
    role_name="Account Manager (Lead)",
    role_type="Senior Hire",
    cost="",
    daily_activities="Pipeline review with recruiters\nClient relationship management\nJob intake calls with clients\nSubmission quality review\nNew client outreach & acquisition\nTeam performance reviews weekly",
    hiring_profile="4+ yrs IT recruiting/account mgmt\nDirect client relationship experience\nBusiness development skills\nStrategic thinker, not just executor",
    key_metrics="50%+ fill rate | < 3 weeks to fill | 90%+ client retention",
    accent_color=BLUE)

# AM focus areas
add_card(slide, Inches(0.6), Inches(6.15), Inches(3.8), Inches(0.85), bg_color=RGBColor(0xEF, 0xF6, 0xFF))
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(3.4), Inches(0.2),
             "ACCOUNT MANAGER OWNS:", font_size=9, color=BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.4), Inches(3.4), Inches(0.5),
             "1. Client relationships & job requirements\n"
             "2. Submission quality & candidate-job fit\n"
             "3. New client acquisition & growth",
             font_size=9, color=CHARCOAL_600)

# Recruiter 1
add_role_card(slide, Inches(4.6), rtop_am, Inches(3.8), rh_am,
    role_name="Recruiter #1",
    role_type="Intern",
    cost="",
    daily_activities="Source candidates for active client jobs\n40+ outreach messages daily\n8+ screening calls per day\nSubmit qualified candidates to AM\nFollow up on pending submissions\nBuild candidate relationships",
    hiring_profile="1-2 yrs IT recruiting\nComfortable with high volume\nGood phone & assessment skills\nOrganized with follow-ups",
    key_metrics="40+ outreach/day | 8+ screens/day | 1+ placement/month",
    accent_color=BLUE)

# Recruiter 2
add_role_card(slide, Inches(8.6), rtop_am, Inches(3.8), rh_am,
    role_name="Recruiter #2",
    role_type="Intern",
    cost="",
    daily_activities="Source candidates for active client jobs\n40+ outreach messages daily\n8+ screening calls per day\nSubmit qualified candidates to AM\nFollow up on pending submissions\nBuild candidate relationships",
    hiring_profile="1-2 yrs IT recruiting\nComfortable with high volume\nGood phone & assessment skills\nOrganized with follow-ups",
    key_metrics="40+ outreach/day | 8+ screens/day | 1+ placement/month",
    accent_color=BLUE)

add_page_number(slide, 7, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 8: BENCH SALES TEAM
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "TEAM STRUCTURE")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "Bench Sales Team", font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(10), Inches(0.4),
             "Full-cycle team. One person pools all jobs for a technology (backend), bench sales submit consultants (frontend).",
             font_size=14, color=CHARCOAL_600)

# How bench works
add_card(slide, Inches(0.6), Inches(2.2), Inches(11.8), Inches(1.0), bg_color=RGBColor(0xEF, 0xF6, 0xFF))
add_text_box(slide, Inches(0.8), Inches(2.25), Inches(5.5), Inches(0.25),
             "BACKEND (Job Pooling)", font_size=11, color=GREEN, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(0.6),
             "E.g., one person collects ALL Guidewire jobs from every\nvendor, portal, and email into one place. Single source of truth.",
             font_size=10, color=CHARCOAL_600)
add_text_box(slide, Inches(6.6), Inches(2.25), Inches(5.5), Inches(0.25),
             "FRONTEND (Submitting & Closing)", font_size=11, color=BLUE, bold=True)
add_text_box(slide, Inches(6.6), Inches(2.5), Inches(5.5), Inches(0.6),
             "Bench sales reps match bench consultants to pooled jobs,\ncustomize resumes, submit, negotiate rates, close placements.",
             font_size=10, color=CHARCOAL_600)

rtop_bs = Inches(3.45)
rh_bs = Inches(3.4)

# Bench Sales Manager
add_role_card(slide, Inches(0.6), rtop_bs, Inches(3.8), rh_bs,
    role_name="Bench Sales Manager (Lead)",
    role_type="Senior Hire",
    cost="",
    daily_activities="Morning huddle with team (15 min)\nReview hot submissions & pipeline\nVendor relationship development\nBench review (who's available, how long)\nRate negotiation & margin management\nCoordinate with Account Manager on shared opps",
    hiring_profile="4+ yrs bench sales or account mgmt\nExisting vendor relationships preferred\nUnderstands MSAs, SOWs, rate cards\nCalm under pressure\nGood at objection handling",
    key_metrics="80%+ bench on billing within 30 days | < 3 week velocity",
    accent_color=BLUE)

# Bench example
add_card(slide, Inches(0.6), Inches(6.15), Inches(3.8), Inches(0.85), bg_color=RGBColor(0xEF, 0xF6, 0xFF))
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(3.4), Inches(0.2),
             "EXAMPLE - GUIDEWIRE TECH:", font_size=9, color=BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.4), Inches(3.4), Inches(0.5),
             "1 person pools ALL Guidewire jobs daily\n"
             "Bench reps submit our Guidewire consultants\n"
             "Same model for Java, Salesforce, etc.",
             font_size=9, color=CHARCOAL_600)

# Bench Sales Rep 1
add_role_card(slide, Inches(4.6), rtop_bs, Inches(3.8), rh_bs,
    role_name="Bench Sales Rep #1",
    role_type="Intern",
    cost="",
    daily_activities="Morning: Requirement hunting (job boards,\n  vendor emails, portal checks)\nMid-day: Match consultants to requirements,\n  customize resume, submit with cover note\nAfternoon: Chase pending, schedule interviews,\n  prep consultants\nEvening: Log all activity, flag blockers",
    hiring_profile="2-3 yrs in bench sales\nHigh energy, persistent\nGood at rate negotiation\nUnderstands C2C, W2, 1099",
    key_metrics="15+ submissions/day | 1+ placement/month",
    accent_color=BLUE)

# Bench Sales Rep 2
add_role_card(slide, Inches(8.6), rtop_bs, Inches(3.8), rh_bs,
    role_name="Bench Sales Rep #2",
    role_type="Intern",
    cost="",
    daily_activities="Morning: Requirement hunting (job boards,\n  vendor emails, portal checks)\nMid-day: Match consultants to requirements,\n  customize resume, submit with cover note\nAfternoon: Chase pending, schedule interviews,\n  prep consultants\nEvening: Log all activity, flag blockers",
    hiring_profile="2-3 yrs in bench sales\nHigh energy, persistent\nGood at rate negotiation\nUnderstands C2C, W2, 1099",
    key_metrics="15+ submissions/day | 1+ placement/month",
    accent_color=BLUE)

add_page_number(slide, 8, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 9: OPT RECRUITING, BDM & DELIVERY
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "TEAM STRUCTURE")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "OPT Recruiting, BDM & Delivery", font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(10), Inches(0.4),
             "These roles feed the other teams. OPT builds the talent pipeline, BDM wins clients, Delivery ensures success.",
             font_size=14, color=CHARCOAL_600)

rtop_s = Inches(2.2)
rh_s = Inches(4.5)

# OPT Recruiter
add_role_card(slide, Inches(0.6), rtop_s, Inches(3.8), rh_s,
    role_name="OPT Recruiter",
    role_type="Intern / Junior",
    cost="",
    daily_activities="Source from OPT boards, university groups\n30+ outreach messages daily\n5+ screening calls per day\nCommunity engagement (Discord, Slack)\nFeed qualified candidates to AM team\nHandoff training leads to Training Sales",
    hiring_profile="2+ yrs recruiting in US staffing\nUnderstands OPT/CPT/H1B basics\nExperience with intl student community\nActive on relevant platforms",
    key_metrics="30+ outreach/day | 5+ qualified leads/week",
    accent_color=GREEN)

# OPT focus
add_card(slide, Inches(0.6), Inches(6.0), Inches(3.8), Inches(1.0), bg_color=RGBColor(0xF0, 0xF9, 0xF4))
add_text_box(slide, Inches(0.8), Inches(6.05), Inches(3.4), Inches(0.2),
             "OPT RECRUITER FEEDS:", font_size=9, color=GREEN, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.25), Inches(3.4), Inches(0.7),
             "1. Training Team - candidates who need upskilling\n"
             "2. AM Team - experienced candidates ready to place\n"
             "3. Bench - graduates who join our W2 bench",
             font_size=9, color=CHARCOAL_600)

# BDM
add_role_card(slide, Inches(4.6), rtop_s, Inches(3.8), rh_s,
    role_name="Business Development Manager",
    role_type="Senior Hire",
    cost="",
    daily_activities="10+ outbound prospecting calls/emails\nFollow up on warm leads daily\nLinkedIn networking & engagement\nVendor portal registrations\n3-5 discovery calls with prospects/week\nMSA/contract negotiations",
    hiring_profile="4+ yrs IT staffing business development\nExisting vendor relationships preferred\nTrack record of closing MSAs\nUnderstands staffing economics\nHunter mentality, not farmer",
    key_metrics="2-3 vendors onboarded/month | 1-2 direct clients/quarter",
    accent_color=GOLD)

# BDM focus
add_card(slide, Inches(4.6), Inches(6.0), Inches(3.8), Inches(1.0), bg_color=RGBColor(0xFD, 0xF8, 0xED))
add_text_box(slide, Inches(4.8), Inches(6.05), Inches(3.4), Inches(0.2),
             "BDM FEEDS:", font_size=9, color=GOLD_DARK, bold=True)
add_text_box(slide, Inches(4.8), Inches(6.25), Inches(3.4), Inches(0.7),
             "1. AM Team - new client accounts & job orders\n"
             "2. Bench Team - vendor partnerships for submissions\n"
             "3. Training - strategic partnerships & referrals",
             font_size=9, color=CHARCOAL_600)

# Delivery Manager
add_role_card(slide, Inches(8.6), rtop_s, Inches(3.8), rh_s,
    role_name="Delivery Manager",
    role_type="Senior Hire",
    cost="",
    daily_activities="Check-in with active consultants (rotating)\nMonitor consultant performance signals\nAddress client/consultant issues immediately\nClient satisfaction calls weekly\nConsultant 1:1s (esp first 30 days)\nRenewal pipeline management",
    hiring_profile="3+ yrs IT staffing delivery/acct mgmt\nStrong relationship management\nProblem-solving mindset\nManaged consultants on client sites\nCalm under pressure, empathetic but firm",
    key_metrics="95%+ retention | 80%+ renewal | < 48hr issue resolution",
    accent_color=BLUE)

# Delivery focus
add_card(slide, Inches(8.6), Inches(6.0), Inches(3.8), Inches(1.0), bg_color=RGBColor(0xEF, 0xF6, 0xFF))
add_text_box(slide, Inches(8.8), Inches(6.05), Inches(3.4), Inches(0.2),
             "DELIVERY ENSURES:", font_size=9, color=BLUE, bold=True)
add_text_box(slide, Inches(8.8), Inches(6.25), Inches(3.4), Inches(0.7),
             "1. Consultant success (first 90 days critical)\n"
             "2. Client satisfaction (proactive check-ins)\n"
             "3. Renewals & extensions (track end dates)",
             font_size=9, color=CHARCOAL_600)

add_page_number(slide, 9, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 10: THE INTERN MODEL - HOW WE SCALE
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "THE INTERN MODEL", color=CHARCOAL_500)
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "How We Scale Teams Using Interns", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.7), Inches(9), Inches(0.4),
             "Each team has ONE senior lead who trains interns. Interns become your recruiters, bench sales reps, and coordinators.",
             font_size=14, color=CHARCOAL_400)

# Why this works
add_card(slide, Inches(0.6), Inches(2.3), Inches(5.8), Inches(2.2), bg_color=DARK_CARD, border_color=CHARCOAL_700)
add_text_box(slide, Inches(0.8), Inches(2.4), Inches(3), Inches(0.3),
             "WHY THIS MODEL WORKS", font_size=12, color=GOLD, bold=True)
why_items = [
    "Lower cost = we can hire more people, faster",
    "Build loyalty - you trained them, they stay",
    "Senior leads are accountable for pod outcomes",
    "After 3 months, interns are fully productive",
    "Interns who perform well get promoted to leads",
]
add_bullet_list(slide, Inches(0.8), Inches(2.75), Inches(5.4), Inches(1.6),
                why_items, font_size=11, color=CHARCOAL_400, bullet_char="\u2022", spacing=Pt(3))

# Intern-to-productive timeline
add_card(slide, Inches(6.8), Inches(2.3), Inches(5.8), Inches(2.2), bg_color=DARK_CARD, border_color=CHARCOAL_700)
add_text_box(slide, Inches(7.0), Inches(2.4), Inches(3), Inches(0.3),
             "INTERN PRODUCTIVITY TIMELINE", font_size=12, color=GOLD, bold=True)

timeline = [
    ("Month 1", "10%", "Learning basics, shadowing senior lead"),
    ("Month 2", "40%", "Guided work, some independent tasks"),
    ("Month 3", "70%", "Mostly independent, senior reviews"),
    ("Month 4+", "100%", "Fully productive, evaluate & promote"),
]

for j, (month, pct, desc) in enumerate(timeline):
    y = Inches(2.8) + j * Inches(0.38)
    add_text_box(slide, Inches(7.0), y, Inches(1), Inches(0.3),
                 month, font_size=10, color=WHITE, bold=True)
    add_text_box(slide, Inches(8.2), y, Inches(0.6), Inches(0.3),
                 pct, font_size=10, color=GOLD, bold=True)
    add_text_box(slide, Inches(8.9), y, Inches(3.5), Inches(0.3),
                 desc, font_size=10, color=CHARCOAL_400)

# Intern growth path (no salaries)
add_card(slide, Inches(0.6), Inches(4.7), Inches(5.8), Inches(2.0), bg_color=DARK_CARD, border_color=CHARCOAL_700)
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(4), Inches(0.3),
             "INTERN GROWTH PATH", font_size=12, color=GOLD, bold=True)

growth_stages = [
    ("Month 1", "LEARNING", "Shadow senior, learn basics, 10% productive"),
    ("Month 2-3", "GUIDED", "Supervised work, building speed, 40-70%"),
    ("Month 4", "INDEPENDENT", "Full workload, meeting targets, 100%"),
    ("Month 6+", "PROMOTION", "Top performers promoted to team leads"),
]

for j, (stage, label, desc) in enumerate(growth_stages):
    y = Inches(5.15) + j * Inches(0.35)
    add_text_box(slide, Inches(0.8), y, Inches(1.2), Inches(0.3),
                 stage, font_size=10, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.1), y, Inches(1.3), Inches(0.3),
                 label, font_size=10, color=GOLD, bold=True)
    add_text_box(slide, Inches(3.5), y, Inches(2.8), Inches(0.3),
                 desc, font_size=10, color=CHARCOAL_400)

# Hiring plan
add_card(slide, Inches(6.8), Inches(4.7), Inches(5.8), Inches(2.0), bg_color=DARK_CARD, border_color=CHARCOAL_700)
add_text_box(slide, Inches(7.0), Inches(4.8), Inches(4), Inches(0.3),
             "HIRING SEQUENCE", font_size=12, color=GOLD, bold=True)

hire_phases = [
    ("Week 1-2", "Trainer, Training Lead, Bench Sales Mgr, BDM (4 seniors)"),
    ("Week 3-4", "Account Manager + 2 Bench Sales Interns + OPT Recruiter"),
    ("Month 2", "Delivery Mgr + 2 Recruiters + Training Sales + Coordinator"),
    ("Month 3", "Full team operational, interns at 70% productivity"),
    ("Month 4+", "Evaluate: promote top interns, scale teams as needed"),
]

for j, (when, what) in enumerate(hire_phases):
    y = Inches(5.15) + j * Inches(0.32)
    add_text_box(slide, Inches(7.0), y, Inches(1.3), Inches(0.25),
                 when, font_size=9, color=WHITE, bold=True)
    add_text_box(slide, Inches(8.4), y, Inches(4), Inches(0.25),
                 what, font_size=9, color=CHARCOAL_400)

add_text_box(slide, Inches(1), Inches(6.95), Inches(11), Inches(0.3),
             "What we look for in interns: Recent grads, hungry to learn, basic communication skills, coachable, persistent.",
             font_size=11, color=CHARCOAL_500)
add_page_number(slide, 10, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 11: PLAN TO HIRE IN SEQUENCE
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "HIRING PLAN")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "Plan to Hire in Sequence", font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(9), Inches(0.3),
             "We hire seniors first. They set up systems and processes. Then we bring in interns who learn under them.",
             font_size=14, color=CHARCOAL_600)

# Phase 1: Week 1-2
add_card(slide, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.0))
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(1.5), Inches(0.3),
             "WEEK 1-2", font_size=14, color=GOLD, bold=True)
add_text_box(slide, Inches(2.3), Inches(2.3), Inches(3), Inches(0.3),
             "Foundation Hires (4 Seniors)", font_size=14, color=BLACK, bold=True)
add_divider(slide, Inches(0.8), Inches(2.6), Inches(5.4))

w1_hires = [
    ("Trainer (Guidewire)", "MOST CRITICAL. Without a great trainer, the model fails."),
    ("Training Lead", "Sets up training ops, enrollment pipeline, student tracking."),
    ("Bench Sales Manager", "Starts building vendor relationships, registers on portals."),
    ("BDM", "Starts client acquisition from day 1. Discovery calls, outreach."),
]
for j, (role, why) in enumerate(w1_hires):
    y = Inches(2.7) + j * Inches(0.35)
    add_text_box(slide, Inches(0.8), y, Inches(2), Inches(0.3),
                 role, font_size=10, color=BLACK, bold=True)
    add_text_box(slide, Inches(2.9), y, Inches(3.3), Inches(0.3),
                 why, font_size=9, color=CHARCOAL_600)

add_text_box(slide, Inches(0.8), Inches(4.0), Inches(5.4), Inches(0.15),
             "WHY THIS ORDER: These 4 generate revenue. Everything else supports them.",
             font_size=9, color=GREEN, bold=True)

# Phase 2: Week 3-4
add_card(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.0))
add_text_box(slide, Inches(7.0), Inches(2.3), Inches(1.5), Inches(0.3),
             "WEEK 3-4", font_size=14, color=GOLD, bold=True)
add_text_box(slide, Inches(8.5), Inches(2.3), Inches(3.5), Inches(0.3),
             "Team Leads + First Interns (4)", font_size=14, color=BLACK, bold=True)
add_divider(slide, Inches(7.0), Inches(2.6), Inches(5.4))

w3_hires = [
    ("Account Manager", "Leads the recruiting team. Fills client jobs BDM wins."),
    ("Bench Sales x2 (Interns)", "First interns. Learn under Bench Sales Manager."),
    ("OPT Recruiter", "Starts building OPT candidate pipeline for all teams."),
]
for j, (role, why) in enumerate(w3_hires):
    y = Inches(2.7) + j * Inches(0.4)
    add_text_box(slide, Inches(7.0), y, Inches(2.2), Inches(0.3),
                 role, font_size=10, color=BLACK, bold=True)
    add_text_box(slide, Inches(9.3), y, Inches(3.1), Inches(0.3),
                 why, font_size=9, color=CHARCOAL_600)

add_text_box(slide, Inches(7.0), Inches(4.0), Inches(5.4), Inches(0.15),
             "WHY THIS ORDER: Seniors have systems ready. Now we add the teams.",
             font_size=9, color=GREEN, bold=True)

# Phase 3: Month 2
add_card(slide, Inches(0.6), Inches(4.4), Inches(5.8), Inches(2.0))
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(1.5), Inches(0.3),
             "MONTH 2", font_size=14, color=GOLD, bold=True)
add_text_box(slide, Inches(2.3), Inches(4.5), Inches(3.5), Inches(0.3),
             "Scale All Teams (5 People)", font_size=14, color=BLACK, bold=True)
add_divider(slide, Inches(0.8), Inches(4.8), Inches(5.4))

m2_hires = [
    ("Delivery Manager", "Senior. By now first placements are landing."),
    ("2 Recruiters (Interns)", "Join AM team. Account Manager has jobs ready."),
    ("Training Sales", "Enrollment sales. Training Lead ready to delegate."),
    ("Training Coordinator", "Session logistics. Trainer is running batches."),
]
for j, (role, why) in enumerate(m2_hires):
    y = Inches(4.9) + j * Inches(0.35)
    add_text_box(slide, Inches(0.8), y, Inches(2.3), Inches(0.3),
                 role, font_size=10, color=BLACK, bold=True)
    add_text_box(slide, Inches(3.2), y, Inches(3), Inches(0.3),
                 why, font_size=9, color=CHARCOAL_600)

# Phase 4: Month 3-4
add_card(slide, Inches(6.8), Inches(4.4), Inches(5.8), Inches(2.0))
add_text_box(slide, Inches(7.0), Inches(4.5), Inches(1.5), Inches(0.3),
             "MONTH 3-4", font_size=14, color=GOLD, bold=True)
add_text_box(slide, Inches(8.5), Inches(4.5), Inches(3.5), Inches(0.3),
             "Full Team Operational", font_size=14, color=BLACK, bold=True)
add_divider(slide, Inches(7.0), Inches(4.8), Inches(5.4))

m3_items = [
    "All teams fully staffed and running",
    "Seniors leading their teams independently",
    "Interns at 70%+ productivity (Month 3)",
    "First training cohort graduates",
    "First placements generating revenue",
    "Evaluate: promote top interns, scale as needed",
]
add_bullet_list(slide, Inches(7.0), Inches(4.9), Inches(5.4), Inches(1.4),
                m3_items, font_size=10, color=CHARCOAL_600, bullet_char="\u2022", spacing=Pt(2))

add_text_box(slide, Inches(1), Inches(6.65), Inches(11), Inches(0.5),
             "The sequence matters: Revenue-generators first, then team leads, then scale with interns.\n"
             "Every intern has a senior already in place to train and manage them.",
             font_size=12, color=CHARCOAL_500, bold=True)
add_page_number(slide, 11, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 12: DAILY ACTIVITIES BY ROLE
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "DAILY OPERATIONS")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "Every Role Has Daily Targets", font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.6), Inches(9), Inches(0.3),
             "If you're not hitting daily numbers, you won't hit monthly targets. Track everything. Report everything.",
             font_size=14, color=CHARCOAL_600)

roles = [
    ("Training Sales", "\u2022 Respond to inquiries < 2 hours\n\u2022 20+ outbound calls/messages\n\u2022 Conduct discovery calls\n\u2022 3-5 enrollment closes/week\n\u2022 Update CRM every interaction"),
    ("Recruiters (AM Team)", "\u2022 Source from boards & LinkedIn\n\u2022 40+ outreach messages/day\n\u2022 8+ screening calls/day\n\u2022 Submit candidates to client jobs\n\u2022 Follow up on pending submissions"),
    ("Bench Sales Reps", "\u2022 Morning: Pool jobs by technology\n\u2022 Mid-day: Match & submit (15+/day)\n\u2022 Afternoon: Follow-up & prep\n\u2022 Evening: Log all activity\n\u2022 Coordinate vendor relationships"),
    ("OPT Recruiter", "\u2022 Source OPT boards & universities\n\u2022 30+ outreach messages/day\n\u2022 5+ screening calls/day\n\u2022 Feed leads to AM & Training\n\u2022 Community engagement daily"),
    ("BDM", "\u2022 10+ outbound prospecting/day\n\u2022 Follow up warm leads\n\u2022 LinkedIn engagement daily\n\u2022 3-5 discovery calls/week\n\u2022 Vendor portal registrations"),
    ("Delivery Manager", "\u2022 Check-in active consultants\n\u2022 Monitor performance signals\n\u2022 Client satisfaction calls weekly\n\u2022 Resolve issues < 48 hours\n\u2022 Renewal pipeline management"),
]

rw2 = Inches(3.7)
rh2 = Inches(1.8)
for i, (role, activities) in enumerate(roles):
    col = i % 3
    row = i // 3
    left = Inches(0.7) + col * (rw2 + Inches(0.15))
    top = Inches(2.1) + row * (rh2 + Inches(0.15))
    add_card(slide, left, top, rw2, rh2)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1),
                 rw2 - Inches(0.4), Inches(0.3), role, font_size=13, color=BLACK, bold=True)
    add_divider(slide, left + Inches(0.2), top + Inches(0.4), rw2 - Inches(0.4))
    add_text_box(slide, left + Inches(0.2), top + Inches(0.5),
                 rw2 - Inches(0.4), Inches(1.2), activities, font_size=9, color=CHARCOAL_600)

add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.7),
             "Senior Leads: Morning standup (15 min), review dashboards, 1:1s weekly, quality audits, strategy with Founder.\n"
             "Everyone: Daily standup (async or sync) - What I did, what I'm doing, blockers. No exceptions.",
             font_size=12, color=CHARCOAL_500, bold=True)
add_page_number(slide, 12, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 13: KEY METRICS & KPIs
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "KEY METRICS")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "Numbers That Matter", font_size=32, color=BLACK, bold=True)

metric_groups = [
    ("TRAINING", GREEN, [
        ("Batch Fill Rate", "100% (9 students)"),
        ("Student Retention", "90%+ complete"),
        ("Top Performer Rate", "20%+ score 90%+"),
        ("Time to 1st Interview", "Within 1 week"),
        ("Enrollment Rate", "30%+ of leads"),
    ]),
    ("RECRUITING", BLUE, [
        ("Outreach/Day", "40+ messages"),
        ("Screening Calls/Day", "8+ calls"),
        ("Qualified Leads/Week", "8+ combined"),
        ("Pipeline Accuracy", "80%+ convert"),
        ("Pool Freshness", "80%+ < 30 days"),
    ]),
    ("BENCH SALES", GOLD, [
        ("Submissions/Day", "15+ each"),
        ("Submit to Interview", "10%+ rate"),
        ("Interview to Offer", "30%+ rate"),
        ("Placements/Month", "2+ combined"),
        ("Bench Utilization", "80%+ in 30 days"),
    ]),
    ("DELIVERY", BLACK, [
        ("Client Satisfaction", "4.5+/5 rating"),
        ("Early Termination", "< 5% rate"),
        ("Renewal Rate", "80%+ contracts"),
        ("Extension Rate", "50%+ extended"),
        ("Issue Resolution", "< 48 hours"),
    ]),
]

mw, mh, mtop = Inches(2.8), Inches(3.2), Inches(1.9)
for i, (cat, cat_color, metrics) in enumerate(metric_groups):
    left = Inches(0.6) + i * (mw + Inches(0.2))
    add_card(slide, left, mtop, mw, mh)
    add_text_box(slide, left + Inches(0.2), mtop + Inches(0.15),
                 mw - Inches(0.4), Inches(0.3), cat, font_size=11, color=cat_color, bold=True)
    add_divider(slide, left + Inches(0.2), mtop + Inches(0.48), mw - Inches(0.4))
    for j, (metric, target) in enumerate(metrics):
        y = mtop + Inches(0.6) + j * Inches(0.48)
        add_text_box(slide, left + Inches(0.2), y, Inches(1.5), Inches(0.25),
                     metric, font_size=10, color=CHARCOAL_600)
        add_text_box(slide, left + Inches(1.6), y, Inches(1.1), Inches(0.25),
                     target, font_size=10, color=BLACK, bold=True, alignment=PP_ALIGN.RIGHT)

add_text_box(slide, Inches(1), Inches(5.4), Inches(11), Inches(1.0),
             "3 CRITICAL SUCCESS METRICS (Month 3-6):\n"
             "1. Batch fill rate 100% (9 students always enrolled)  |  "
             "2. Bench placements 2+/month  |  "
             "3. Client pipeline growing every week\n"
             "If we hit these three, everything else follows. Miss them, and we have a problem.",
             font_size=12, color=CHARCOAL_500, bold=True)
add_page_number(slide, 13, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 14: HOW WE WORK
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "HOW WE WORK")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "Tools, Rituals & Communication", font_size=32, color=BLACK, bold=True)

tools = [
    ("Slack", "Real-time communication.\nDaily updates, quick questions.\nRespond same day.", "Communication"),
    ("Linear", "Task & project tracking.\nAll work starts with a ticket.\nTrack progress daily.", "Project Mgmt"),
    ("InTime Platform", "Our own staffing software.\nCRM, ATS, Bench tracking.\nUse it for all operations.", "Operations"),
    ("Google Workspace", "Email, Docs, Sheets, Calendar.\nExternal communication.\nDocument everything.", "Collaboration"),
]

tw2, th2, ttop2 = Inches(2.7), Inches(1.65), Inches(1.7)
for i, (name, desc, category) in enumerate(tools):
    left = Inches(0.7) + i * (tw2 + Inches(0.2))
    add_card(slide, left, ttop2, tw2, th2)
    add_text_box(slide, left + Inches(0.2), ttop2 + Inches(0.12),
                 tw2 - Inches(0.4), Inches(0.3), name, font_size=14, color=BLACK, bold=True)
    add_text_box(slide, left + Inches(0.2), ttop2 + Inches(0.38),
                 tw2 - Inches(0.4), Inches(0.2), category, font_size=9, color=GOLD, bold=True)
    add_text_box(slide, left + Inches(0.2), ttop2 + Inches(0.65),
                 tw2 - Inches(0.4), Inches(0.9), desc, font_size=10, color=CHARCOAL_600)

# Rituals
add_text_box(slide, Inches(0.8), Inches(3.6), Inches(5), Inches(0.3),
             "RITUALS & CADENCES", font_size=12, color=BLACK, bold=True)
rituals = [
    "Daily Standup (15 min): What I did, what I'm doing, blockers",
    "Weekly Team Sync (30-60 min): Progress, priorities, blockers",
    "Weekly 1:1 with Manager (30 min): Coaching, feedback, support",
    "Bi-Weekly Sprint Planning: Next 2 weeks of work",
    "Monthly All Hands (60 min): Company updates, wins, learnings",
    "Monthly Metrics Review: How we're doing against targets",
    "Quarterly OKR Setting: Goals for next quarter",
]
add_bullet_list(slide, Inches(0.8), Inches(3.95), Inches(5.5), Inches(2.8),
                rituals, font_size=10, color=CHARCOAL_600, bullet_char="\u2022", spacing=Pt(3))

# Comm rules
add_text_box(slide, Inches(6.8), Inches(3.6), Inches(5), Inches(0.3),
             "COMMUNICATION RULES", font_size=12, color=BLACK, bold=True)
comm = [
    "Over-communicate context. Assume others don't know.",
    "Be explicit: FYI vs Need by EOD vs Blocking.",
    "Close the loop: Acknowledge, respond, resolve.",
    "Bad news travels fast. Good news can wait.",
    "Escalate early, never late. Flag blockers immediately.",
    "If it wasn't written down, it didn't happen.",
    "Disagree openly, commit fully, follow through completely.",
]
add_bullet_list(slide, Inches(6.8), Inches(3.95), Inches(5.5), Inches(2.8),
                comm, font_size=10, color=CHARCOAL_600, bullet_char="\u2022", spacing=Pt(3))

add_page_number(slide, 14, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 15: DO's & DON'Ts
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "EXPECTATIONS")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "What I Expect From Everyone", font_size=32, color=BLACK, bold=True)
add_text_box(slide, Inches(1), Inches(1.6), Inches(9), Inches(0.3),
             "I have limited time. You are my hands. This is what matters to me.",
             font_size=15, color=CHARCOAL_600)

do_items = [
    "Take ownership - see problem, fix it or flag it",
    "Hit daily numbers - every role has daily targets",
    "Communicate proactively - done, blocked, need help",
    "Ask questions early - don't wait days when stuck",
    "Follow SOPs and processes - they exist for a reason",
    "Track everything in the system - no mental notes",
    "Be accountable - own results, good or bad",
    "Learn fast - US staffing is learn-on-the-job",
    "Build relationships - candidates, clients, team",
    "Think like an owner - this is OUR company",
]

dont_items = [
    "Don't wait to be told what to do",
    "Don't hide bad news - share it immediately",
    "Don't skip logging activity in CRM",
    "Don't make promises you can't keep",
    "Don't blame others - own the outcome",
    "Don't go silent - over-communicate",
    "Don't cut corners on quality",
    "Don't treat this as just a job",
]

add_card(slide, Inches(0.7), Inches(2.1), Inches(5.8), Inches(4.6), bg_color=RGBColor(0xF0, 0xF9, 0xF4))
add_text_box(slide, Inches(1), Inches(2.2), Inches(2), Inches(0.3), "DO", font_size=16, color=GREEN, bold=True)
add_bullet_list(slide, Inches(1), Inches(2.6), Inches(5.2), Inches(3.8),
                do_items, font_size=11, color=CHARCOAL_700, bullet_char="\u2714", spacing=Pt(3))

add_card(slide, Inches(6.8), Inches(2.1), Inches(5.8), Inches(4.6), bg_color=RGBColor(0xFE, 0xF2, 0xF2))
add_text_box(slide, Inches(7.1), Inches(2.2), Inches(2), Inches(0.3), "DON'T", font_size=16, color=RED, bold=True)
add_bullet_list(slide, Inches(7.1), Inches(2.6), Inches(5.2), Inches(3.8),
                dont_items, font_size=11, color=CHARCOAL_700, bullet_char="\u2718", spacing=Pt(3))

add_page_number(slide, 15, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 16: WHAT I EXPECT FROM OPERATIONS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "OPERATIONS EXPECTATIONS", color=CHARCOAL_500)
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "What I Expect From My Operations Person", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.7), Inches(9), Inches(0.4),
             "You are my right hand. When I'm not available, you keep things running. Here's what that means.",
             font_size=15, color=CHARCOAL_400)

# 6 pillars of ops expectations
ops_expects = [
    ("Be My Eyes", "Know what's happening\nacross every team, every day.\nI should never be surprised.\nYou see everything before I do.",
     "Daily dashboard review + summary to me"),
    ("Keep People Accountable", "Track daily activity numbers.\nFollow up on missed targets.\nEscalate before things break.\nNo one should coast on your watch.",
     "Weekly performance report per person"),
    ("Run the Cadences", "Own daily standups, weekly syncs.\nMake sure meetings happen.\nCapture action items and follow up.\nRituals don't run themselves.",
     "100% meeting completion, notes shared"),
    ("Handle the Details", "Student onboarding, session logistics,\ntimesheet tracking, vendor paperwork.\nInvoicing, payments, compliance.\nNothing falls through the cracks.",
     "Zero missed deadlines or payments"),
    ("Solve Before Escalating", "Try to fix it yourself first.\nBring solutions, not just problems.\nWhen you escalate, give me context\nand your recommendation.",
     "80% resolved without escalation"),
    ("Build Systems", "Create checklists, templates, SOPs.\nDocument processes so anyone can follow.\nAutomate what's repetitive.\nMake the operation run without you.",
     "Everything documented and repeatable"),
]

ow, oh = Inches(3.7), Inches(1.9)
for i, (title, desc, metric) in enumerate(ops_expects):
    col = i % 3
    row = i // 3
    left = Inches(0.6) + col * (ow + Inches(0.15))
    top = Inches(2.3) + row * (oh + Inches(0.15))
    add_card(slide, left, top, ow, oh, bg_color=DARK_CARD, border_color=CHARCOAL_700)

    add_text_box(slide, left + Inches(0.2), top + Inches(0.1),
                 ow - Inches(0.4), Inches(0.25), title, font_size=13, color=GOLD, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.4),
                 ow - Inches(0.4), Inches(1.0), desc, font_size=10, color=CHARCOAL_400)
    add_text_box(slide, left + Inches(0.2), top + Inches(1.5),
                 ow - Inches(0.4), Inches(0.3), metric, font_size=8, color=GREEN, bold=True)

add_text_box(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.5),
             "Bottom line: If the operation runs smoothly without me checking in every hour, you're doing your job.\n"
             "If I have to ask 'what's happening?' - something is wrong.",
             font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 16, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 17: FIRST 90 DAYS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "YOUR FIRST 90 DAYS")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "From Learning to Leading", font_size=32, color=BLACK, bold=True)

phases_90 = [
    ("Days 1-7", "ONBOARD", "Learn the Basics", [
        "HR paperwork, system access",
        "Read Company Bible cover to cover",
        "Meet team lead and all team members",
        "Read US IT Staffing 101 guide",
        "Understand the 7 funnels model",
        "Shadow team members in daily work",
    ]),
    ("Days 8-30", "LEARN", "Understand Business", [
        "Complete role-specific handbook",
        "Observe 5+ real work scenarios",
        "Practice tasks with supervision",
        "Daily check-in with manager",
        "Complete first tasks independently",
        "30-day review with manager",
    ]),
    ("Days 31-60", "BUILD", "Start Contributing", [
        "Handle 80% of tasks independently",
        "Carry increasing workload",
        "Solve problems before escalating",
        "Suggest process improvements",
        "Build cross-functional relationships",
        "60-day review: set Month 3 goals",
    ]),
    ("Days 61-90", "PERFORM", "Work Independently", [
        "Meet or exceed your role's KPIs",
        "No supervision needed",
        "Mentor newer team members",
        "Take on stretch assignments",
        "Lead small improvement projects",
        "90-day review: career dev plan",
    ]),
]

pw_90, ph_90, ptop_90 = Inches(2.7), Inches(4.3), Inches(1.8)
phase_colors = [BLUE, GREEN, GOLD, BLACK]

for i, (timeline, phase, title, items) in enumerate(phases_90):
    left = Inches(0.7) + i * (pw_90 + Inches(0.25))
    add_card(slide, left, ptop_90, pw_90, ph_90)
    add_text_box(slide, left + Inches(0.2), ptop_90 + Inches(0.12),
                 Inches(1.2), Inches(0.25), timeline, font_size=10, color=CHARCOAL_500)
    add_text_box(slide, left + pw_90 - Inches(1.2), ptop_90 + Inches(0.12),
                 Inches(1), Inches(0.25), phase, font_size=10, color=phase_colors[i],
                 bold=True, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, left + Inches(0.2), ptop_90 + Inches(0.42),
                 pw_90 - Inches(0.4), Inches(0.3), title, font_size=13, color=BLACK, bold=True)
    add_divider(slide, left + Inches(0.2), ptop_90 + Inches(0.78), pw_90 - Inches(0.4))
    add_bullet_list(slide, left + Inches(0.2), ptop_90 + Inches(0.9),
                    pw_90 - Inches(0.4), ph_90 - Inches(1.1),
                    items, font_size=10, color=CHARCOAL_600, bullet_char="\u2022", spacing=Pt(3))

add_text_box(slide, Inches(1), Inches(6.35), Inches(11), Inches(0.5),
             "By Month 4, interns reach 100% productivity. We evaluate: retain, promote, or replace.\n"
             "Your first 90 days determine your trajectory. Invest in learning now - it compounds.",
             font_size=12, color=CHARCOAL_500, bold=True)
add_page_number(slide, 17, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 18: SUCCESS CRITERIA
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "SUCCESS CRITERIA")
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "What Good Looks Like at Day 30, 60, 90",
             font_size=32, color=BLACK, bold=True)

role_criteria = [
    ("Recruiters", [
        ("Day 30", "10 candidates sourced/week\n3 screens/week\nUnderstand process"),
        ("Day 60", "25 candidates sourced/week\n10 screens/week\nExecute independently"),
        ("Day 90", "50 candidates sourced/week\n15+ screens/week\nImprove process"),
    ]),
    ("Account Managers", [
        ("Day 30", "Know all clients\n1 submission/week\nUnderstand client needs"),
        ("Day 60", "Manage 2-3 clients\n5 submissions/week\nBuild relationships"),
        ("Day 90", "Full client load\n10+ submissions/week\nGrow accounts"),
    ]),
    ("Bench Sales", [
        ("Day 30", "Know all consultants\n10 vendor contacts\nUnderstand process"),
        ("Day 60", "Active marketing\n30+ vendor relationships\nExecute daily"),
        ("Day 90", "1+ placement/month\nConsistent placements\nOptimize process"),
    ]),
    ("Operations", [
        ("Day 30", "Know all processes\nMinimal errors\nSupervised work"),
        ("Day 60", "Handle routine tasks\nQuality work\nMostly independent"),
        ("Day 90", "Full ownership\nProcess improvements\nFully independent"),
    ]),
]

rw3, rh3, rtop3 = Inches(2.8), Inches(3.5), Inches(1.8)
for i, (role, milestones) in enumerate(role_criteria):
    left = Inches(0.6) + i * (rw3 + Inches(0.2))
    add_card(slide, left, rtop3, rw3, rh3)
    add_text_box(slide, left + Inches(0.2), rtop3 + Inches(0.12),
                 rw3 - Inches(0.4), Inches(0.3), role, font_size=14, color=BLACK, bold=True)
    add_divider(slide, left + Inches(0.2), rtop3 + Inches(0.45), rw3 - Inches(0.4))
    day_colors = [BLUE, AMBER, GREEN]
    for j, (day, criteria) in enumerate(milestones):
        y = rtop3 + Inches(0.55) + j * Inches(0.95)
        add_text_box(slide, left + Inches(0.2), y,
                     Inches(0.7), Inches(0.2), day, font_size=9, color=day_colors[j], bold=True)
        add_text_box(slide, left + Inches(0.2), y + Inches(0.2),
                     rw3 - Inches(0.4), Inches(0.7), criteria, font_size=9, color=CHARCOAL_600)

add_text_box(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.6),
             "Performance = IMPACT, not hours. What you accomplished > how long you sat at your desk.\n"
             "Check-ins: Daily (Week 1) | Weekly (Month 1) | Bi-weekly (Month 2) | Monthly (Month 3+)",
             font_size=12, color=CHARCOAL_500, bold=True)
add_page_number(slide, 18, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 19: ESSENTIAL READING & GROWTH
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gold_accent_line(slide)

add_section_label(slide, Inches(1), Inches(0.6), "RESOURCES & GROWTH", color=CHARCOAL_500)

add_text_box(slide, Inches(1), Inches(1.0), Inches(5), Inches(0.6),
             "Essential Reading", font_size=28, color=WHITE, bold=True)

readings = [
    ("Company Bible", "Values, culture, how we work, hire, decide"),
    ("Business Model", "7 funnels, roles, accountability, metrics"),
    ("Phase 1 Ops Plan", "16-person team, daily activities, financials"),
    ("Financial Model", "Revenue, unit economics, P&L, break-even"),
    ("First 90 Days", "Your onboarding roadmap day-by-day"),
    ("US IT Staffing 101", "Industry basics, visa types, terminology"),
    ("Your Role Handbook", "SOPs and expectations for your role"),
]

for i, (doc, desc) in enumerate(readings):
    y = Inches(1.7) + i * Inches(0.55)
    add_text_box(slide, Inches(1.2), y, Inches(2.2), Inches(0.3),
                 doc, font_size=12, color=GOLD, bold=True)
    add_text_box(slide, Inches(3.5), y, Inches(3), Inches(0.3),
                 desc, font_size=10, color=CHARCOAL_400)

# Growth path
add_text_box(slide, Inches(7.5), Inches(1.0), Inches(5), Inches(0.6),
             "Your Growth Path", font_size=28, color=WHITE, bold=True)

growth = [
    ("Month 1-3", "LEARNING",
     "Learn the ropes, shadow seniors,\nhit basic daily activity targets."),
    ("Month 4-6", "PRODUCTIVE",
     "Fully independent, meeting KPIs,\ncarry full workload, prove yourself."),
    ("Month 7-12", "TEAM PLAYER",
     "Exceed targets, help newer members,\nearn promotion based on performance."),
    ("Year 2+", "FUTURE LEAD",
     "Lead a pod of interns yourself,\nown a function end-to-end."),
]

gw2, gh2 = Inches(5), Inches(1.1)
for i, (timeline, stage, desc) in enumerate(growth):
    top = Inches(1.75) + i * (gh2 + Inches(0.1))
    add_card(slide, Inches(7.3), top, gw2, gh2, bg_color=DARK_CARD, border_color=CHARCOAL_700)
    stage_colors = [BLUE, GREEN, GOLD, WHITE]
    add_text_box(slide, Inches(7.5), top + Inches(0.08),
                 Inches(1.5), Inches(0.25), timeline, font_size=10, color=CHARCOAL_400)
    add_text_box(slide, Inches(9.2), top + Inches(0.08),
                 Inches(2.8), Inches(0.25), stage, font_size=10, color=stage_colors[i], bold=True,
                 alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(7.5), top + Inches(0.38),
                 gw2 - Inches(0.4), Inches(0.65), desc, font_size=9, color=CHARCOAL_400)

add_text_box(slide, Inches(1), Inches(6.4), Inches(11), Inches(0.5),
             "All documents are in docs/functional/ directory. Your manager will guide you through them.\n"
             "Performance is rewarded. Top performers get promoted and earn more responsibility.",
             font_size=12, color=CHARCOAL_500)
add_page_number(slide, 19, TOTAL_SLIDES)

# ══════════════════════════════════════════════
# SLIDE 20: CLOSING
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_gold_accent_line(slide)

add_text_box(slide, Inches(1.5), Inches(1.3), Inches(10), Inches(1.5),
             "\"The recruiter recruits.\nThe account manager manages accounts.\nThe operator handles exceptions.\nEveryone does their actual job.\nNobody fights the software.\"",
             font_size=26, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.1), Inches(10), Inches(0.4),
             "- The InTime Vision", font_size=14, color=GOLD, alignment=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    left=Inches(5.5), top=Inches(3.8), width=Inches(2), height=Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = GOLD
shape.line.fill.background()

add_text_box(slide, Inches(2), Inches(4.2), Inches(9), Inches(1.0),
             "You're joining at the right time.\nThe model is clear. The team is forming.\nNow we need people who will execute.",
             font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.5),
             "Your commitment + our system = success for everyone.",
             font_size=16, color=CHARCOAL_400, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(6.3), Inches(9), Inches(0.3),
             "I N T I M E  |  US IT Staffing  |  Hyderabad Operations  |  February 2026",
             font_size=11, color=CHARCOAL_500, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 20, TOTAL_SLIDES)

# ──────────────────────────────────────────────
# SAVE
# ──────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)),
                           "InTime_New_Hire_Onboarding.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {TOTAL_SLIDES}")
