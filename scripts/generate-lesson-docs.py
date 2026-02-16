#!/usr/bin/env python3
"""
Generate Word documents per lesson from PPT content for manual review/editing.
Each doc includes slide content, speaker notes, video references, and assignment links.

Dependencies: pip install python-pptx python-docx
"""

import json
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn

def sanitize_text(text):
    """Remove control characters that are invalid in XML (used by python-docx)."""
    if not text:
        return text
    # Remove NULL bytes and control chars (0x00-0x08, 0x0B-0x0C, 0x0E-0x1F) except tab/newline/CR
    return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)


MATERIAL_DIR = "/Users/sumanthrajkumarnagolu/Documents/Material"
VIDEO_DIR = "/Users/sumanthrajkumarnagolu/Projects/intime-v3/public/academy/guidewire/videos"
ASSIGNMENT_MANIFEST = "/Users/sumanthrajkumarnagolu/Projects/intime-v3/public/academy/guidewire/assignment-manifest.json"
OUTPUT_DIR = "/Users/sumanthrajkumarnagolu/Projects/intime-v3/output/lesson-docs"

CHAPTER_MAP = {
    "Chapter 1 - Guidewire Cloud Overview": {"id": 1, "slug": "ch01"},
    "Chapter 2 - Surepath Overview": {"id": 2, "slug": "ch02"},
    "Chapter 3 - InsuranceSuite Implementation Tools": {"id": 3, "slug": "ch03"},
    "Chapter 4 - PolicyCenter Introduction": {"id": 4, "slug": "ch04"},
    "Chapter 5 - Claim Center Introduction": {"id": 5, "slug": "ch05"},
    "Chapter 6 - Billing Center Introduction": {"id": 6, "slug": "ch06"},
    "Chapter 7 - InsuranceSuite Developer Fundamentals": {"id": 7, "slug": "ch07"},
    "Chapter 8 - policy center configuration": {"id": 8, "slug": "ch08"},
    "Chapter 9 - ClaimCenter Configuration": {"id": 9, "slug": "ch09"},
    "Chapter 10 - BillingCenter Configuration": {"id": 10, "slug": "ch10"},
    "Chapter 11 - Introduction to Integration": {"id": 11, "slug": "ch11"},
    "Chapter 12 - Advanced product Designer": {"id": 12, "slug": "ch12"},
    "Chapter 13 - Rating Introduction": {"id": 13, "slug": "ch13"},
    "Chapter 14 - Rating Configuration": {"id": 14, "slug": "ch14"},
}

# Assignment folder -> chapter slug mapping
ASSIGNMENT_CHAPTER_MAP = {
    "ch04": "PolicyCenter Introduction",
    "ch05": "ClaimCenter Introduction",
    "ch06": "BillingCenter Introduction",
    "ch07": "InsuranceSuite Fundamentals",
    "ch08": "PolicyCenter Configuration",
    "ch09": "ClaimCenter Configuration",
    "ch10": "BillingCenter Configuration",
    "ch11": "InsuranceSuite Integration",
}


def set_cell_shading(cell, color_hex):
    """Set cell background color."""
    shading_elm = cell._tc.get_or_add_tcPr()
    shading = shading_elm.makeelement(
        qn('w:shd'),
        {qn('w:fill'): color_hex, qn('w:val'): 'clear'}
    )
    shading_elm.append(shading)


def add_shaded_paragraph(doc, text, shade_color="D9D9D9", font_size=9, italic=True):
    """Add a paragraph with gray shading (used for speaker notes)."""
    para = doc.add_paragraph()
    para.paragraph_format.space_before = Pt(4)
    para.paragraph_format.space_after = Pt(4)
    run = para.add_run(sanitize_text(text))
    run.font.size = Pt(font_size)
    run.font.italic = italic
    run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # Add shading to paragraph
    pPr = para._p.get_or_add_pPr()
    shading = pPr.makeelement(
        qn('w:shd'),
        {qn('w:fill'): shade_color, qn('w:val'): 'clear'}
    )
    pPr.append(shading)
    return para


def get_lesson_number(folder_name, pptx_filename):
    """Extract lesson number from folder or filename."""
    patterns = [
        r'IS_Fund[_\s]*(\d+)',
        r'In_policy[_\s]*(\d+)',
        r'In_Claim[_\s]*(\d+)',
        r'BillingCenter[_\s]*(\d+)',
        r'In_Integration[_\s]*(\d+)',
        r'Ra_Intro[_\s]*(\d+)',
        r'Ra_Conf[_\s]*(\d+)',
        r'PP[_\s]*(\d+)',
        r'CC[_\s]*(\d+)',
        r'^(\d+)\s*[-–]',
    ]
    for source in [folder_name, pptx_filename]:
        for pattern in patterns:
            m = re.search(pattern, source)
            if m:
                return int(m.group(1))
    return 0


def extract_lesson_title(pptx_filename):
    """Extract a clean lesson title from the PPTX filename."""
    name = os.path.splitext(pptx_filename)[0]
    patterns = [
        r'^IS_Fund[_\s]*\d+[_\s]*[-–]\s*',
        r'^IS_Fund[_\s]*\d+[_\s]*',
        r'^In_policy[_\s]*\d+[_\s]*[-–]?\s*',
        r'^In_policy[_\s]*\d+[_\s]*',
        r'^PC_Intro[_\s]*\d+[_\s]*[-–]?\s*',
        r'^IS_Claim[_\s]*\d+[_\s]*[-–]?\s*',
        r'^IS_Claim[_\s]*__?\d+[_\s]*[-–]?\s*',
        r'^BillingCenter[_\s]*\d+[_\s]*[-–]?\s*',
        r'^In_Integration[_\s]*\d+[_\s]*[-–]?\s*',
        r'^Ra_Intro[_\s]*\d+[_\s]*[-–]?\s*',
        r'^Ra_Conf[_\s]*\d+[_\s]*[-–]?\s*',
        r'^PP[_\s]*\d+[_\s]*[-–]?\s*Configuration[_\s]*[-–]?\s*',
        r'^PP[_\s]*\d+[_\s]*[-–]?\s*',
        r'^CC[_\s]*\d+[_\s]*[-–]?\s*',
        r'^\d+[_\s]*[-–]\s*',
    ]
    clean = name
    for pattern in patterns:
        clean = re.sub(pattern, '', clean, flags=re.IGNORECASE)
    clean = clean.strip(' -_')
    return clean if clean else name


def find_videos_for_lesson(chapter_slug, lesson_folder):
    """Find video files for a lesson from the videos directory."""
    video_chapter_dir = os.path.join(VIDEO_DIR, chapter_slug)
    if not os.path.isdir(video_chapter_dir):
        return []

    # Also check the source material lesson folder for videos
    videos = []
    all_videos = sorted(os.listdir(video_chapter_dir))

    if lesson_folder:
        # Match videos whose names contain the lesson folder prefix
        folder_lower = lesson_folder.lower().replace(' ', '_')
        for v in all_videos:
            if v.lower().endswith(('.mp4', '.mkv', '.avi', '.wmv')):
                # Match by lesson folder prefix
                v_lower = v.lower()
                if folder_lower in v_lower or v_lower.startswith(folder_lower.split('_')[0] if '_' in folder_lower else folder_lower[:5]):
                    videos.append(v)

    # If no matches, try matching by lesson number pattern
    if not videos and lesson_folder:
        num = get_lesson_number(lesson_folder, "")
        if num > 0:
            num_str = f"{num:02d}"
            for v in all_videos:
                if v.lower().endswith(('.mp4', '.mkv', '.avi', '.wmv')):
                    if f"_{num_str}_" in v or f"_{num_str}." in v or v.startswith(f"{num_str}_") or v.startswith(f"{num_str} "):
                        videos.append(v)

    return videos


def find_lesson_folders(chapter_dir):
    """Find all lesson folders in a chapter directory, handling nested structures."""
    lesson_folders = []

    for item in sorted(os.listdir(chapter_dir)):
        item_path = os.path.join(chapter_dir, item)
        if not os.path.isdir(item_path) or item in ['.DS_Store', 'backups', 'videos']:
            continue

        has_ppt = any(
            f.endswith(('.pptx', '.pptm')) and not f.startswith('~$')
            for f in os.listdir(item_path)
            if os.path.isfile(os.path.join(item_path, f))
        )

        if has_ppt:
            lesson_folders.append((item, item_path))
        else:
            # Check subdirectories (nested structure like Ch8, Ch11)
            for sub_item in sorted(os.listdir(item_path)):
                sub_path = os.path.join(item_path, sub_item)
                if os.path.isdir(sub_path):
                    has_sub_ppt = any(
                        f.endswith(('.pptx', '.pptm')) and not f.startswith('~$')
                        for f in os.listdir(sub_path)
                        if os.path.isfile(os.path.join(sub_path, f))
                    )
                    if has_sub_ppt:
                        lesson_folders.append((sub_item, sub_path))

    return lesson_folders


def generate_lesson_doc(pptx_path, chapter_name, chapter_slug, lesson_folder, lesson_number, assignment_data):
    """Generate a Word document from a PPTX file."""
    filename = os.path.basename(pptx_path)
    title = extract_lesson_title(filename)

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"    ERROR opening {filename}: {e}")
        return None

    # Create Word document
    doc = Document()

    # Set default font
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)

    # ===== HEADER: Chapter + Lesson info =====
    header_para = doc.add_heading(level=1)
    run = header_para.add_run(f"Chapter {chapter_slug.replace('ch', '')}: {chapter_name}")
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

    lesson_para = doc.add_heading(level=2)
    run = lesson_para.add_run(f"Lesson {lesson_number}: {title}")
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Source file reference
    meta = doc.add_paragraph()
    meta.paragraph_format.space_after = Pt(2)
    run = meta.add_run(f"Source: {filename}")
    run.font.size = Pt(8)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # ===== VIDEO REFERENCES =====
    videos = find_videos_for_lesson(chapter_slug, lesson_folder)
    if videos:
        doc.add_heading("Video References", level=3)
        for i, video in enumerate(videos, 1):
            vp = doc.add_paragraph(style='List Number')
            run = vp.add_run(f"[Video {i}: {video}]")
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)
    else:
        doc.add_heading("Video References", level=3)
        p = doc.add_paragraph()
        run = p.add_run("No videos found for this lesson.")
        run.font.italic = True
        run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    doc.add_paragraph()  # spacer

    # ===== SLIDE CONTENT =====
    doc.add_heading("Slide Content", level=2)

    quiz_questions = []
    total_slides = len(prs.slides)

    for slide_idx, slide in enumerate(prs.slides, 1):
        # --- Slide header ---
        slide_heading = doc.add_heading(level=3)
        slide_title = ""

        # Get title
        for shape in slide.shapes:
            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                ph = shape.placeholder_format
                if ph and ph.idx == 0 and shape.has_text_frame:
                    slide_title = shape.text_frame.text.strip()
                    break
            if not slide_title and shape.has_text_frame and 'title' in (shape.name or '').lower():
                slide_title = shape.text_frame.text.strip()

        if not slide_title:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide_title = shape.text_frame.text.strip()
                    break

        slide_title = sanitize_text(slide_title)
        run = slide_heading.add_run(f"Slide {slide_idx}/{total_slides}: {slide_title or '(Untitled)'}")
        run.font.size = Pt(12)

        # --- Slide body content ---
        for shape in slide.shapes:
            # Skip title placeholder (already used)
            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                ph = shape.placeholder_format
                if ph and ph.idx == 0:
                    continue

            # Tables
            if shape.has_table:
                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)

                if rows > 0 and cols > 0:
                    doc_table = doc.add_table(rows=rows, cols=cols, style='Table Grid')
                    for r_idx, row in enumerate(table.rows):
                        for c_idx, cell in enumerate(row.cells):
                            doc_cell = doc_table.cell(r_idx, c_idx)
                            doc_cell.text = sanitize_text(cell.text.strip())
                            # Bold header row
                            if r_idx == 0:
                                for paragraph in doc_cell.paragraphs:
                                    for run in paragraph.runs:
                                        run.bold = True
                                set_cell_shading(doc_cell, "E8E8E8")
                    doc.add_paragraph()  # spacer after table
                continue

            # Images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                p = doc.add_paragraph()
                run = p.add_run("[Image placeholder]")
                run.font.italic = True
                run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
                continue

            # Charts
            if shape.has_chart:
                p = doc.add_paragraph()
                run = p.add_run("[Chart placeholder]")
                run.font.italic = True
                run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
                continue

            # Text content
            if shape.has_text_frame:
                # Skip if this is just the title repeated
                if shape.text_frame.text.strip() == slide_title:
                    continue

                for para in shape.text_frame.paragraphs:
                    text = sanitize_text(para.text.strip())
                    if not text:
                        continue

                    level = para.level if hasattr(para, 'level') else 0
                    is_bold = any(run.font.bold for run in para.runs if run.font.bold is not None)

                    # Check for quiz/checkpoint content
                    text_lower = text.lower()
                    is_quiz = any(kw in text_lower for kw in [
                        'quiz', 'checkpoint', 'question:', 'true or false',
                        'which of the following', 'select the correct',
                        'what is the', 'how do you', 'assessment',
                    ])

                    if is_quiz:
                        quiz_questions.append({
                            "slide": slide_idx,
                            "text": text,
                        })
                        p = doc.add_paragraph()
                        run = p.add_run(f"[Quiz Question] {text}")
                        run.bold = True
                        run.font.color.rgb = RGBColor(0xCC, 0x66, 0x00)
                    elif is_bold and level == 0:
                        p = doc.add_paragraph()
                        run = p.add_run(text)
                        run.bold = True
                    else:
                        indent = "  " * level
                        bullet_char = "  " if level == 0 else "    " * level
                        if level > 0:
                            p = doc.add_paragraph(style='List Bullet')
                        else:
                            p = doc.add_paragraph()
                        run = p.add_run(text)
                        run.font.size = Pt(10 if level > 0 else 11)

        # --- Speaker Notes ---
        if slide.has_notes_slide:
            notes_text = sanitize_text(slide.notes_slide.notes_text_frame.text.strip())
            if notes_text:
                add_shaded_paragraph(
                    doc,
                    f"Speaker Notes: {notes_text}",
                    shade_color="E8E8E8",
                    font_size=9,
                    italic=True,
                )

    # ===== QUIZ SECTION (consolidated) =====
    if quiz_questions:
        doc.add_page_break()
        doc.add_heading("Quiz / Checkpoint Questions", level=2)
        for i, q in enumerate(quiz_questions, 1):
            p = doc.add_paragraph()
            run = p.add_run(f"Q{i} (Slide {q['slide']}): ")
            run.bold = True
            run = p.add_run(q["text"])
            run.font.size = Pt(10)

    # ===== ASSIGNMENT REFERENCE =====
    if assignment_data and chapter_slug in assignment_data.get("chapters", {}):
        chapter_assignments = assignment_data["chapters"][chapter_slug]["assignments"]
        # Find assignment matching this lesson number
        matching = [a for a in chapter_assignments if a["number"] == lesson_number]
        if matching:
            doc.add_heading("Assignment", level=2)
            for a in matching:
                p = doc.add_paragraph()
                run = p.add_run(f"[Assignment {a['number']:02d}]: {a['title']}")
                run.font.color.rgb = RGBColor(0x00, 0x88, 0x00)
                run.bold = True
                p2 = doc.add_paragraph()
                run2 = p2.add_run(f"File: {a['originalFilename']}")
                run2.font.size = Pt(9)
                run2.font.italic = True

    # ===== SUMMARY FOOTER =====
    doc.add_paragraph()  # spacer
    footer = doc.add_paragraph()
    footer.paragraph_format.space_before = Pt(12)
    run = footer.add_run(f"--- End of Lesson {lesson_number}: {title} | {total_slides} slides | {len(videos)} videos ---")
    run.font.size = Pt(8)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER

    return doc, title


def process_chapter(chapter_folder, config, assignment_data):
    """Process all PPTX files in a chapter and generate Word docs."""
    chapter_dir = os.path.join(MATERIAL_DIR, chapter_folder)
    chapter_id = config["id"]
    chapter_slug = config["slug"]
    chapter_name = chapter_folder.replace(f"Chapter {chapter_id} - ", "")

    if not os.path.isdir(chapter_dir):
        print(f"  SKIP: {chapter_folder} (not found)")
        return 0

    out_dir = os.path.join(OUTPUT_DIR, chapter_slug)
    os.makedirs(out_dir, exist_ok=True)

    print(f"\n{'=' * 60}")
    print(f"Chapter {chapter_id}: {chapter_name}")
    print(f"{'=' * 60}")

    lesson_folders = find_lesson_folders(chapter_dir)
    docs_generated = 0

    # Process lesson folders
    for folder_name, folder_path in lesson_folders:
        pptx_files = sorted([
            f for f in os.listdir(folder_path)
            if f.endswith(('.pptx', '.pptm')) and not f.startswith('~$')
        ])

        for pptx in pptx_files:
            pptx_path = os.path.join(folder_path, pptx)
            lesson_number = get_lesson_number(folder_name, pptx)
            print(f"  Lesson {lesson_number:02d}: {pptx}")

            result = generate_lesson_doc(
                pptx_path, chapter_name, chapter_slug,
                folder_name, lesson_number, assignment_data
            )

            if result:
                doc, title = result
                safe_title = re.sub(r'[^\w\s-]', '', title)[:50].strip()
                doc_filename = f"Lesson_{lesson_number:02d}_{safe_title}.docx"
                doc_path = os.path.join(out_dir, doc_filename)
                doc.save(doc_path)
                docs_generated += 1
                print(f"    -> {doc_filename}")

    # Process root-level PPTs
    root_ppts = sorted([
        f for f in os.listdir(chapter_dir)
        if f.endswith(('.pptx', '.pptm')) and not f.startswith('~$')
        and os.path.isfile(os.path.join(chapter_dir, f))
    ])

    for pptx in root_ppts:
        pptx_path = os.path.join(chapter_dir, pptx)
        lesson_number = get_lesson_number("", pptx)
        print(f"  Lesson {lesson_number:02d}: {pptx} (root)")

        result = generate_lesson_doc(
            pptx_path, chapter_name, chapter_slug,
            None, lesson_number, assignment_data
        )

        if result:
            doc, title = result
            safe_title = re.sub(r'[^\w\s-]', '', title)[:50].strip()
            doc_filename = f"Lesson_{lesson_number:02d}_{safe_title}.docx"
            doc_path = os.path.join(out_dir, doc_filename)
            doc.save(doc_path)
            docs_generated += 1
            print(f"    -> {doc_filename}")

    return docs_generated


def process_docx_chapter(chapter_folder, config, assignment_data):
    """Generate lesson docs from DOCX source files (for chapters like Ch12)."""
    chapter_dir = os.path.join(MATERIAL_DIR, chapter_folder)
    chapter_id = config["id"]
    chapter_slug = config["slug"]
    chapter_name = chapter_folder.replace(f"Chapter {chapter_id} - ", "")

    if not os.path.isdir(chapter_dir):
        print(f"  SKIP: {chapter_folder} (not found)")
        return 0

    from docx import Document as ReadDocument

    out_dir = os.path.join(OUTPUT_DIR, chapter_slug)
    os.makedirs(out_dir, exist_ok=True)

    print(f"\n{'=' * 60}")
    print(f"Chapter {chapter_id}: {chapter_name} (DOCX source)")
    print(f"{'=' * 60}")

    docx_files = sorted([
        f for f in os.listdir(chapter_dir)
        if f.endswith('.docx') and not f.startswith('~$')
    ])

    docs_generated = 0

    for i, docx_file in enumerate(docx_files, 1):
        docx_path = os.path.join(chapter_dir, docx_file)
        print(f"  Lesson {i:02d}: {docx_file}")

        try:
            src_doc = ReadDocument(docx_path)
        except Exception as e:
            print(f"    ERROR: {e}")
            continue

        # Create output Word doc
        doc = Document()
        style = doc.styles['Normal']
        style.font.name = 'Calibri'
        style.font.size = Pt(11)

        title = os.path.splitext(docx_file)[0]
        title = re.sub(r'^\d+\s*[-–]\s*', '', title).strip()

        header_para = doc.add_heading(level=1)
        run = header_para.add_run(f"Chapter {chapter_slug.replace('ch', '')}: {chapter_name}")
        run.font.size = Pt(18)

        lesson_para = doc.add_heading(level=2)
        run = lesson_para.add_run(f"Lesson {i}: {title}")
        run.font.size = Pt(14)

        # Video references
        video_chapter_dir = os.path.join(VIDEO_DIR, chapter_slug)
        videos = []
        if os.path.isdir(video_chapter_dir):
            videos = sorted([
                f for f in os.listdir(video_chapter_dir)
                if f.lower().endswith(('.mp4', '.mkv'))
            ])

        if videos:
            doc.add_heading("Video References", level=3)
            for j, video in enumerate(videos, 1):
                vp = doc.add_paragraph(style='List Number')
                run = vp.add_run(f"[Video {j}: {video}]")
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)

        doc.add_heading("Content", level=2)

        # Copy content from source DOCX
        for para in src_doc.paragraphs:
            text = sanitize_text(para.text.strip())
            if not text:
                continue

            if para.style and para.style.name and 'Heading' in para.style.name:
                level = 3
                if '1' in para.style.name:
                    level = 3
                elif '2' in para.style.name:
                    level = 4
                doc.add_heading(text, level=level)
            else:
                p = doc.add_paragraph()
                run = p.add_run(text)
                is_bold = any(r.bold for r in para.runs) if para.runs else False
                if is_bold:
                    run.bold = True

        safe_title = re.sub(r'[^\w\s-]', '', title)[:50].strip()
        doc_filename = f"Lesson_{i:02d}_{safe_title}.docx"
        doc_path = os.path.join(out_dir, doc_filename)
        doc.save(doc_path)
        docs_generated += 1
        print(f"    -> {doc_filename}")

    return docs_generated


def main():
    print("=" * 60)
    print("GUIDEWIRE ACADEMY - LESSON DOC GENERATOR")
    print("=" * 60)

    os.makedirs(OUTPUT_DIR, exist_ok=True)

    # Load assignment manifest if available
    assignment_data = None
    if os.path.isfile(ASSIGNMENT_MANIFEST):
        with open(ASSIGNMENT_MANIFEST, 'r') as f:
            assignment_data = json.load(f)
        print(f"Loaded assignment manifest: {assignment_data.get('totalAssignments', 0)} assignments")
    else:
        print("No assignment manifest found (run organize-assignments.py first for assignment references)")

    total_docs = 0

    for chapter_folder, config in CHAPTER_MAP.items():
        chapter_dir = os.path.join(MATERIAL_DIR, chapter_folder)

        if not os.path.isdir(chapter_dir):
            print(f"\n  SKIP: {chapter_folder} (not found)")
            continue

        # Check if chapter has PPTX or DOCX
        has_pptx = False
        for root, dirs, files in os.walk(chapter_dir):
            if any(f.endswith(('.pptx', '.pptm')) and not f.startswith('~$') for f in files):
                has_pptx = True
                break

        has_docx = any(
            f.endswith('.docx') and not f.startswith('~$')
            for f in os.listdir(chapter_dir)
            if os.path.isfile(os.path.join(chapter_dir, f))
        )

        if has_pptx:
            count = process_chapter(chapter_folder, config, assignment_data)
            total_docs += count
        elif has_docx:
            count = process_docx_chapter(chapter_folder, config, assignment_data)
            total_docs += count
        else:
            print(f"\n  SKIP: {chapter_folder} (no PPTX or DOCX files)")

    print(f"\n{'=' * 60}")
    print(f"GENERATION COMPLETE")
    print(f"{'=' * 60}")
    print(f"Total Word documents generated: {total_docs}")
    print(f"Output: {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
