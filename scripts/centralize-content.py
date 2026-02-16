#!/usr/bin/env python3
"""
Centralize all content layers into content lesson JSONs.

For each lesson, extracts original PPT speaker notes and ensures all 4 layers exist:
  - originalNotes: Raw text from PPT speaker notes (freshly extracted)
  - notes: Enhanced/updated notes (current value, preserved as-is)
  - narration: AI-generated narration (preserved where exists)
  - ocrText: OCR text from slide images (preserved where exists)

Also reports differences between original PPT notes and current enhanced notes.
"""

import json
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx required")
    print("  pip3 install python-pptx")
    sys.exit(1)

MATERIAL_DIR = Path("/Users/sumanthrajkumarnagolu/Documents/Material")
CONTENT_DIR = Path("public/academy/guidewire/content")

CHAPTER_MAP = {
    "ch04": "Chapter 4 - PolicyCenter Introduction",
    "ch05": "Chapter 5 - Claim Center Introduction",
    "ch06": "Chapter 6 - Billing Center Introduction",
    "ch07": "Chapter 7 - InsuranceSuite Developer Fundamentals",
    "ch08": "Chapter 8 - policy center configuration",
    "ch09": "Chapter 9 - ClaimCenter Configuration",
    "ch11": "Chapter 11 - Introduction to Integration",
    "ch13": "Chapter 13 - Rating Introduction",
    "ch14": "Chapter 14 - Rating Configuration",
}


def extract_ppt_notes(pptx_path: str) -> dict:
    """Extract speaker notes from every slide in a PPTX. Returns {slide_number: notes_text}."""
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"    ERROR opening PPTX: {e}")
        return {}

    notes = {}
    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        text = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        notes[slide_num] = text
    return notes


def find_pptx(chapter_dir: Path, source_file: str, source_folder: str = None) -> str:
    """Find the PPTX file, trying multiple strategies."""
    # Strategy 1: Direct path
    if source_folder:
        path = chapter_dir / source_folder / source_file
        if path.exists():
            return str(path)

    # Strategy 2: Recursive search
    for root, dirs, files in os.walk(str(chapter_dir)):
        if source_file in files:
            return os.path.join(root, source_file)

    return None


def main():
    print("=" * 80)
    print("CENTRALIZE CONTENT: Extract Original PPT Notes + Validate All Layers")
    print("=" * 80)

    total_lessons = 0
    total_slides = 0
    total_originals_added = 0
    total_notes_differ = 0
    total_notes_enhanced = 0  # Had no original, now has enhanced
    total_notes_expanded = 0  # Original exists but enhanced is different
    total_notes_identical = 0
    ppt_not_found = []
    diff_report = []

    for ch_slug, ch_name in CHAPTER_MAP.items():
        content_ch_dir = CONTENT_DIR / ch_slug
        chapter_dir = MATERIAL_DIR / ch_name
        index_path = content_ch_dir / "index.json"

        if not content_ch_dir.exists() or not index_path.exists():
            continue

        if not chapter_dir.exists():
            print(f"\n  WARNING: Material dir not found for {ch_slug}: {chapter_dir}")
            continue

        with open(index_path) as f:
            index_data = json.load(f)

        print(f"\n{'='*70}")
        print(f"Chapter: {ch_slug} - {ch_name}")
        print(f"{'='*70}")

        for lesson_info in index_data.get("lessons", []):
            lnum = lesson_info["lessonId"].split("-l")[1]
            lname = f"lesson-{lnum}"
            lesson_path = content_ch_dir / f"{lname}.json"

            if not lesson_path.exists():
                continue

            with open(lesson_path) as f:
                content = json.load(f)

            source_file = content.get("sourceFile", "")
            source_folder = content.get("sourceFolder", "")

            if not source_file:
                print(f"  {lname}: SKIP (no sourceFile)")
                continue

            # Find and extract PPT notes
            pptx_path = find_pptx(chapter_dir, source_file, source_folder)

            if not pptx_path:
                ppt_not_found.append(f"{ch_slug}/{lname}: {source_file}")
                print(f"  {lname}: PPT NOT FOUND ({source_file})")
                # Still ensure all slides have the fields
                for slide in content.get("slides", []):
                    slide.setdefault("originalNotes", "")
                    slide.setdefault("notes", "")
                    slide.setdefault("narration", "")
                    slide.setdefault("ocrText", "")
                    slide.setdefault("ocrWordCount", 0)
                    slide.setdefault("ocrConfidence", 0)
                with open(lesson_path, 'w') as f:
                    json.dump(content, f, indent=2, ensure_ascii=False)
                total_lessons += 1
                continue

            ppt_notes = extract_ppt_notes(pptx_path)
            total_lessons += 1

            lesson_originals = 0
            lesson_differ = 0
            lesson_enhanced = 0
            lesson_expanded = 0
            lesson_identical = 0

            for slide in content.get("slides", []):
                snum = slide.get("slideNumber", 0)
                total_slides += 1

                original = ppt_notes.get(snum, "")
                current_notes = slide.get("notes", "")

                # Set originalNotes
                slide["originalNotes"] = original
                if original:
                    lesson_originals += 1
                    total_originals_added += 1

                # Ensure all fields exist
                slide.setdefault("notes", "")
                slide.setdefault("narration", "")
                slide.setdefault("ocrText", "")
                slide.setdefault("ocrWordCount", 0)
                slide.setdefault("ocrConfidence", 0)

                # Compare original vs current
                if original and current_notes:
                    if original.strip() == current_notes.strip():
                        lesson_identical += 1
                        total_notes_identical += 1
                    else:
                        lesson_expanded += 1
                        total_notes_expanded += 1
                        # Check if current is an expansion (contains original)
                        if original.strip() in current_notes.strip():
                            diff_type = "EXPANDED"
                        else:
                            diff_type = "MODIFIED"
                        diff_report.append({
                            "location": f"{ch_slug}/{lname}/slide-{snum:02d}",
                            "type": diff_type,
                            "originalLen": len(original),
                            "currentLen": len(current_notes),
                        })
                elif not original and current_notes:
                    lesson_enhanced += 1
                    total_notes_enhanced += 1
                    diff_report.append({
                        "location": f"{ch_slug}/{lname}/slide-{snum:02d}",
                        "type": "NEW (no PPT notes, enhanced added)",
                        "originalLen": 0,
                        "currentLen": len(current_notes),
                    })
                elif original and not current_notes:
                    # Original exists but current is empty - copy original
                    slide["notes"] = original
                    lesson_differ += 1
                    total_notes_differ += 1

            # Write updated content
            with open(lesson_path, 'w') as f:
                json.dump(content, f, indent=2, ensure_ascii=False)

            ppt_slide_count = len(ppt_notes)
            content_slide_count = len(content.get("slides", []))
            has_narration = any(s.get("narration") for s in content.get("slides", []))
            has_ocr = any(s.get("ocrText") for s in content.get("slides", []))

            layers = []
            if lesson_originals > 0:
                layers.append(f"orig:{lesson_originals}")
            layers.append(f"notes:{sum(1 for s in content['slides'] if s.get('notes'))}")
            if has_narration:
                layers.append("narration")
            if has_ocr:
                layers.append("ocr")

            status_parts = []
            if lesson_identical:
                status_parts.append(f"{lesson_identical} identical")
            if lesson_enhanced:
                status_parts.append(f"{lesson_enhanced} enhanced")
            if lesson_expanded:
                status_parts.append(f"{lesson_expanded} expanded")

            print(f"  {lname}: {content_slide_count} slides, PPT:{ppt_slide_count} | [{', '.join(layers)}] | {', '.join(status_parts) if status_parts else 'all new'}")

    # Write diff report
    report_path = CONTENT_DIR / "_notes_diff_report.json"
    report = {
        "summary": {
            "totalLessons": total_lessons,
            "totalSlides": total_slides,
            "originalNotesAdded": total_originals_added,
            "notesIdentical": total_notes_identical,
            "notesEnhanced": total_notes_enhanced,
            "notesExpanded": total_notes_expanded,
            "notesCopiedFromOriginal": total_notes_differ,
            "pptNotFound": len(ppt_not_found),
        },
        "pptNotFound": ppt_not_found,
        "differences": diff_report,
    }
    with open(report_path, 'w') as f:
        json.dump(report, f, indent=2, ensure_ascii=False)

    print()
    print("=" * 80)
    print("CENTRALIZATION COMPLETE")
    print("=" * 80)
    print(f"  Lessons processed:        {total_lessons}")
    print(f"  Total slides:             {total_slides}")
    print(f"  Original notes extracted: {total_originals_added}")
    print()
    print(f"  Notes comparison:")
    print(f"    Identical (PPT = current):    {total_notes_identical}")
    print(f"    Enhanced (no PPT, has notes):  {total_notes_enhanced}")
    print(f"    Expanded/Modified:             {total_notes_expanded}")
    print(f"    Copied from original:          {total_notes_differ}")
    print()
    if ppt_not_found:
        print(f"  PPTs not found: {len(ppt_not_found)}")
        for p in ppt_not_found:
            print(f"    - {p}")
    print()
    print(f"  Diff report: {report_path}")
    print()
    print("  Each slide now has:")
    print("    - originalNotes  (raw from PPT speaker notes)")
    print("    - notes          (enhanced/updated version)")
    print("    - narration      (AI-generated, where available)")
    print("    - ocrText        (from slide image OCR)")


if __name__ == "__main__":
    main()
