#!/usr/bin/env python3
"""
Academy Content Pipeline
========================
Processes Guidewire training PPTs from source Material into clean portal content.

Steps:
  1. Copy & organize PPTs from Material into working directory
  2. Analyze & categorize every slide (color detection + position heuristics)
  3. Transcribe ALL slides (extract notes, OCR images via existing slide PNGs)
  4. Extract Q&A pairs from Question+Answer slide pairs
  5. Delete non-content slides from PPT copies → cleaned PPTs
  6. Export slide images from cleaned PPTs
  7. Build final lesson JSONs with structured metadata

Usage:
  python3 scripts/academy-pipeline.py [--step N] [--chapter chNN] [--dry-run]

  --step N       Run only step N (1-7)
  --chapter chNN Process only one chapter (e.g., ch04)
  --dry-run      Show what would be done without writing files
"""

import argparse
import json
import os
import re
import shutil
import sys
from collections import defaultdict
from copy import deepcopy
from io import BytesIO
from pathlib import Path

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# =============================================================================
# Configuration
# =============================================================================

MATERIAL_DIR = Path('/Users/sumanthrajkumarnagolu/Documents/Material')
PROJECT_DIR = Path('/Users/sumanthrajkumarnagolu/Projects/intime-v3')
ACADEMY_DIR = PROJECT_DIR / 'public' / 'academy' / 'guidewire'

# Working directory for PPT copies (we modify these, never the originals)
PPTS_DIR = ACADEMY_DIR / 'ppts'
# Output directories
CONTENT_DIR = ACADEMY_DIR / 'content'
SLIDES_DIR = ACADEMY_DIR / 'slides'

# Chapter mapping: Material folder name → chapter slug
CHAPTER_MAP = {
    'Chapter 1 - Guidewire Cloud Overview': ('ch01', 1),
    'Chapter 2 - Surepath Overview': ('ch02', 2),
    'Chapter 3 - InsuranceSuite Implementation Tools': ('ch03', 3),
    'Chapter 4 - PolicyCenter Introduction': ('ch04', 4),
    'Chapter 5 - Claim Center Introduction': ('ch05', 5),
    'Chapter 6 - Billing Center Introduction': ('ch06', 6),
    'Chapter 7 - InsuranceSuite Developer Fundamentals': ('ch07', 7),
    'Chapter 8 - policy center configuration': ('ch08', 8),
    'Chapter 9 - ClaimCenter Configuration': ('ch09', 9),
    'Chapter 10 - BillingCenter Configuration': ('ch10', 10),
    'Chapter 11 - Introduction to Integration': ('ch11', 11),
    'Chapter 12 - Advanced product Designer': ('ch12', 12),
    'Chapter 13 - Rating Introduction': ('ch13', 13),
    'Chapter 14 - Rating Configuration': ('ch14', 14),
}

# Lesson folder prefixes per chapter (to sort lesson folders correctly)
LESSON_PREFIX = {
    'ch04': 'In_policy_',
    'ch05': 'In_Claim_',
    'ch06': 'BillingCenter_',
    'ch07': 'IS_Fund_',
    'ch08': 'PP_',
    'ch09': '',  # numbered descriptive names
    'ch11': 'In_Integration_',
    'ch12': '',  # flat (videos/ + docx)
    'ch13': 'Ra_Intro_',
    'ch14': 'Ra_Conf_',
}

# =============================================================================
# Title Extraction
# =============================================================================

def extract_lesson_title(ppt_name):
    """Extract a clean lesson title from a PPT filename.
    Strips known prefixes, suffixes, and file extensions."""
    # Remove file extension
    title = ppt_name.rsplit('.', 1)[0] if '.' in ppt_name else ppt_name

    # For backup ch06 PPTs: "Chapter 6 - Billingcenter Introduction - Lesson - 01 - BillingLifeCycle"
    # Extract the last meaningful segment after "Lesson - NN - "
    ch6_match = re.match(r'Chapter\s+\d+.*?Lesson\s*-?\s*\d+\s*-\s*(.+)', title, re.IGNORECASE)
    if ch6_match:
        title = ch6_match.group(1)
        title = re.sub(r'-ppt$', '', title, flags=re.IGNORECASE)

    # Strip common prefixes (order matters: try most specific first)
    prefix_patterns = [
        r'^PC_Intro_\d+[_\s]*',            # PC_Intro_01_Accounts
        r'^In_policy_\d+\s*_?\s*',          # In_policy_02_Policy Transactions
        r'^IS_Claim__?\d+[_\s]*',           # IS_Claim_01_ or IS_Claim__02_
        r'^IS_?\s*Fund\s*[_-]?\s*\d+\s*[_\s-]*',  # IS_Fund_03-, IS_ Fund_04-, IS_Fund -18
        r'^PP_\d+[_\s]*',                   # PP_02_Configuring
        r'^In_Integration_\d+[_\s]*',       # In_Integration_01_
        r'^Ra_Intro_\d+[_\s]*',             # Ra_Intro_01_
        r'^Ra_Conf_\d+[_\s]*',              # Ra_Conf_01_
        r'^Chapter\s+\d+.*?-\s*',           # Remaining chapter prefix patterns
        r'^\d{2}\s*-\s*',                   # Leading "01 - " from ch09
    ]
    for pattern in prefix_patterns:
        title = re.sub(pattern, '', title, flags=re.IGNORECASE)

    # Strip suffixes
    title = re.sub(r'_Configuration$', '', title, flags=re.IGNORECASE)
    title = re.sub(r'-ppt$', '', title, flags=re.IGNORECASE)

    # Clean up
    title = title.strip(' _-')

    # Split CamelCase words (e.g., BillingLifeCycle → Billing Life Cycle)
    # but preserve known compound words
    if '_' not in title and ' ' not in title and len(title) > 1:
        # This is likely a CamelCase name from ch06 backups
        title = re.sub(r'([a-z])([A-Z])', r'\1 \2', title)
        title = re.sub(r'([A-Z]+)([A-Z][a-z])', r'\1 \2', title)

    # Replace underscores with spaces
    title = title.replace('_', ' ')

    # Remove double spaces
    title = re.sub(r'\s+', ' ', title).strip()

    # Title case if all lowercase or all uppercase
    if title == title.lower() or title == title.upper():
        title = title.title()

    # Fix common casing issues
    replacements = {
        'Policycenter': 'PolicyCenter',
        'Claimcenter': 'ClaimCenter',
        'Billingcenter': 'BillingCenter',
        'Guidewire': 'Guidewire',
        'Gosu': 'Gosu',
        'Xml': 'XML',
        'Restful': 'RESTful',
        'Soap': 'SOAP',
        'Api': 'API',
        'Apd': 'APD',
    }
    for wrong, right in replacements.items():
        title = re.sub(r'\b' + wrong + r'\b', right, title, flags=re.IGNORECASE)

    return title if title else ppt_name.rsplit('.', 1)[0]


# =============================================================================
# Slide Type Detection
# =============================================================================

def get_slide_image(slide):
    """Extract the embedded image from a slide (PPTs are image-only)."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return Image.open(BytesIO(shape.image.blob)).convert('RGB')
    return None


def get_slide_notes(slide):
    """Extract speaker notes text from a slide."""
    if slide.has_notes_slide:
        return slide.notes_slide.notes_text_frame.text.strip()
    return ''


def detect_slide_colors(img):
    """Compute teal% and green% for a slide image."""
    arr = np.array(img)[:, :, :3]
    h, w = arr.shape[:2]
    total = h * w

    # Teal/cyan: R<100, G>140, B>180
    teal_pct = ((arr[:, :, 0] < 100) & (arr[:, :, 1] > 140) & (arr[:, :, 2] > 180)).sum() / total * 100

    # Green/mint (answer slides): R<100, G>170, B in 140-200
    green_pct = (
        (arr[:, :, 0] < 100) & (arr[:, :, 1] > 170) &
        (arr[:, :, 2] > 130) & (arr[:, :, 2] < 200)
    ).sum() / total * 100

    return teal_pct, green_pct


def categorize_slide(slide, index, total_slides):
    """
    Categorize a slide based on color detection, position, and notes.

    Returns one of:
      'title', 'objectives', 'content', 'question', 'answer',
      'review', 'demo', 'vm_instructions', 'end'
    """
    img = get_slide_image(slide)
    notes = get_slide_notes(slide)
    notes_lower = notes.lower()

    # Color detection
    teal_pct, green_pct = 0.0, 0.0
    if img:
        teal_pct, green_pct = detect_slide_colors(img)

    # --- High confidence color-based detection ---

    # Question slides: >70% teal
    if teal_pct > 70:
        return 'question'

    # Answer slides: >70% green/mint
    if green_pct > 70:
        return 'answer'

    # --- Position + color heuristics ---

    # Slide 1 is always the title slide
    if index == 0:
        return 'title'

    # Slide 2 (index 1) with partial teal (30-50%) = objectives
    if index == 1 and teal_pct > 25:
        return 'objectives'

    # Last few slides with partial teal (30-50%) = review/summary
    if index >= total_slides - 5 and teal_pct > 25:
        return 'review'

    # --- Notes-based detection ---

    # VM Instructions
    if 'virtual machine' in notes_lower or 'start your vm' in notes_lower:
        return 'vm_instructions'

    # Demo slides (notes mention demo/demonstration)
    if notes_lower.startswith('demo:') or notes_lower.startswith('demonstration:'):
        return 'demo'

    # End slide: last slide with no notes and no meaningful content
    if index == total_slides - 1 and not notes:
        return 'end'

    return 'content'


# =============================================================================
# Step 1: Copy & Organize PPTs
# =============================================================================

def _extract_folder_num(folder):
    """Extract numeric ordering from a folder name."""
    nums = re.findall(r'(\d+)', folder.name)
    return int(nums[0]) if nums else 0


def _find_ppt_in_folder(folder):
    """Find the first PPT file in a folder, ignoring temp files."""
    ppts = [
        f for f in folder.iterdir()
        if f.suffix.lower() in ('.pptx', '.pptm') and not f.name.startswith('~$')
    ]
    return ppts[0] if ppts else None


def _find_videos_in_folder(folder):
    """Find video files in a folder."""
    return sorted([
        f for f in folder.iterdir()
        if f.suffix.lower() in ('.mp4', '.mkv', '.avi', '.wmv')
    ])


def _get_lesson_folders(chapter_path):
    """Get sorted lesson subfolders from a chapter directory."""
    folders = []
    for item in chapter_path.iterdir():
        if not item.is_dir():
            continue
        if item.name.startswith('.') or item.name in ('backups', 'videos', 'Assigments'):
            continue
        folders.append(item)
    folders.sort(key=_extract_folder_num)
    return folders


def _resolve_chapter_path(slug, chapter_path):
    """Handle chapters with extra nesting levels.
    Returns the actual directory containing lesson folders."""
    if slug == 'ch08':
        # ch08: Chapter 8/.../Policy Center 10.0 Configuration/PP_*
        for item in chapter_path.iterdir():
            if item.is_dir() and 'configuration' in item.name.lower():
                return item
    elif slug == 'ch11':
        # ch11: Chapter 11/Introduction to Integration/In_Integration_*
        for item in chapter_path.iterdir():
            if item.is_dir() and 'integration' in item.name.lower():
                return item
    return chapter_path


def _find_ch06_ppts():
    """Special handling for ch06: PPTs split between main folder and backups."""
    main_path = MATERIAL_DIR / 'Chapter 6 - Billing Center Introduction'
    backup_path = MATERIAL_DIR / 'backups' / 'Chapter 5 - Billingcenter Introduction'

    # Build a map of lesson_num → ppt_path from main folder
    lesson_ppts = {}  # lesson_num → { ppt_path, lesson_folder, videos }

    if main_path.exists():
        for folder in _get_lesson_folders(main_path):
            num = _extract_folder_num(folder)
            ppt = _find_ppt_in_folder(folder)
            videos = _find_videos_in_folder(folder)
            if ppt:
                lesson_ppts[num] = {
                    'ppt_path': ppt,
                    'ppt_name': ppt.name,
                    'lesson_folder': folder.name,
                    'videos': [v.name for v in videos],
                }

    # Fill gaps from backups
    if backup_path.exists():
        for folder in _get_lesson_folders(backup_path):
            # Extract lesson number from backup folder names like
            # "Chapter 6 - Billingcenter Introduction - Lesson - 01 - BillingLifeCycle"
            nums = re.findall(r'Lesson\s*-?\s*(\d+)', folder.name)
            if not nums:
                nums = re.findall(r'(\d+)', folder.name)
            if not nums:
                continue
            num = int(nums[0])

            # Only use backup if main folder doesn't have a PPT for this lesson
            if num not in lesson_ppts:
                ppt = _find_ppt_in_folder(folder)
                if ppt:
                    # For backups, the PPT filename is often truncated.
                    # Extract title from the folder name which has the full name.
                    folder_title_match = re.search(r'Lesson\s*-?\s*\d+\s*-\s*(.+)', folder.name)
                    override_title = None
                    if folder_title_match:
                        raw = folder_title_match.group(1).strip()
                        # Split CamelCase and clean up
                        raw = re.sub(r'([a-z])([A-Z])', r'\1 \2', raw)
                        raw = re.sub(r'([A-Z]+)([A-Z][a-z])', r'\1 \2', raw)
                        override_title = raw

                    lesson_ppts[num] = {
                        'ppt_path': ppt,
                        'ppt_name': ppt.name,
                        'lesson_folder': folder.name,
                        'videos': [],
                        'override_title': override_title,
                    }

    # Convert to sorted list with sequential lesson numbers
    lessons = []
    for num in sorted(lesson_ppts):
        info = lesson_ppts[num]
        lessons.append({
            'lesson_num': num,  # Use actual folder number, not sequential
            **info,
        })
    return lessons


def find_ppts_in_material():
    """Scan Material directory and return organized list of PPTs per chapter/lesson."""
    results = {}  # chapter_slug → [{ lesson_num, ppt_path, lesson_folder, videos }]

    for folder_name, (slug, ch_id) in CHAPTER_MAP.items():
        # Skip chapters with no PPTs
        if slug in ('ch01', 'ch02', 'ch03', 'ch10', 'ch12'):
            continue

        # ch06 has special handling (split between main and backups)
        if slug == 'ch06':
            lessons = _find_ch06_ppts()
            if lessons:
                results[slug] = lessons
            continue

        chapter_path = MATERIAL_DIR / folder_name
        if not chapter_path.exists():
            print(f'  WARNING: Chapter folder not found: {folder_name}')
            continue

        # Handle extra nesting
        chapter_path = _resolve_chapter_path(slug, chapter_path)

        # Get lesson folders
        lesson_folders = _get_lesson_folders(chapter_path)
        lessons = []

        for lesson_folder in lesson_folders:
            num = _extract_folder_num(lesson_folder)
            ppt = _find_ppt_in_folder(lesson_folder)
            videos = _find_videos_in_folder(lesson_folder)

            if ppt:
                lessons.append({
                    'lesson_num': num,  # Use folder number, not sequential
                    'ppt_path': ppt,
                    'ppt_name': ppt.name,
                    'lesson_folder': lesson_folder.name,
                    'videos': [v.name for v in videos],
                })

        if lessons:
            results[slug] = lessons

    return results


def step1_copy_ppts(chapter_filter=None, dry_run=False):
    """Copy PPTs from Material into organized working directory.
    Clears existing PPT copies first to avoid stale files."""
    print('\n' + '=' * 60)
    print('STEP 1: Copy & Organize PPTs')
    print('=' * 60)

    ppt_map = find_ppts_in_material()

    total_copied = 0
    for slug in sorted(ppt_map):
        if chapter_filter and slug != chapter_filter:
            continue

        lessons = ppt_map[slug]
        chapter_out = PPTS_DIR / slug

        # Clear existing PPT copies for this chapter
        if chapter_out.exists() and not dry_run:
            shutil.rmtree(chapter_out)

        print(f'\n  {slug}: {len(lessons)} lessons')

        for lesson in lessons:
            dest = chapter_out / f'lesson-{lesson["lesson_num"]:02d}.pptx'
            title = lesson.get('override_title') or extract_lesson_title(lesson['ppt_name'])
            print(f'    Lesson {lesson["lesson_num"]:2d}: {title} ({lesson["ppt_name"][:40]})')

            if not dry_run:
                chapter_out.mkdir(parents=True, exist_ok=True)
                shutil.copy2(lesson['ppt_path'], dest)
                total_copied += 1

    print(f'\n  Total: {total_copied} PPTs copied')
    return ppt_map


# =============================================================================
# Step 2: Analyze & Categorize All Slides
# =============================================================================

def step2_categorize(chapter_filter=None, dry_run=False):
    """Analyze every slide in every PPT and categorize them."""
    print('\n' + '=' * 60)
    print('STEP 2: Analyze & Categorize Slides')
    print('=' * 60)

    results = {}  # slug → lesson_num → [{ slide_num, category, notes, has_image }]
    stats = defaultdict(int)

    for chapter_dir in sorted(PPTS_DIR.iterdir()):
        if not chapter_dir.is_dir():
            continue
        slug = chapter_dir.name
        if chapter_filter and slug != chapter_filter:
            continue

        results[slug] = {}

        for ppt_file in sorted(chapter_dir.glob('lesson-*.pptx')):
            lesson_num = int(ppt_file.stem.split('-')[1])
            prs = Presentation(str(ppt_file))
            total = len(prs.slides)

            slides_info = []
            for i, slide in enumerate(prs.slides):
                category = categorize_slide(slide, i, total)
                notes = get_slide_notes(slide)
                img = get_slide_image(slide)

                slides_info.append({
                    'slide_index': i,
                    'slide_num': i + 1,
                    'category': category,
                    'notes': notes,
                    'has_image': img is not None,
                })
                stats[category] += 1

            results[slug][lesson_num] = slides_info

            # Print summary for this lesson
            cats = defaultdict(int)
            for s in slides_info:
                cats[s['category']] += 1
            cat_str = ', '.join(f'{k}={v}' for k, v in sorted(cats.items()))
            print(f'  {slug}/lesson-{lesson_num:02d}: {total} slides → {cat_str}')

    print(f'\n  Category totals:')
    for cat in sorted(stats):
        print(f'    {cat:16s}: {stats[cat]}')

    # Save categorization results
    if not dry_run:
        out_path = ACADEMY_DIR / 'pipeline-categorization.json'
        with open(out_path, 'w') as f:
            json.dump(results, f, indent=2)
        print(f'\n  Saved to {out_path}')

    return results


# =============================================================================
# Step 3: Transcribe All Slides
# =============================================================================

def step3_transcribe(categorization, chapter_filter=None, dry_run=False):
    """
    Transcribe all slides — save notes and image data for every slide
    INCLUDING non-content slides (so we preserve objectives text, review
    summaries, etc. for indexing).
    """
    print('\n' + '=' * 60)
    print('STEP 3: Transcribe All Slides')
    print('=' * 60)

    transcriptions = {}  # slug → lesson_num → { slides, metadata }

    for chapter_dir in sorted(PPTS_DIR.iterdir()):
        if not chapter_dir.is_dir():
            continue
        slug = chapter_dir.name
        if chapter_filter and slug != chapter_filter:
            continue

        transcriptions[slug] = {}

        for ppt_file in sorted(chapter_dir.glob('lesson-*.pptx')):
            lesson_num = int(ppt_file.stem.split('-')[1])
            prs = Presentation(str(ppt_file))
            cat_data = categorization.get(slug, {}).get(lesson_num, [])

            slides_data = []
            objectives_text = []
            review_text = []
            vm_instructions = []

            for i, slide in enumerate(prs.slides):
                cat_info = cat_data[i] if i < len(cat_data) else {}
                category = cat_info.get('category', 'content')
                notes = get_slide_notes(slide)

                slide_record = {
                    'slide_num': i + 1,
                    'category': category,
                    'notes': notes,
                }

                # Collect special metadata
                if category == 'objectives' and notes:
                    objectives_text.append(notes)
                elif category == 'review' and notes:
                    review_text.append(notes)
                elif category == 'vm_instructions' and notes:
                    vm_instructions.append(notes)

                slides_data.append(slide_record)

            transcriptions[slug][lesson_num] = {
                'total_slides': len(prs.slides),
                'slides': slides_data,
                'objectives': objectives_text,
                'review_summary': review_text,
                'vm_instructions': vm_instructions,
            }

            obj_count = len(objectives_text)
            review_count = len(review_text)
            notes_count = sum(1 for s in slides_data if s['notes'])
            print(f'  {slug}/lesson-{lesson_num:02d}: {len(slides_data)} slides, {notes_count} with notes, obj={obj_count}, review={review_count}')

    if not dry_run:
        out_path = ACADEMY_DIR / 'pipeline-transcriptions.json'
        with open(out_path, 'w') as f:
            json.dump(transcriptions, f, indent=2)
        print(f'\n  Saved to {out_path}')

    return transcriptions


# =============================================================================
# Step 4: Extract Q&A Pairs
# =============================================================================

def step4_extract_qa(categorization, chapter_filter=None, dry_run=False):
    """
    Find Question+Answer slide pairs and extract the Q&A.
    Question slides (teal) are followed by Answer slides (green).
    Since text is baked into images, we use the existing slide PNGs
    and the OCR data if available, or save image paths for later OCR.
    """
    print('\n' + '=' * 60)
    print('STEP 4: Extract Question & Answer Pairs')
    print('=' * 60)

    qa_pairs = {}  # slug → lesson_num → [{ question_slide, answer_slide, question_img, answer_img }]
    total_pairs = 0

    for slug in sorted(categorization):
        if chapter_filter and slug != chapter_filter:
            continue

        qa_pairs[slug] = {}

        for lesson_num in sorted(categorization[slug]):
            slides = categorization[slug][lesson_num]
            pairs = []

            i = 0
            while i < len(slides):
                if slides[i]['category'] == 'question':
                    pair = {
                        'question_slide_num': slides[i]['slide_num'],
                        'question_notes': slides[i].get('notes', ''),
                    }
                    # Check if next slide is an answer
                    if i + 1 < len(slides) and slides[i + 1]['category'] == 'answer':
                        pair['answer_slide_num'] = slides[i + 1]['slide_num']
                        pair['answer_notes'] = slides[i + 1].get('notes', '')
                        i += 2
                    else:
                        pair['answer_slide_num'] = None
                        pair['answer_notes'] = ''
                        i += 1
                    pairs.append(pair)
                    total_pairs += 1
                else:
                    i += 1

            if pairs:
                qa_pairs[slug][lesson_num] = pairs
                print(f'  {slug}/lesson-{lesson_num:02d}: {len(pairs)} Q&A pairs')

    print(f'\n  Total Q&A pairs found: {total_pairs}')

    if not dry_run:
        out_path = ACADEMY_DIR / 'pipeline-qa-pairs.json'
        with open(out_path, 'w') as f:
            json.dump(qa_pairs, f, indent=2)
        print(f'\n  Saved to {out_path}')

    return qa_pairs


# =============================================================================
# Step 5: Delete Non-Content Slides from PPT Copies
# =============================================================================

def step5_clean_ppts(categorization, chapter_filter=None, dry_run=False):
    """
    Remove non-content slides from PPT copies.
    Keeps only 'content' slides. Deletes: title, objectives, question,
    answer, review, demo, vm_instructions, end.
    """
    print('\n' + '=' * 60)
    print('STEP 5: Clean PPTs (remove non-content slides)')
    print('=' * 60)

    total_removed = 0
    total_kept = 0

    for chapter_dir in sorted(PPTS_DIR.iterdir()):
        if not chapter_dir.is_dir():
            continue
        slug = chapter_dir.name
        if chapter_filter and slug != chapter_filter:
            continue

        for ppt_file in sorted(chapter_dir.glob('lesson-*.pptx')):
            lesson_num = int(ppt_file.stem.split('-')[1])
            cat_data = categorization.get(slug, {}).get(lesson_num, [])

            # Find indices to DELETE (in reverse order so indices stay valid)
            to_delete = []
            to_keep = 0
            for info in cat_data:
                if info['category'] != 'content':
                    to_delete.append(info['slide_index'])
                else:
                    to_keep += 1

            if not to_delete:
                print(f'  {slug}/lesson-{lesson_num:02d}: nothing to remove (all content)')
                continue

            print(f'  {slug}/lesson-{lesson_num:02d}: removing {len(to_delete)} slides, keeping {to_keep}')
            total_removed += len(to_delete)
            total_kept += to_keep

            if dry_run:
                continue

            # Open PPT, delete slides in reverse order
            prs = Presentation(str(ppt_file))
            for idx in sorted(to_delete, reverse=True):
                rId = prs.slides._sldIdLst[idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[idx]

            prs.save(str(ppt_file))

    print(f'\n  Total: removed {total_removed} slides, kept {total_kept}')


# =============================================================================
# Step 6: Export Slide Images from Cleaned PPTs
# =============================================================================

def step6_export_images(chapter_filter=None, dry_run=False):
    """
    Export slide images from cleaned PPTs.
    Since each slide is a single embedded picture, we extract the image blob directly.
    Clears existing slide directories first to avoid stale images.
    """
    print('\n' + '=' * 60)
    print('STEP 6: Export Slide Images from Cleaned PPTs')
    print('=' * 60)

    total_images = 0

    for chapter_dir in sorted(PPTS_DIR.iterdir()):
        if not chapter_dir.is_dir():
            continue
        slug = chapter_dir.name
        if chapter_filter and slug != chapter_filter:
            continue

        # Clear existing slides for this chapter
        slides_chapter_dir = SLIDES_DIR / slug
        if slides_chapter_dir.exists() and not dry_run:
            shutil.rmtree(slides_chapter_dir)
            print(f'  Cleared existing slides for {slug}')

        for ppt_file in sorted(chapter_dir.glob('lesson-*.pptx')):
            lesson_num = int(ppt_file.stem.split('-')[1])
            prs = Presentation(str(ppt_file))

            out_dir = SLIDES_DIR / slug / f'lesson-{lesson_num:02d}'
            count = 0

            for i, slide in enumerate(prs.slides):
                img = get_slide_image(slide)
                if img:
                    if not dry_run:
                        out_dir.mkdir(parents=True, exist_ok=True)
                        img_path = out_dir / f'slide-{i + 1:02d}.png'
                        img.save(str(img_path), 'PNG')
                    count += 1
                    total_images += 1

            print(f'  {slug}/lesson-{lesson_num:02d}: {count} images exported')

    print(f'\n  Total: {total_images} slide images')


# =============================================================================
# Step 7: Build Final Lesson JSONs
# =============================================================================

def step7_build_jsons(categorization, transcriptions, qa_pairs, ppt_map,
                      chapter_filter=None, dry_run=False):
    """
    Build clean lesson JSON files with:
    - Content slides only (notes, slide number)
    - Extracted objectives (from transcription of objectives slides)
    - Extracted Q&A pairs (from question/answer slides)
    - Review/takeaways (from transcription of review slides)
    - Video references
    - Assignment references
    """
    print('\n' + '=' * 60)
    print('STEP 7: Build Final Lesson JSONs')
    print('=' * 60)

    for slug in sorted(categorization):
        if chapter_filter and slug != chapter_filter:
            continue

        ch_id = int(slug[2:])
        chapter_out = CONTENT_DIR / slug

        # Clear old lesson JSONs for this chapter (preserve non-lesson files)
        if chapter_out.exists() and not dry_run:
            for old_file in chapter_out.glob('lesson-*.json'):
                old_file.unlink()

        for lesson_num in sorted(categorization[slug]):
            cat_data = categorization[slug][lesson_num]
            trans_data = transcriptions.get(slug, {}).get(lesson_num, {})
            lesson_qa = qa_pairs.get(slug, {}).get(lesson_num, [])

            # Content slides only (renumbered sequentially)
            content_slides = []
            content_idx = 0
            for info in cat_data:
                if info['category'] != 'content':
                    continue
                content_idx += 1
                notes = info.get('notes', '')
                notes_lower = notes.lower()
                # Detect demo slides within content category
                slide_type = 'content'
                if notes_lower.startswith('demo') or 'demonstration' in notes_lower[:50]:
                    slide_type = 'demo'
                content_slides.append({
                    'slideNumber': content_idx,
                    'originalSlideNumber': info['slide_num'],
                    'slideType': slide_type,
                    'notes': notes,
                    'hasImage': info.get('has_image', True),
                    'bodyParagraphs': [],
                    'hasTable': False,
                    'hasChart': False,
                    'tableData': None,
                    'mediaReferences': [],
                })

            # Build Q&A data
            question_data = []
            for pair in lesson_qa:
                question_data.append({
                    'originalQuestionSlide': pair['question_slide_num'],
                    'originalAnswerSlide': pair.get('answer_slide_num'),
                    'questionNotes': pair.get('question_notes', ''),
                    'answerNotes': pair.get('answer_notes', ''),
                    # These will be filled by OCR step later
                    'question': '',
                    'answer': '',
                })

            # Get video refs from ppt_map
            ppt_lessons = ppt_map.get(slug, [])
            lesson_info = next((l for l in ppt_lessons if l['lesson_num'] == lesson_num), None)
            videos = []
            if lesson_info:
                for vi, vname in enumerate(lesson_info['videos']):
                    videos.append({
                        'index': vi,
                        'filename': vname,
                        'path': f'/academy/guidewire/videos/{slug}/{vname}',
                    })

            # Estimate reading time (2 min per content slide + 5 min per video)
            est_minutes = max(10, len(content_slides) * 2 + len(videos) * 5)

            lesson_json = {
                'lessonId': f'{slug}-l{lesson_num:02d}',
                'chapterId': ch_id,
                'chapterSlug': slug,
                'lessonNumber': lesson_num,
                'title': (lesson_info.get('override_title') or extract_lesson_title(lesson_info['ppt_name'])) if lesson_info else f'Lesson {lesson_num}',
                'sourceFile': lesson_info['ppt_name'] if lesson_info else '',
                'sourceFolder': lesson_info['lesson_folder'] if lesson_info else '',
                'totalSlides': len(content_slides),
                'originalTotalSlides': trans_data.get('total_slides', 0),
                'estimatedMinutes': est_minutes,
                'slides': content_slides,
                'videos': videos,
                'objectives': trans_data.get('objectives', []),
                'reviewSummary': trans_data.get('review_summary', []),
                'questionAnswers': question_data,
                'vmInstructions': trans_data.get('vm_instructions', []),
                'hasNotes': any(s['notes'] for s in content_slides),
                'hasTables': False,
                'hasImages': True,
            }

            if not dry_run:
                chapter_out.mkdir(parents=True, exist_ok=True)
                out_path = chapter_out / f'lesson-{lesson_num:02d}.json'
                with open(out_path, 'w') as f:
                    json.dump(lesson_json, f, indent=2, ensure_ascii=False)

            print(f'  {slug}/lesson-{lesson_num:02d}: {len(content_slides)} content slides, {len(question_data)} Q&A, {len(videos)} videos')

    # Build chapter index files
    if not dry_run:
        for slug in sorted(categorization):
            if chapter_filter and slug != chapter_filter:
                continue
            ch_dir = CONTENT_DIR / slug
            lessons = sorted(ch_dir.glob('lesson-*.json'))
            index = {
                'chapterSlug': slug,
                'totalLessons': len(lessons),
                'lessons': [],
            }
            for lf in lessons:
                with open(lf) as f:
                    data = json.load(f)
                index['lessons'].append({
                    'lessonId': data['lessonId'],
                    'lessonNumber': data['lessonNumber'],
                    'title': data['title'],
                    'totalSlides': data['totalSlides'],
                    'estimatedMinutes': data['estimatedMinutes'],
                    'hasVideos': len(data['videos']) > 0,
                    'hasQuestions': len(data['questionAnswers']) > 0,
                })
            with open(ch_dir / 'index.json', 'w') as f:
                json.dump(index, f, indent=2)
            print(f'  {slug}/index.json: {len(lessons)} lessons indexed')


# =============================================================================
# Main
# =============================================================================

def main():
    parser = argparse.ArgumentParser(description='Academy Content Pipeline')
    parser.add_argument('--step', type=int, help='Run only this step (1-7)')
    parser.add_argument('--chapter', type=str, help='Process only this chapter (e.g., ch04)')
    parser.add_argument('--dry-run', action='store_true', help='Show what would be done')
    args = parser.parse_args()

    chapter = args.chapter
    dry_run = args.dry_run

    if dry_run:
        print('*** DRY RUN — no files will be written ***\n')

    # Shared state between steps
    ppt_map = None
    categorization = None
    transcriptions = None
    qa_pairs = None

    steps_to_run = [args.step] if args.step else [1, 2, 3, 4, 5, 6, 7]

    if 1 in steps_to_run:
        ppt_map = step1_copy_ppts(chapter, dry_run)

    if 2 in steps_to_run:
        if not PPTS_DIR.exists():
            print('ERROR: Run step 1 first to copy PPTs')
            sys.exit(1)
        categorization = step2_categorize(chapter, dry_run)

    if 3 in steps_to_run:
        # Load categorization if not already in memory
        if categorization is None:
            cat_path = ACADEMY_DIR / 'pipeline-categorization.json'
            if cat_path.exists():
                with open(cat_path) as f:
                    categorization = json.load(f)
                # Convert lesson_num keys back to int
                for slug in categorization:
                    categorization[slug] = {int(k): v for k, v in categorization[slug].items()}
            else:
                print('ERROR: Run step 2 first')
                sys.exit(1)
        transcriptions = step3_transcribe(categorization, chapter, dry_run)

    if 4 in steps_to_run:
        if categorization is None:
            cat_path = ACADEMY_DIR / 'pipeline-categorization.json'
            if cat_path.exists():
                with open(cat_path) as f:
                    categorization = json.load(f)
                for slug in categorization:
                    categorization[slug] = {int(k): v for k, v in categorization[slug].items()}
            else:
                print('ERROR: Run step 2 first')
                sys.exit(1)
        qa_pairs = step4_extract_qa(categorization, chapter, dry_run)

    if 5 in steps_to_run:
        if categorization is None:
            cat_path = ACADEMY_DIR / 'pipeline-categorization.json'
            if cat_path.exists():
                with open(cat_path) as f:
                    categorization = json.load(f)
                for slug in categorization:
                    categorization[slug] = {int(k): v for k, v in categorization[slug].items()}
            else:
                print('ERROR: Run step 2 first')
                sys.exit(1)
        step5_clean_ppts(categorization, chapter, dry_run)

    if 6 in steps_to_run:
        step6_export_images(chapter, dry_run)

    if 7 in steps_to_run:
        # Load all required data
        if categorization is None:
            cat_path = ACADEMY_DIR / 'pipeline-categorization.json'
            with open(cat_path) as f:
                categorization = json.load(f)
            for slug in categorization:
                categorization[slug] = {int(k): v for k, v in categorization[slug].items()}

        if transcriptions is None:
            trans_path = ACADEMY_DIR / 'pipeline-transcriptions.json'
            if trans_path.exists():
                with open(trans_path) as f:
                    transcriptions = json.load(f)
                for slug in transcriptions:
                    transcriptions[slug] = {int(k): v for k, v in transcriptions[slug].items()}
            else:
                transcriptions = {}

        if qa_pairs is None:
            qa_path = ACADEMY_DIR / 'pipeline-qa-pairs.json'
            if qa_path.exists():
                with open(qa_path) as f:
                    qa_pairs = json.load(f)
                for slug in qa_pairs:
                    qa_pairs[slug] = {int(k): v for k, v in qa_pairs[slug].items()}
            else:
                qa_pairs = {}

        if ppt_map is None:
            ppt_map = find_ppts_in_material()

        step7_build_jsons(categorization, transcriptions, qa_pairs, ppt_map, chapter, dry_run)

    print('\n✓ Pipeline complete')


if __name__ == '__main__':
    main()
